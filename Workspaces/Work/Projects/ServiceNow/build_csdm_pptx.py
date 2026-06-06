from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

LAYERS = [
    {
        "title": "Business Offering",
        "desc": "What the business receives — a named, business-facing service",
        "tag": "BUSINESS",
        "bg": RGBColor(0xD4, 0xED, 0xDA),
        "border": RGBColor(0x5A, 0x9E, 0x6F),
        "tag_bg": RGBColor(0x5A, 0x9E, 0x6F),
    },
    {
        "title": "Technical Offering",
        "desc": "What IT provides to fulfill that business need",
        "tag": "IT",
        "bg": RGBColor(0xCC, 0xE5, 0xFF),
        "border": RGBColor(0x4A, 0x86, 0xC8),
        "tag_bg": RGBColor(0x4A, 0x86, 0xC8),
    },
    {
        "title": "Service Instance",
        "desc": "The live, running service in a specific environment",
        "tag": "OPERATIONAL",
        "bg": RGBColor(0xFF, 0xF3, 0xCD),
        "border": RGBColor(0xC8, 0xA1, 0x35),
        "tag_bg": RGBColor(0xC8, 0xA1, 0x35),
    },
    {
        "title": "Application Instance",
        "desc": "A specific deployment of the app (Prod, Dev, Test…)",
        "tag": "DEPLOYED",
        "bg": RGBColor(0xFF, 0xE5, 0xD0),
        "border": RGBColor(0xC8, 0x70, 0x40),
        "tag_bg": RGBColor(0xC8, 0x70, 0x40),
    },
    {
        "title": "Business Application",
        "desc": "The software product itself",
        "tag": "SOFTWARE",
        "bg": RGBColor(0xF8, 0xD7, 0xDA),
        "border": RGBColor(0xB8, 0x4A, 0x52),
        "tag_bg": RGBColor(0xB8, 0x4A, 0x52),
    },
]

def hex_color(r, g, b):
    return "{:02X}{:02X}{:02X}".format(r, g, b)

def add_rounded_rect(slide, x, y, w, h, bg_color, border_color, radius=Pt(8)):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 5,  # rounded rect = 5
        x, y, w, h
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    line = shape.line
    line.color.rgb = border_color
    line.width = Pt(1.5)
    # Set corner radius via XML
    sp = shape._element
    spPr = sp.find('.//' + pptx.oxml.ns.qn('p:spPr'))
    if spPr is None:
        spPr = sp.find('.//' + pptx.oxml.ns.qn('spPr'))
    prstGeom = spPr.find(pptx.oxml.ns.qn('a:prstGeom'))
    if prstGeom is not None:
        avLst = prstGeom.find(pptx.oxml.ns.qn('a:avLst'))
        if avLst is not None:
            gd = avLst.find(pptx.oxml.ns.qn('a:gd'))
            if gd is not None:
                gd.set('fmla', 'val 30000')
    return shape

def add_text_box(slide, text, x, y, w, h, font_size, bold=False, color=RGBColor(0,0,0), align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── Slide 1: Title ──────────────────────────────────────────────────────────
    blank_layout = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(blank_layout)

    bg = slide1.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x5F)

    add_text_box(slide1, "CSDM Service Layers",
                 Inches(1), Inches(2.5), Inches(11.33), Inches(1.2),
                 Pt(48), bold=True, color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)

    add_text_box(slide1, "Understanding how Business Applications connect to Business Value",
                 Inches(1.5), Inches(3.8), Inches(10.33), Inches(0.8),
                 Pt(22), bold=False, color=RGBColor(0xCC,0xDD,0xEE), align=PP_ALIGN.CENTER)

    add_text_box(slide1, "Read the stack bottom to top: Technical → Business",
                 Inches(2), Inches(5.0), Inches(9.33), Inches(0.6),
                 Pt(16), bold=False, color=RGBColor(0xAA,0xBB,0xCC), align=PP_ALIGN.CENTER)

    # ── Slide 2: Stack Diagram ──────────────────────────────────────────────────
    slide2 = prs.slides.add_slide(blank_layout)
    bg2 = slide2.background.fill
    bg2.solid()
    bg2.fore_color.rgb = RGBColor(0xF9, 0xF7, 0xF2)

    add_text_box(slide2, "CSDM: The Service Layer Stack",
                 Inches(0.4), Inches(0.15), Inches(12), Inches(0.55),
                 Pt(28), bold=True, color=RGBColor(0x1E,0x3A,0x5F), align=PP_ALIGN.CENTER)

    add_text_box(slide2, "↑  Technical  →  Business  ↑     (read bottom to top)",
                 Inches(0.4), Inches(0.68), Inches(12), Inches(0.35),
                 Pt(13), bold=False, color=RGBColor(0x88,0x88,0x88), align=PP_ALIGN.CENTER)

    box_h = Inches(0.92)
    gap = Inches(0.18)
    box_w = Inches(9.6)
    left_x = Inches(1.87)
    top_y = Inches(1.15)

    # Draw layers top-to-bottom (Business Offering at top = index 0)
    for i, layer in enumerate(LAYERS):
        y = top_y + i * (box_h + gap)

        # Main box
        shape = slide2.shapes.add_shape(5, left_x, y, box_w, box_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = layer["bg"]
        shape.line.color.rgb = layer["border"]
        shape.line.width = Pt(1.8)

        # Title text
        title_box = slide2.shapes.add_textbox(left_x + Inches(0.2), y + Inches(0.08),
                                               Inches(5.5), Inches(0.42))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = layer["title"]
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
        run.font.name = "Calibri"

        # Desc text
        desc_box = slide2.shapes.add_textbox(left_x + Inches(0.2), y + Inches(0.48),
                                              Inches(7.0), Inches(0.38))
        tf2 = desc_box.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = layer["desc"]
        run2.font.size = Pt(13)
        run2.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
        run2.font.name = "Calibri"

        # Tag pill (right side)
        tag_box = slide2.shapes.add_textbox(left_x + Inches(8.05), y + Inches(0.28),
                                             Inches(1.35), Inches(0.36))
        tag_shape = slide2.shapes.add_shape(5,
                                             left_x + Inches(7.95),
                                             y + Inches(0.25),
                                             Inches(1.45), Inches(0.42))
        tag_shape.fill.solid()
        tag_shape.fill.fore_color.rgb = layer["tag_bg"]
        tag_shape.line.color.rgb = layer["tag_bg"]

        tf3 = tag_box.text_frame
        p3 = tf3.paragraphs[0]
        p3.alignment = PP_ALIGN.CENTER
        run3 = p3.add_run()
        run3.text = layer["tag"]
        run3.font.size = Pt(11)
        run3.font.bold = True
        run3.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run3.font.name = "Calibri"

        # Arrow between layers (except after last)
        if i < len(LAYERS) - 1:
            arrow_y = y + box_h + Inches(0.01)
            arrow_box = slide2.shapes.add_textbox(left_x, arrow_y, box_w, gap + Inches(0.04))
            tf_a = arrow_box.text_frame
            p_a = tf_a.paragraphs[0]
            p_a.alignment = PP_ALIGN.CENTER
            run_a = p_a.add_run()
            run_a.text = "↑"
            run_a.font.size = Pt(14)
            run_a.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
            run_a.font.name = "Calibri"

    # Footer
    add_text_box(slide2, "CSDM — Common Service Data Model  |  ServiceNow",
                 Inches(0.4), Inches(7.1), Inches(12), Inches(0.35),
                 Pt(11), bold=False, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.CENTER)

    out = r"C:\Users\manuel.b.vazquez\Documents\Claude AI OS\Workspaces\Work\Projects\ServiceNow\CSDM-Service-Layers.pptx"
    prs.save(out)
    print(f"Saved: {out}")

if __name__ == "__main__":
    build()
