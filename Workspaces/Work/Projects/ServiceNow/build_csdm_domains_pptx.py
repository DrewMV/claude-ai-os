from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pptx.oxml.ns as ns

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Domain definitions ────────────────────────────────────────────────────────
DOMAINS = [
    {
        "name": "Foundation",
        "subtitle": "Base referential data — required before all other domains",
        "color": RGBColor(0x34, 0x49, 0x5E),
        "light": RGBColor(0xEC, 0xF0, 0xF1),
        "tag": "PREREQUISITE",
        "components": [
            ("Business Process", "cmdb_ci_business_process", "Core process definitions used across the CMDB"),
            ("Value Stream", "cmn_value_stream", "End-to-end flow delivering value to customers"),
            ("Product Models", "cmdb_model", "Hardware, software, and application model catalog"),
            ("Contracts", "ast_contract", "Vendor and service agreements"),
            ("Organizational Data", "core_company / business_unit / cmn_department", "Company, BU, Department, Location, Users, Groups"),
            ("Life Cycle Control", "life_cycle_mapping / life_cycle_control", "Standard CSDM life-cycle stage and status pairs"),
        ],
    },
    {
        "name": "Ideation & Strategy",
        "subtitle": "Ideas and strategies for new or improved services",
        "color": RGBColor(0x6C, 0x3A, 0x83),
        "light": RGBColor(0xF5, 0xEA, 0xF8),
        "tag": "STRATEGIC",
        "components": [
            ("Product Idea", "sn_align_core_product_idea", "Concept for a new or enhanced digital product"),
            ("Planning Item", "sn_align_core_planning_item", "Work item aligning ideas to delivery plans"),
            ("Strategic Plan", "sn_gf_plan", "Long-range plan linked to organizational goals"),
            ("Strategic Priority", "sn_gf_strategy", "High-level priority driving investment decisions"),
            ("Goal", "sn_gf_goal", "Measurable target tied to a strategic priority"),
        ],
    },
    {
        "name": "Design & Planning",
        "subtitle": "Logical designs for digital products being built or bought",
        "color": RGBColor(0x1A, 0x5C, 0x96),
        "light": RGBColor(0xE8, 0xF3, 0xFF),
        "tag": "DESIGN",
        "components": [
            ("Business Capability", "cmdb_ci_business_capability", "High-level capability required to execute the business model"),
            ("Business Application", "cmdb_ci_business_app", "Purchased or internally developed application (logical)"),
            ("Information Object", "cmdb_ci_information_object", "Type of data exchanged (PII, PCI DSS, HIPAA, etc.)"),
        ],
    },
    {
        "name": "Build & Integration",
        "subtitle": "Logical development details during the software build effort",
        "color": RGBColor(0x1A, 0x7A, 0x6E),
        "light": RGBColor(0xE6, 0xF7, 0xF5),
        "tag": "DEVELOPMENT",
        "components": [
            ("SDLC Component", "cmdb_ci_sdlc_component", "Unique code effort — microservice, API, or software element"),
            ("DevOps Pipeline Data", "Various DevOps tables", "Development pipeline and CI/CD model objects"),
            ("AI Digital Asset", "alm_ai_digital_asset", "Deployable AI source: System, Model, Dataset, Prompt assets"),
        ],
    },
    {
        "name": "Service Delivery",
        "subtitle": "Operational domain — deployed apps, infrastructure, technology services",
        "color": RGBColor(0xC0, 0x6A, 0x10),
        "light": RGBColor(0xFF, 0xF0, 0xE0),
        "tag": "OPERATIONAL",
        "components": [
            ("Technology Mgmt Service", "cmdb_ci_service_technical", "Identifies the technology provider your business consumes"),
            ("Technology Mgmt Offering", "service_offering (Technical)", "Stratification of technical service by environment/pricing"),
            ("Service Instance", "cmdb_ci_service_auto", "Central operational entity — deployed system or app stack"),
            ("Application", "cmdb_ci_appl", "Discoverable running instance of code on a host"),
            ("API Components", "cmdb_ci_api / api_gateway / api_backend", "API, API Gateway, and API Backend CIs"),
            ("Infrastructure CIs", "cmdb_ci_*", "Physical and logical components: servers, DBs, network gear"),
        ],
    },
    {
        "name": "Service Consumption",
        "subtitle": "Portfolios and catalogs where consumers request business services",
        "color": RGBColor(0x27, 0x7A, 0x3C),
        "light": RGBColor(0xE8, 0xF8, 0xED),
        "tag": "CONSUMER-FACING",
        "components": [
            ("Business Service", "cmdb_ci_service_business", "Broad business objective dependent on underlying technology"),
            ("Business Service Offering", "service_offering (Business)", "Specific offerings consumers select — SLAs, pricing, support"),
            ("Service Catalog", "sc_catalog", "Where business offerings are published and browsed"),
            ("Business Service Portfolio", "service_portfolio", "Hierarchical collection of business services"),
        ],
    },
    {
        "name": "Manage Portfolio",
        "subtitle": "Cross-domain visibility layer for service owners",
        "color": RGBColor(0x7D, 0x1F, 0x2A),
        "light": RGBColor(0xFD, 0xEB, 0xED),
        "tag": "GOVERNANCE",
        "components": [
            ("SPM Taxonomy Node", "spm_taxonomy_node", "Connects Business Services and Tech Services to the portfolio"),
            ("Cross-Domain Visibility", "All CSDM domains", "Service owners oversee apps, services, and instances holistically"),
        ],
    },
]

# ── Helpers ───────────────────────────────────────────────────────────────────

def rgb_to_hex(rgb):
    return "{:02X}{:02X}{:02X}".format(rgb[0], rgb[1], rgb[2])

def add_solid_rect(slide, x, y, w, h, fill_rgb, line_rgb=None, line_pt=0):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb and line_pt:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape

def add_label(slide, text, x, y, w, h, size, bold=False,
              color=RGBColor(0,0,0), align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb

def add_tag_pill(slide, text, x, y, bg_rgb):
    w, h = Inches(1.55), Inches(0.30)
    add_solid_rect(slide, x, y, w, h, bg_rgb)
    add_label(slide, text, x, y + Inches(0.02), w, h,
              size=9, bold=True, color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)

# ── Slide builders ────────────────────────────────────────────────────────────

def build_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)

    # Decorative stripe
    add_solid_rect(slide, 0, Inches(5.8), SLIDE_W, Inches(0.08),
                   RGBColor(0x3A, 0xA0, 0xFF))

    add_label(slide, "Common Service Data Model",
              Inches(1), Inches(1.6), Inches(11.33), Inches(0.8),
              size=42, bold=True, color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)

    add_label(slide, "CSDM",
              Inches(1), Inches(2.5), Inches(11.33), Inches(0.6),
              size=28, bold=False, color=RGBColor(0x3A,0xA0,0xFF), align=PP_ALIGN.CENTER)

    add_label(slide, "Domains, Key Components & ServiceNow Tables",
              Inches(1.5), Inches(3.3), Inches(10.33), Inches(0.5),
              size=18, color=RGBColor(0xAA,0xCC,0xEE), align=PP_ALIGN.CENTER)

    add_label(slide, "Source: ServiceNow Australia Release Documentation",
              Inches(1), Inches(6.1), Inches(11.33), Inches(0.35),
              size=12, italic=True, color=RGBColor(0x66,0x88,0xAA), align=PP_ALIGN.CENTER)


def build_overview_slide(prs):
    """7-box domain map on one slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0xF4, 0xF6, 0xF9)

    add_label(slide, "CSDM — The 7 Domains at a Glance",
              Inches(0.3), Inches(0.12), Inches(12.73), Inches(0.55),
              size=26, bold=True, color=RGBColor(0x1E,0x2D,0x40), align=PP_ALIGN.CENTER)

    # Layout: 4 top, 3 bottom
    BOX_W, BOX_H = Inches(2.9), Inches(2.3)
    GAP = Inches(0.22)
    TOP_Y = Inches(0.85)
    BOT_Y = Inches(3.38)

    top_x_start = (SLIDE_W - 4*BOX_W - 3*GAP) / 2
    bot_x_start = (SLIDE_W - 3*BOX_W - 2*GAP) / 2

    positions = []
    for i in range(4):
        positions.append((top_x_start + i*(BOX_W+GAP), TOP_Y))
    for i in range(3):
        positions.append((bot_x_start + i*(BOX_W+GAP), BOT_Y))

    for idx, domain in enumerate(DOMAINS):
        x, y = positions[idx]
        # Box
        add_solid_rect(slide, x, y, BOX_W, BOX_H,
                       domain["light"], domain["color"], line_pt=1.5)
        # Top color bar
        add_solid_rect(slide, x, y, BOX_W, Inches(0.32), domain["color"])
        # Domain number
        add_label(slide, str(idx+1), x + Inches(0.12), y + Inches(0.02),
                  Inches(0.25), Inches(0.28), size=14, bold=True,
                  color=RGBColor(0xFF,0xFF,0xFF))
        # Domain name
        add_label(slide, domain["name"],
                  x + Inches(0.35), y + Inches(0.02),
                  BOX_W - Inches(0.45), Inches(0.30),
                  size=13, bold=True, color=RGBColor(0xFF,0xFF,0xFF))
        # Subtitle
        add_label(slide, domain["subtitle"],
                  x + Inches(0.12), y + Inches(0.38),
                  BOX_W - Inches(0.24), Inches(0.55),
                  size=10, italic=True, color=domain["color"])
        # Component bullets (first 3 only, space-limited)
        for ci, comp in enumerate(domain["components"][:4]):
            add_label(slide, f"• {comp[0]}",
                      x + Inches(0.14), y + Inches(0.95 + ci * 0.30),
                      BOX_W - Inches(0.28), Inches(0.30),
                      size=10, color=RGBColor(0x22,0x22,0x22))

    # Footer
    add_label(slide, "Each domain is detailed on the following slides",
              Inches(0.3), Inches(7.1), Inches(12.73), Inches(0.30),
              size=11, italic=True, color=RGBColor(0x88,0x88,0x88), align=PP_ALIGN.CENTER)


def build_domain_slide(prs, domain, index):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0xFA, 0xFB, 0xFC)

    # Left color bar
    add_solid_rect(slide, 0, 0, Inches(0.18), SLIDE_H, domain["color"])

    # Header band
    add_solid_rect(slide, Inches(0.18), 0, SLIDE_W - Inches(0.18), Inches(1.15), domain["color"])

    # Domain number badge
    badge = add_solid_rect(slide, Inches(0.3), Inches(0.18),
                           Inches(0.6), Inches(0.6), RGBColor(0xFF,0xFF,0xFF))
    add_label(slide, str(index), Inches(0.3), Inches(0.18), Inches(0.6), Inches(0.6),
              size=22, bold=True, color=domain["color"], align=PP_ALIGN.CENTER)

    # Domain name
    add_label(slide, domain["name"],
              Inches(1.1), Inches(0.10), Inches(9.5), Inches(0.60),
              size=30, bold=True, color=RGBColor(0xFF,0xFF,0xFF))

    # Tag pill
    r, g, b = domain["color"]
    darker = RGBColor(max(0, r-40), max(0, g-40), max(0, b-40))
    add_tag_pill(slide, domain["tag"], Inches(11.2), Inches(0.22), darker)

    # Subtitle
    add_label(slide, domain["subtitle"],
              Inches(1.1), Inches(0.72), Inches(10.5), Inches(0.38),
              size=14, italic=True, color=RGBColor(0xFF,0xFF,0xFF))

    # Components table
    col_w = [Inches(2.8), Inches(3.2), Inches(5.6)]
    col_x = [Inches(0.28), Inches(3.18), Inches(6.48)]
    row_h = Inches(0.68)
    table_top = Inches(1.28)

    # Header row
    headers = ["Component", "Table / Class", "Description"]
    for ci, hdr in enumerate(headers):
        add_solid_rect(slide, col_x[ci], table_top, col_w[ci], Inches(0.42), domain["color"])
        add_label(slide, hdr, col_x[ci] + Inches(0.08), table_top + Inches(0.04),
                  col_w[ci] - Inches(0.1), Inches(0.38),
                  size=12, bold=True, color=RGBColor(0xFF,0xFF,0xFF))

    max_rows = min(len(domain["components"]), 7)
    for ri, comp in enumerate(domain["components"][:max_rows]):
        y = table_top + Inches(0.42) + ri * row_h
        row_bg = domain["light"] if ri % 2 == 0 else RGBColor(0xFF,0xFF,0xFF)
        for ci in range(3):
            add_solid_rect(slide, col_x[ci], y, col_w[ci], row_h,
                           row_bg, domain["color"], line_pt=0.5)
        texts = [comp[0], comp[1], comp[2]]
        sizes = [13, 11, 12]
        bolds = [True, False, False]
        colors = [domain["color"], RGBColor(0x55,0x55,0x55), RGBColor(0x22,0x22,0x22)]
        for ci in range(3):
            add_label(slide, texts[ci],
                      col_x[ci] + Inches(0.1), y + Inches(0.06),
                      col_w[ci] - Inches(0.14), row_h - Inches(0.1),
                      size=sizes[ci], bold=bolds[ci], color=colors[ci])

    # Footer
    add_label(slide, f"Domain {index} of 7  |  CSDM — Common Service Data Model",
              Inches(0.28), Inches(7.12), Inches(12.5), Inches(0.30),
              size=10, italic=True, color=RGBColor(0xAA,0xAA,0xAA))


def build_closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0x0D, 0x1B, 0x2A)

    add_solid_rect(slide, 0, Inches(3.55), SLIDE_W, Inches(0.08), RGBColor(0x3A,0xA0,0xFF))

    add_label(slide, "How the Domains Connect",
              Inches(1), Inches(1.5), Inches(11.33), Inches(0.7),
              size=34, bold=True, color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)

    flow = "Foundation  →  Ideation & Strategy  →  Design & Planning  →  Build & Integration  →  Service Delivery  →  Service Consumption"
    add_label(slide, flow,
              Inches(0.5), Inches(2.4), Inches(12.33), Inches(0.55),
              size=14, color=RGBColor(0x3A,0xA0,0xFF), align=PP_ALIGN.CENTER)

    add_label(slide, "↕",
              Inches(1), Inches(3.0), Inches(11.33), Inches(0.4),
              size=20, color=RGBColor(0x66,0x88,0xAA), align=PP_ALIGN.CENTER)

    add_label(slide, "Manage Portfolio  (cross-domain visibility layer)",
              Inches(1), Inches(3.45), Inches(11.33), Inches(0.45),
              size=16, bold=True, color=RGBColor(0xFF,0xCC,0x66), align=PP_ALIGN.CENTER)

    bullets = [
        "• Foundation data must exist before any other domain is populated",
        "• Ideation → Design → Build represents the product lifecycle (logical, non-operational)",
        "• Service Delivery is the operational heart — running instances, infrastructure, offerings",
        "• Service Consumption is business-facing — what end users see and request",
        "• Manage Portfolio gives service owners holistic visibility across all layers",
    ]
    for bi, b in enumerate(bullets):
        add_label(slide, b, Inches(1.5), Inches(4.1 + bi * 0.48),
                  Inches(10.33), Inches(0.45),
                  size=14, color=RGBColor(0xCC,0xDD,0xEE))


def build_ba_types_slide(prs):
    """Business Application types, classification, and hierarchy within Design & Planning."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0xF4, 0xF6, 0xF9)

    D_BLUE = RGBColor(0x1A, 0x5C, 0x96)
    L_BLUE = RGBColor(0xE8, 0xF3, 0xFF)
    MID_BLUE = RGBColor(0x4A, 0x86, 0xC8)

    # Header band
    add_solid_rect(slide, 0, 0, SLIDE_W, Inches(1.05), D_BLUE)
    add_label(slide, "Design & Planning — Business Application Types & Relationships",
              Inches(0.3), Inches(0.08), Inches(12.73), Inches(0.52),
              size=22, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    add_label(slide, "All are Business Applications (cmdb_ci_business_app) — differentiated by classification attributes and relationships",
              Inches(0.5), Inches(0.60), Inches(12.33), Inches(0.35),
              size=12, italic=True, color=RGBColor(0xBB, 0xD4, 0xF0), align=PP_ALIGN.CENTER)

    # ── Left panel: Hierarchy diagram ────────────────────────────────────────
    # Platform BA — ServiceNow
    SN_X, SN_Y, SN_W, SN_H = Inches(0.3), Inches(1.25), Inches(4.2), Inches(1.1)
    add_solid_rect(slide, SN_X, SN_Y, SN_W, SN_H, L_BLUE, D_BLUE, line_pt=2.5)
    add_solid_rect(slide, SN_X, SN_Y, SN_W, Inches(0.32), D_BLUE)
    add_label(slide, "PLATFORM",
              SN_X + Inches(0.12), SN_Y + Inches(0.04), Inches(1.2), Inches(0.26),
              size=9, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_label(slide, "ServiceNow",
              SN_X + Inches(0.15), SN_Y + Inches(0.37), SN_W - Inches(0.25), Inches(0.40),
              size=20, bold=True, color=D_BLUE)
    add_label(slide, "Application Type: COTS  |  Install Type: Cloud",
              SN_X + Inches(0.15), SN_Y + Inches(0.76), SN_W - Inches(0.25), Inches(0.26),
              size=10, italic=True, color=RGBColor(0x44, 0x44, 0x44))

    # Connector line going down then splitting
    SPLIT_Y = SN_Y + SN_H + Inches(0.22)
    MID_X = SN_X + SN_W / 2

    # Vertical stem
    add_solid_rect(slide, MID_X - Inches(0.02), SN_Y + SN_H,
                   Inches(0.04), Inches(0.22), RGBColor(0x88, 0x88, 0x88))
    # Horizontal bar
    CHILD_W = Inches(1.55)
    CHILD_GAP = Inches(0.22)
    TOTAL_CHILD_W = 3 * CHILD_W + 2 * CHILD_GAP
    CHILD_START_X = SN_X + (SN_W - TOTAL_CHILD_W) / 2
    add_solid_rect(slide, CHILD_START_X + CHILD_W / 2, SPLIT_Y - Inches(0.02),
                   TOTAL_CHILD_W - CHILD_W, Inches(0.04), RGBColor(0x88, 0x88, 0x88))

    # Three product BAs
    PRODUCTS = [
        ("ITSM", "IT Service\nManagement"),
        ("ITOM", "IT Operations\nManagement"),
        ("HRSD", "HR Service\nDelivery"),
    ]
    CHILD_Y = SPLIT_Y + Inches(0.15)
    CHILD_H = Inches(1.15)
    for pi, (code, full) in enumerate(PRODUCTS):
        cx = CHILD_START_X + pi * (CHILD_W + CHILD_GAP)
        # Vertical drop to each child
        add_solid_rect(slide, cx + CHILD_W / 2 - Inches(0.02), SPLIT_Y - Inches(0.02),
                       Inches(0.04), Inches(0.17), RGBColor(0x88, 0x88, 0x88))
        # Child box
        add_solid_rect(slide, cx, CHILD_Y, CHILD_W, CHILD_H,
                       RGBColor(0xD6, 0xEA, 0xFF), MID_BLUE, line_pt=1.8)
        add_solid_rect(slide, cx, CHILD_Y, CHILD_W, Inches(0.28), MID_BLUE)
        add_label(slide, "PRODUCT",
                  cx + Inches(0.08), CHILD_Y + Inches(0.03), CHILD_W - Inches(0.1), Inches(0.22),
                  size=8, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        add_label(slide, code,
                  cx + Inches(0.08), CHILD_Y + Inches(0.32), CHILD_W - Inches(0.1), Inches(0.36),
                  size=18, bold=True, color=D_BLUE, align=PP_ALIGN.CENTER)
        add_label(slide, full,
                  cx + Inches(0.08), CHILD_Y + Inches(0.66), CHILD_W - Inches(0.1), Inches(0.44),
                  size=9, color=RGBColor(0x33, 0x33, 0x33), align=PP_ALIGN.CENTER)

    # "All are cmdb_ci_business_app" label
    add_label(slide, "All four records live in cmdb_ci_business_app",
              SN_X, CHILD_Y + CHILD_H + Inches(0.10), SN_W, Inches(0.28),
              size=10, italic=True, color=RGBColor(0x55, 0x55, 0x55), align=PP_ALIGN.CENTER)

    # ── Right panel: Two ways to model relationship ───────────────────────────
    RP_X = Inches(4.85)

    add_label(slide, "Two ways to express the relationship in CSDM",
              RP_X, Inches(1.25), Inches(8.2), Inches(0.35),
              size=14, bold=True, color=D_BLUE)

    # Option 1
    add_solid_rect(slide, RP_X, Inches(1.68), Inches(8.2), Inches(1.55),
                   RGBColor(0xEA, 0xF4, 0xFF), D_BLUE, line_pt=1.2)
    add_label(slide, "Option 1 — Platform Attribute",
              RP_X + Inches(0.15), Inches(1.74), Inches(7.9), Inches(0.36),
              size=13, bold=True, color=D_BLUE)
    add_label(slide,
              'On each product BA (ITSM, ITOM, HRSD), set the Platform field to "ServiceNow".\n'
              "Quick to implement. Useful for filtering and reporting by platform.",
              RP_X + Inches(0.15), Inches(2.10), Inches(7.9), Inches(0.55),
              size=11, color=RGBColor(0x22, 0x22, 0x22))
    add_label(slide, "ITSM.platform = \"ServiceNow\"",
              RP_X + Inches(0.15), Inches(2.64), Inches(7.9), Inches(0.30),
              size=10, italic=True, color=RGBColor(0x44, 0x44, 0x88))

    # Option 2
    add_solid_rect(slide, RP_X, Inches(3.38), Inches(8.2), Inches(1.55),
                   RGBColor(0xEA, 0xF4, 0xFF), D_BLUE, line_pt=1.2)
    add_label(slide, "Option 2 — Business Application Reference",
              RP_X + Inches(0.15), Inches(3.44), Inches(7.9), Inches(0.36),
              size=13, bold=True, color=D_BLUE)
    add_label(slide,
              "On each product BA, use the Business Application reference field to point directly\n"
              "to the ServiceNow platform BA record. Creates an explicit, traceable CI relationship.",
              RP_X + Inches(0.15), Inches(3.80), Inches(7.9), Inches(0.55),
              size=11, color=RGBColor(0x22, 0x22, 0x22))
    add_label(slide, "ITSM.business_application → ServiceNow (cmdb_ci_business_app record)",
              RP_X + Inches(0.15), Inches(4.35), Inches(7.9), Inches(0.30),
              size=10, italic=True, color=RGBColor(0x44, 0x44, 0x88))

    # Classification attributes table
    add_label(slide, "Business Application Classification Attributes",
              RP_X, Inches(5.08), Inches(8.2), Inches(0.35),
              size=13, bold=True, color=D_BLUE)

    attrs = [
        ("Application Type", "Nature of software", "COTS, Custom, Open Source"),
        ("Architecture Type", "Structural pattern", "Web Based, Microservices, Monolith"),
        ("Install Type", "Where it lives", "Cloud, On-Premise, Hybrid"),
        ("Platform", "Underlying platform", "ServiceNow, Salesforce, SAP"),
    ]
    COL_W = [Inches(2.1), Inches(2.1), Inches(3.8)]
    COL_X = [RP_X, RP_X + Inches(2.15), RP_X + Inches(4.30)]
    ROW_H = Inches(0.36)
    HDR_Y = Inches(5.46)

    hdrs = ["Attribute", "What it describes", "Example values"]
    for ci, hdr in enumerate(hdrs):
        add_solid_rect(slide, COL_X[ci], HDR_Y, COL_W[ci], ROW_H, D_BLUE)
        add_label(slide, hdr, COL_X[ci] + Inches(0.06), HDR_Y + Inches(0.05),
                  COL_W[ci] - Inches(0.1), ROW_H - Inches(0.08),
                  size=10, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    for ri, row in enumerate(attrs):
        row_y = HDR_Y + ROW_H + ri * ROW_H
        row_bg = L_BLUE if ri % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
        for ci, val in enumerate(row):
            add_solid_rect(slide, COL_X[ci], row_y, COL_W[ci], ROW_H,
                           row_bg, D_BLUE, line_pt=0.5)
            add_label(slide, val, COL_X[ci] + Inches(0.06), row_y + Inches(0.05),
                      COL_W[ci] - Inches(0.1), ROW_H - Inches(0.06),
                      size=10, bold=(ci == 0), color=RGBColor(0x11, 0x11, 0x11))

    # Footer
    add_label(slide, "CSDM — Design & Planning Domain  |  Business Application Classification",
              Inches(0.28), Inches(7.12), Inches(12.5), Inches(0.30),
              size=10, italic=True, color=RGBColor(0xAA, 0xAA, 0xAA))


def build_relationship_slide(prs):
    """Design ↔ Service Delivery relationship: Business App → Service Instance → Service Delivery Network."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0xF4, 0xF6, 0xF9)

    # Title
    add_label(slide, "Key Relationship: Design → Service Delivery",
              Inches(0.3), Inches(0.12), Inches(12.73), Inches(0.52),
              size=24, bold=True, color=RGBColor(0x1E, 0x2D, 0x40), align=PP_ALIGN.CENTER)
    add_label(slide, "How a logical Business Application becomes an operational Service Instance backed by infrastructure",
              Inches(0.5), Inches(0.65), Inches(12.33), Inches(0.35),
              size=13, italic=True, color=RGBColor(0x55, 0x55, 0x55), align=PP_ALIGN.CENTER)

    # ── Domain label banners ─────────────────────────────────────────────────
    # Design & Planning banner (left)
    add_solid_rect(slide, Inches(0.28), Inches(1.15), Inches(4.0), Inches(0.38),
                   RGBColor(0x1A, 0x5C, 0x96))
    add_label(slide, "Design & Planning Domain",
              Inches(0.28), Inches(1.15), Inches(4.0), Inches(0.38),
              size=12, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

    # Service Delivery banner (right)
    add_solid_rect(slide, Inches(5.6), Inches(1.15), Inches(7.4), Inches(0.38),
                   RGBColor(0xC0, 0x6A, 0x10))
    add_label(slide, "Service Delivery Domain",
              Inches(5.6), Inches(1.15), Inches(7.4), Inches(0.38),
              size=12, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

    # ── Business Application box ──────────────────────────────────────────────
    BA_X, BA_Y, BA_W, BA_H = Inches(0.28), Inches(1.7), Inches(4.0), Inches(1.5)
    add_solid_rect(slide, BA_X, BA_Y, BA_W, BA_H,
                   RGBColor(0xE8, 0xF3, 0xFF), RGBColor(0x1A, 0x5C, 0x96), line_pt=2)
    add_label(slide, "Business Application",
              BA_X + Inches(0.15), BA_Y + Inches(0.12),
              BA_W - Inches(0.25), Inches(0.40),
              size=16, bold=True, color=RGBColor(0x1A, 0x5C, 0x96))
    add_label(slide, "cmdb_ci_business_app",
              BA_X + Inches(0.15), BA_Y + Inches(0.52),
              BA_W - Inches(0.25), Inches(0.28),
              size=10, italic=True, color=RGBColor(0x44, 0x44, 0x44))
    add_label(slide, "Logical — not yet deployed.\nThe software product as planned or purchased.",
              BA_X + Inches(0.15), BA_Y + Inches(0.80),
              BA_W - Inches(0.25), Inches(0.60),
              size=11, color=RGBColor(0x33, 0x33, 0x33))

    # ── Arrow: Business App → Service Instance ────────────────────────────────
    ARR1_X = BA_X + BA_W + Inches(0.08)
    ARR1_Y = BA_Y + BA_H / 2 - Inches(0.22)
    add_label(slide, "→", ARR1_X, ARR1_Y, Inches(0.7), Inches(0.44),
              size=28, bold=True, color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)
    add_label(slide, "deployed as",
              ARR1_X, ARR1_Y + Inches(0.38), Inches(0.7), Inches(0.28),
              size=9, italic=True, color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)

    # ── Service Instance box ──────────────────────────────────────────────────
    SI_X, SI_Y, SI_W, SI_H = Inches(5.6), Inches(1.7), Inches(3.6), Inches(1.5)
    add_solid_rect(slide, SI_X, SI_Y, SI_W, SI_H,
                   RGBColor(0xFF, 0xF0, 0xE0), RGBColor(0xC0, 0x6A, 0x10), line_pt=2)
    add_label(slide, "Service Instance",
              SI_X + Inches(0.15), SI_Y + Inches(0.12),
              SI_W - Inches(0.25), Inches(0.40),
              size=16, bold=True, color=RGBColor(0xC0, 0x6A, 0x10))
    add_label(slide, "cmdb_ci_service_auto",
              SI_X + Inches(0.15), SI_Y + Inches(0.52),
              SI_W - Inches(0.25), Inches(0.28),
              size=10, italic=True, color=RGBColor(0x44, 0x44, 0x44))
    add_label(slide, "Operational — live and running.\nThe deployed system in a specific environment.",
              SI_X + Inches(0.15), SI_Y + Inches(0.80),
              SI_W - Inches(0.25), Inches(0.60),
              size=11, color=RGBColor(0x33, 0x33, 0x33))

    # ── Arrow: Service Instance → Service Delivery Network (downward) ─────────
    ARR2_X = SI_X + SI_W / 2 - Inches(0.5)
    ARR2_Y = SI_Y + SI_H + Inches(0.05)
    add_label(slide, "↓", ARR2_X, ARR2_Y, Inches(1.0), Inches(0.40),
              size=28, bold=True, color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)
    add_label(slide, "depends on",
              ARR2_X, ARR2_Y + Inches(0.38), Inches(1.0), Inches(0.28),
              size=9, italic=True, color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)

    # ── Service Delivery Network box (below Service Instance) ─────────────────
    NET_X, NET_Y, NET_W, NET_H = Inches(4.5), Inches(4.0), Inches(5.8), Inches(2.85)
    add_solid_rect(slide, NET_X, NET_Y, NET_W, NET_H,
                   RGBColor(0xFF, 0xF8, 0xEC), RGBColor(0xC0, 0x6A, 0x10), line_pt=1.5)
    add_label(slide, "Service Delivery Network",
              NET_X + Inches(0.15), NET_Y + Inches(0.10),
              NET_W - Inches(0.25), Inches(0.40),
              size=15, bold=True, color=RGBColor(0xC0, 0x6A, 0x10))
    add_label(slide, "Infrastructure CIs that the Service Instance runs on",
              NET_X + Inches(0.15), NET_Y + Inches(0.50),
              NET_W - Inches(0.25), Inches(0.28),
              size=10, italic=True, color=RGBColor(0x55, 0x55, 0x55))

    infra = [
        ("Servers / Compute", "cmdb_ci_server"),
        ("Databases", "cmdb_ci_database"),
        ("Network Devices", "cmdb_ci_netgear"),
        ("Applications (running)", "cmdb_ci_appl"),
        ("Storage", "cmdb_ci_storage_device"),
        ("Cloud Resources", "cmdb_ci_cloud_*"),
    ]
    col2_x = NET_X + Inches(2.85)
    for ri, (name, tbl) in enumerate(infra):
        row_x = NET_X + Inches(0.15) if ri < 3 else col2_x
        row_y = NET_Y + Inches(0.88) + (ri % 3) * Inches(0.58)
        add_solid_rect(slide, row_x, row_y, Inches(2.55), Inches(0.50),
                       RGBColor(0xFF, 0xED, 0xD5), RGBColor(0xC0, 0x6A, 0x10), line_pt=0.8)
        add_label(slide, name, row_x + Inches(0.1), row_y + Inches(0.04),
                  Inches(2.35), Inches(0.26), size=11, bold=True,
                  color=RGBColor(0x7A, 0x40, 0x10))
        add_label(slide, tbl, row_x + Inches(0.1), row_y + Inches(0.27),
                  Inches(2.35), Inches(0.20), size=9, italic=True,
                  color=RGBColor(0x77, 0x77, 0x77))

    # ── Note callout ─────────────────────────────────────────────────────────
    add_solid_rect(slide, Inches(0.28), Inches(3.6), Inches(4.0), Inches(3.25),
                   RGBColor(0xEE, 0xF5, 0xFF), RGBColor(0x1A, 0x5C, 0x96), line_pt=1.2)
    add_label(slide, "Key Distinction",
              Inches(0.38), Inches(3.68), Inches(3.8), Inches(0.35),
              size=12, bold=True, color=RGBColor(0x1A, 0x5C, 0x96))
    notes = [
        "Business Application is logical —",
        "it exists in Design & Planning",
        "before anything is deployed.",
        "",
        "Service Instance is operational —",
        "it IS the deployed, running system",
        "in a specific environment (Prod,",
        "Dev, Test, etc.).",
        "",
        "The Service Delivery Network",
        "represents everything the Service",
        "Instance depends on to function.",
    ]
    for ni, note in enumerate(notes):
        add_label(slide, note, Inches(0.38), Inches(4.08) + ni * Inches(0.24),
                  Inches(3.7), Inches(0.26),
                  size=11, color=RGBColor(0x22, 0x22, 0x55))

    # Footer
    add_label(slide, "CSDM — Design & Planning  |  Service Delivery Domain Relationship",
              Inches(0.28), Inches(7.12), Inches(12.5), Inches(0.30),
              size=10, italic=True, color=RGBColor(0xAA, 0xAA, 0xAA))


# ── Main ──────────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_title_slide(prs)
    build_overview_slide(prs)
    for i, domain in enumerate(DOMAINS):
        build_domain_slide(prs, domain, i+1)
    build_ba_types_slide(prs)
    build_relationship_slide(prs)
    build_closing_slide(prs)

    out = r"C:\Users\manuel.b.vazquez\Documents\Claude AI OS\Workspaces\Work\Projects\ServiceNow\CSDM-Domains.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({2 + len(DOMAINS) + 3} slides)")

if __name__ == "__main__":
    build()
