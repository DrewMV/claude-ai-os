"""
SUPERSEDED (reconciliation note, 2026-07-27).
This 38-decision vote deck was already replaced by build_ccb_review_deck.py for the
7/21 meeting, and neither matches the class slides in the deck actually DELIVERED to
PPL - "CMDB-Data-Dictionary-CCB-2026 (July 21 2026) - DELIVERED to PPL.pptx" - which
uses the CI Class Reference "Managed Attributes" format (Server `environment` =
Mandatory). Authoritative content mirror: cmdb-health-completeness-correctness-stories.md.
Retained for history only.

CMDB Data Dictionary - CCB Vote Deck (combined).

Built for the 7/21/2026 CMDB CCB. ONE deck for the whole hour:
  1  Title
  2  At a glance + how to vote (overview + the two pre-vote flags)
  3  EDUCATION - the four decision types decoded (Mandatory / Audited / Owner / Values)
  4  EDUCATION - IRE / reconciliation precedence in 60 seconds
  5-13  One slide per class - all 9 classes, all 38 decisions, mirroring the
         Data Dictionary CCB Decision Sheet 1:1 (# | Attribute | Decision |
         Context | Team suggestion | CCB ruling[blank for live capture]).

Source of truth for the 38 items: data-dictionary-ccb-decision-sheet.md (38 items,
9 classes). Team-suggestion column is the delivery team's proposed default, NOT a
decision; the CCB ruling column is authoritative and captured live on 7/21.

Theme mirrors the CMDB-CSDM deck family (Class Reference / PI-3 / CO6) so it reads
as a set.

Run:     python build_ccb_vote_deck.py
Output:  CMDB-Data-Dictionary-CCB-Vote-<EDITION>.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

EDITION   = "2026-07-21"
TITLE     = "CMDB Data Dictionary - CCB Vote"
SUBTITLE  = "38 governance decisions across 9 CI classes  -  CMP Stage 1"
AUTHOR    = "Manuel Vazquez  -  Scrum Master, CMDB-CSDM"
FOOTLABEL = "PPL CMDB-CSDM  |  Data Dictionary CCB Vote (CMP Stage 1)  |  Jul 21, 2026"
CAPTION   = ("Team suggestion = delivery team's proposed default (NOT a decision).  "
             "CCB ruling is authoritative - the chair captures Accept / Amend / Defer live.")

# --------------------------------------------------------------------------- THEME
DARK   = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0x2E, 0x6D, 0xB4)
LIGHT  = RGBColor(0xF4, 0xF6, 0xF9)
TEXT   = RGBColor(0x2D, 0x2D, 0x2D)
MUTED  = RGBColor(0x6B, 0x72, 0x80)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
AMBER  = RGBColor(0xE0, 0x8A, 0x00)
RED    = RGBColor(0xC0, 0x39, 0x2B)
RULE   = RGBColor(0xDD, 0xE3, 0xEA)
ACCENT_TINT = RGBColor(0xE9, 0xF0, 0xF8)
AMBER_TINT  = RGBColor(0xFB, 0xF1, 0xDD)
RED_TINT    = RGBColor(0xF9, 0xE7, 0xE5)
EMU_W, EMU_H = Inches(13.333), Inches(7.5)
FONT = "Segoe UI"
MONO = "Consolas"

# Decision type -> colour (used on the decoder slide and in every class table)
DEC_COLOR = {
    "Mandatory?": ACCENT,
    "Audited?":   GREEN,
    "Owner?":     RED,
    "Values?":    AMBER,
}

# --------------------------------------------------------------------------- DATA
# One dict per class. items = (num, attribute, decision, context, team_suggestion)
CLASSES = [
    dict(
        name="Business Application", table="cmdb_ci_business_app",
        owner="Todd Dierksheide", owner_role="CCB Class Manager - Business Application",
        items=[
            (1, "Application Category", "Mandatory?", "Choice; authoritative = EA / LeanIX",         "Mandatory = Y (core classification)"),
            (2, "Life Cycle Stage",     "Mandatory?", "Choice (CSDM stages); authoritative = EA",     "Mandatory = Y (lifecycle reporting)"),
            (3, "IT Application Owner",  "Mandatory?", "sys_user; via Data Certification",             "Mandatory = Y (ownership / audit)"),
            (4, "Business Owner",        "Mandatory?", "sys_user; via Data Certification",             "Mandatory = Y (Dec audit remediation)"),
        ],
    ),
    dict(
        name="Service Instance", table="(no dedicated CCB Class Manager mapped)",
        owner="none assigned", owner_role="CCB must assign a Certification Owner",
        warn="Ownerless class. Assign an owner (EA / Platform - e.g. Gaurav Parmar or Sonika Das) "
             "before items 5-6 can be ratified, or defer the class.",
        items=[
            (5, "Name",                  "Owner?",     "No dedicated CCB Class Mgr mapped",  "Assign class owner (EA / Platform)"),
            (6, "Environment",           "Owner?",     "Prod / QA / Dev / Test",             "Same owner as #5"),
            (7, "Service Owner (owned_by)", "Mandatory?", "sys_user",                        "Mandatory = Y"),
            (8, "ITSM Support Group",    "Mandatory?", "sys_user_group",                     "Mandatory = Y (governance)"),
            (9, "Change Approval Group", "Mandatory?", "sys_user_group",                     "Mandatory = Y (governance)"),
        ],
    ),
    dict(
        name="Computer", table="cmdb_ci_computer",
        owner="Monica Green / Paul Becker", owner_role="CCB Class Managers - Computers (Physical / Virtual)",
        items=[
            (10, "Assigned To",       "Mandatory?", "sys_user; Asset Management",        "- (PPL call)"),
            (11, "Operating System",  "Audited?",   "SCCM / Intune authoritative",       "Align to Dec audit scope"),
            (12, "MAC Address",       "Audited?",   "Endpoint connector / Discovery",    "Align to Dec audit scope"),
        ],
    ),
    dict(
        name="Server (parent)", table="cmdb_ci_server",
        owner="Ray Reuter", owner_role="CCB Class Manager - Servers & Storage CIs",
        items=[
            (13, "Operational Status", "Values?",  "Build / Operational / Non-Operational / Retired", "Confirm full list vs OOB choices"),
            (14, "IP Address",         "Audited?", "Discovery authoritative",              "Align to Dec audit scope"),
            (15, "MAC Address",        "Audited?", "Discovery authoritative",              "Align to Dec audit scope"),
            (16, "CPU Count",          "Audited?", "vCenter override (else Discovery)",    "Align to Dec audit scope"),
            (17, "RAM (MB)",           "Audited?", "vCenter override (else Discovery)",    "Align to Dec audit scope"),
        ],
    ),
    dict(
        name="Windows Server (child of Server)", table="cmdb_ci_win_server",
        owner="Ray Reuter", owner_role="CCB Class Manager - Servers & Storage CIs",
        items=[
            (18, "OS Version",          "Audited?",   "SCCM overrides Discovery",             "Align to Dec audit scope"),
            (19, "OS Domain",           "Audited?",   "SCCM overrides Discovery",             "Align to Dec audit scope"),
            (20, "OS Service Pack",     "Audited?",   "SCCM overrides Discovery",             "Align to Dec audit scope"),
            (21, "ITSM Support Group",  "Mandatory?", "'Wintel Infrastructure Ops'; not auto-writable", "Mandatory = Y (governance)"),
            (22, "Change Approval Group","Mandatory?","'Wintel CAB'; not auto-writable",      "Mandatory = Y (governance)"),
        ],
    ),
    dict(
        name="Linux Server (child of Server)", table="cmdb_ci_linux_server",
        owner="Ray Reuter", owner_role="CCB Class Manager - Servers & Storage CIs",
        items=[
            (23, "OS Architecture",     "Audited?",   "Discovery via SSH authoritative",      "Align to Dec audit scope"),
            (24, "Kernel Release",      "Audited?",   "Discovery via SSH authoritative",      "Align to Dec audit scope"),
            (25, "ITSM Support Group",  "Mandatory?", "'Unix / Linux Administration Team'",   "Mandatory = Y (governance)"),
            (26, "Change Approval Group","Mandatory?","'Linux CAB'",                          "Mandatory = Y (governance)"),
        ],
    ),
    dict(
        name="Database (parent)", table="cmdb_ci_database",
        owner="Ray Reuter", owner_role="CCB Class Manager - Database CIs",
        items=[
            (27, "Version",            "Audited?", "Discovery authoritative",   "Align to Dec audit scope"),
            (28, "TCP Port",           "Audited?", "Discovery authoritative",   "Align to Dec audit scope"),
            (29, "Operational Status", "Values?",  "Operational, Retired, ...", "Confirm full list"),
        ],
    ),
    dict(
        name="Oracle Database (child)", table="cmdb_ci_db_ora_instance",
        owner="Ray Reuter", owner_role="CCB Class Manager - Database CIs",
        items=[
            (30, "SID",            "Audited?",   "Oracle probes",       "Align to Dec audit scope"),
            (31, "Oracle Edition", "Audited?",   "Oracle probes",       "Align to Dec audit scope"),
            (32, "Oracle Version", "Audited?",   "Oracle probes",       "Align to Dec audit scope"),
            (33, "Support Group",  "Mandatory?", "'Oracle DBA Team'",   "Mandatory = Y (governance)"),
        ],
    ),
    dict(
        name="SQL Server Database (child)", table="cmdb_ci_db_mssql_instance",
        owner="Ray Reuter", owner_role="CCB Class Manager - Database CIs",
        items=[
            (34, "Instance Name", "Audited?",   "Discovery / SCCM",                  "Align to Dec audit scope"),
            (35, "Clustered",     "Audited?",   "Discovery authoritative",           "Align to Dec audit scope"),
            (36, "SQL Edition",   "Audited?",   "SCCM authoritative (edition / license)", "Align to Dec audit scope"),
            (37, "TCP Port",      "Audited?",   "Discovery overrides SCCM",          "Align to Dec audit scope"),
            (38, "Support Group", "Mandatory?", "'SQL DBA Team'",                    "Mandatory = Y (governance)"),
        ],
    ),
]

# --------------------------------------------------------------------------- HELPERS
def _set(run, size, color, bold=False, italic=False, name=FONT):
    run.font.size = Pt(size); run.font.color.rgb = color
    run.font.bold = bold; run.font.italic = italic; run.font.name = name


def rect(slide, x, y, w, h, color, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def textbox(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(4); tf.margin_top = tf.margin_bottom = Pt(2)
    for i, spec in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        text, size, color, bold, italic = spec
        r = p.add_run(); r.text = text; _set(r, size, color, bold, italic)
    return tb


def header(slide, title, subtitle=None):
    rect(slide, 0, 0, EMU_W, Inches(1.05), DARK)
    rect(slide, 0, Inches(1.05), EMU_W, Inches(0.06), ACCENT)
    textbox(slide, Inches(0.55), Inches(0.16), Inches(12.2), Inches(0.55),
            [(title, 23, WHITE, True, False)], MSO_ANCHOR.MIDDLE)
    if subtitle:
        textbox(slide, Inches(0.57), Inches(0.66), Inches(12.2), Inches(0.32),
                [(subtitle, 12, RGBColor(0xC9, 0xD6, 0xE6), False, False)], MSO_ANCHOR.MIDDLE)
    textbox(slide, Inches(0.55), Inches(7.08), Inches(10), Inches(0.32),
            [(FOOTLABEL, 9, MUTED, False, False)])


def cell(t, i, j, lines, size, color, bold=False, align=PP_ALIGN.LEFT,
         fill=None, anchor=None, name=FONT, italic=False):
    c = t.cell(i, j)
    if fill is not None:
        c.fill.solid(); c.fill.fore_color.rgb = fill
    if anchor is not None:
        c.vertical_anchor = anchor
    c.margin_left = c.margin_right = Pt(5); c.margin_top = c.margin_bottom = Pt(2)
    tf = c.text_frame; tf.word_wrap = True
    if isinstance(lines, str):
        lines = [lines]
    for k, ln in enumerate(lines):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.alignment = align; p.space_before = Pt(0); p.space_after = Pt(1)
        r = p.add_run(); r.text = ln; _set(r, size, color, bold, italic, name=name)


# --------------------------------------------------------------------------- SLIDE 1: TITLE
def add_title_slide(prs, blank):
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, EMU_W, EMU_H, DARK)
    rect(s, 0, Inches(4.55), EMU_W, Inches(0.06), ACCENT)
    textbox(s, Inches(0.9), Inches(2.55), Inches(11.5), Inches(1.1),
            [(TITLE, 40, WHITE, True, False)], MSO_ANCHOR.BOTTOM)
    textbox(s, Inches(0.92), Inches(4.75), Inches(11.5), Inches(0.5),
            [(SUBTITLE, 17, RGBColor(0xC9, 0xD6, 0xE6), False, False)])
    textbox(s, Inches(0.92), Inches(5.35), Inches(11.5), Inches(0.4),
            [("CMDB Change Control Board  -  Chair: Josh Sterling  -  July 21, 2026",
              12.5, RGBColor(0x9F, 0xB3, 0xCC), False, False)])
    textbox(s, Inches(0.92), Inches(6.7), Inches(11.5), Inches(0.4),
            [(AUTHOR, 11, RGBColor(0x9F, 0xB3, 0xCC), False, False)])


# --------------------------------------------------------------------------- SLIDE 2: AT A GLANCE
GLANCE_ROWS = [
    ("1. Business Application",      "Todd Dierksheide",          "4",  "Mandatory x4"),
    ("2. Service Instance",          "!! none assigned",          "5",  "Owner x2, Mandatory x3"),
    ("3. Computer",                  "Monica Green / Paul Becker","3",  "Mandatory x1, Audited x2"),
    ("4. Server (parent)",           "Ray Reuter",                "5",  "Values x1, Audited x4"),
    ("5. Windows Server",            "Ray Reuter",                "5",  "Audited x3, Mandatory x2"),
    ("6. Linux Server",              "Ray Reuter",                "4",  "Audited x2, Mandatory x2"),
    ("7. Database (parent)",         "Ray Reuter",                "3",  "Audited x2, Values x1"),
    ("8. Oracle Database",           "Ray Reuter",                "4",  "Audited x3, Mandatory x1"),
    ("9. SQL Server Database",       "Ray Reuter",                "5",  "Audited x4, Mandatory x1"),
]


def add_glance_slide(prs, blank):
    s = prs.slides.add_slide(blank)
    header(s, "At a Glance  -  38 Decisions, How to Vote",
           "Group the vote by class; the chair records each ruling in the CCB ruling column")

    # left: overview table
    n = len(GLANCE_ROWS) + 2  # header + rows + total
    t = s.shapes.add_table(n, 4, Inches(0.55), Inches(1.35), Inches(7.55), Inches(0.36 * n)).table
    for j, cw in enumerate([Inches(2.55), Inches(2.35), Inches(0.80), Inches(1.85)]):
        t.columns[j].width = cw
    for j, h in enumerate(["Class", "Class Manager", "#", "Decision types"]):
        cell(t, 0, j, h, 10, WHITE, True, PP_ALIGN.CENTER if j == 2 else PP_ALIGN.LEFT, DARK,
             anchor=MSO_ANCHOR.MIDDLE)
    for i, (cls, mgr, num, types) in enumerate(GLANCE_ROWS, start=1):
        fill = LIGHT if i % 2 else WHITE
        warn = mgr.startswith("!!")
        cell(t, i, 0, cls,  9,   DARK, True,  PP_ALIGN.LEFT,   fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 1, mgr,  8.5, RED if warn else TEXT, warn, PP_ALIGN.LEFT, fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 2, num,  9.5, ACCENT, True, PP_ALIGN.CENTER, fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 3, types,8,   MUTED, False, PP_ALIGN.LEFT,   fill, MSO_ANCHOR.MIDDLE)
    r = n - 1
    cell(t, r, 0, "Total", 10, WHITE, True, PP_ALIGN.LEFT, DARK, MSO_ANCHOR.MIDDLE)
    cell(t, r, 1, "",      10, WHITE, True, PP_ALIGN.LEFT, DARK, MSO_ANCHOR.MIDDLE)
    cell(t, r, 2, "38",    11, WHITE, True, PP_ALIGN.CENTER, DARK, MSO_ANCHOR.MIDDLE)
    cell(t, r, 3, "M14 / A20 / O2 / V2", 8.5, WHITE, True, PP_ALIGN.LEFT, DARK, MSO_ANCHOR.MIDDLE)
    for rr in range(n):
        t.rows[rr].height = Inches(0.36)

    # right column: how to vote + two flags
    rx, rw = 8.35, 4.40
    rect(s, Inches(rx), Inches(1.35), Inches(rw), Inches(0.34), ACCENT)
    textbox(s, Inches(rx + 0.1), Inches(1.35), Inches(rw - 0.2), Inches(0.34),
            [("How to vote each item", 11, WHITE, True, False)], MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(rx + 0.05), Inches(1.76), Inches(rw), Inches(0.9),
            [("Accept  -  as the team suggested", 10.5, GREEN, True, False),
             ("Amend  -  note the change in the ruling cell", 10.5, AMBER, True, False),
             ("Defer  -  record owner + target date", 10.5, RED, True, False),
             ("A class is accepted once all its items are ruled.", 9, MUTED, False, True)])

    # flag 1
    rect(s, Inches(rx), Inches(2.86), Inches(rw), Inches(1.12), RED_TINT, line=RED)
    rect(s, Inches(rx), Inches(2.86), Inches(0.08), Inches(1.12), RED)
    textbox(s, Inches(rx + 0.18), Inches(2.90), Inches(rw - 0.28), Inches(1.06),
            [("!  Resolve first - Service Instance", 10.5, RED, True, False),
             ("Class has no manager (items 5-6). Assign an owner "
              "(EA / Platform) or defer - items can't be ratified without one.",
              9, TEXT, False, False)])

    # flag 2
    rect(s, Inches(rx), Inches(4.10), Inches(rw), Inches(1.12), AMBER_TINT, line=AMBER)
    rect(s, Inches(rx), Inches(4.10), Inches(0.08), Inches(1.12), AMBER)
    textbox(s, Inches(rx + 0.18), Inches(4.14), Inches(rw - 0.28), Inches(1.06),
            [("!  Sequence - Ray Reuter load", 10.5, AMBER, True, False),
             ("Ray owns 26 of 38 (all Server + Database). Take Server / "
              "Database first so they aren't stranded if time runs short.",
              9, TEXT, False, False)])

    # bottom caption
    textbox(s, Inches(0.55), Inches(6.6), Inches(12.2), Inches(0.3),
            [(CAPTION, 9, MUTED, False, True)])


# --------------------------------------------------------------------------- SLIDE 3: DECISION-TYPE DECODER
DECODER = [
    ("Mandatory?", "14 decisions", ACCENT, ACCENT_TINT,
     "Is this attribute REQUIRED on the class?",
     "Y = a record can't be certified without it. Sets the governance / certification bar."),
    ("Audited?", "20 decisions", GREEN, RGBColor(0xE7, 0xF2, 0xEC),
     "Include this attribute in AUDIT / certification scope?",
     "Y = it counts toward the 90% completeness / correctness target and appears on audit dashboards."),
    ("Owner?", "2 decisions", RED, RED_TINT,
     "Assign a Certification Owner / Class Manager.",
     "Only Service Instance (5-6), currently unassigned. Name an owner (EA / Platform) or defer the class."),
    ("Values?", "2 decisions", AMBER, AMBER_TINT,
     "Confirm the allowed-value (choice) list.",
     "e.g. Operational Status: Build / Operational / Non-Operational / Retired. Confirm vs the OOB choices."),
]


def add_decoder_slide(prs, blank):
    s = prs.slides.add_slide(blank)
    header(s, "Education  -  The Four Decision Types",
           "Every one of the 38 items is one of these four. Rule each Accept / Amend / Defer.")
    x0, y0, cw, ch, gap = 0.55, 1.45, 6.02, 2.42, 0.16
    coords = [(x0, y0), (x0 + cw + gap, y0),
              (x0, y0 + ch + gap), (x0 + cw + gap, y0 + ch + gap)]
    for (dec, count, col, tint, question, answer), (x, y) in zip(DECODER, coords):
        rect(s, Inches(x), Inches(y), Inches(cw), Inches(ch), tint, line=col)
        rect(s, Inches(x), Inches(y), Inches(0.10), Inches(ch), col)
        textbox(s, Inches(x + 0.25), Inches(y + 0.14), Inches(cw - 1.6), Inches(0.5),
                [(dec, 18, col, True, False)])
        rect(s, Inches(x + cw - 1.55), Inches(y + 0.20), Inches(1.32), Inches(0.36), col)
        textbox(s, Inches(x + cw - 1.55), Inches(y + 0.20), Inches(1.32), Inches(0.36),
                [(count, 10, WHITE, True, False)], MSO_ANCHOR.MIDDLE, PP_ALIGN.CENTER)
        textbox(s, Inches(x + 0.25), Inches(y + 0.78), Inches(cw - 0.5), Inches(0.7),
                [(question, 12.5, DARK, True, False)])
        textbox(s, Inches(x + 0.25), Inches(y + 1.48), Inches(cw - 0.5), Inches(0.85),
                [(answer, 10.5, TEXT, False, False)])
    textbox(s, Inches(0.55), Inches(6.62), Inches(12.2), Inches(0.3),
            [("Capture each as Accept (as suggested) / Amend (note the change) / Defer (owner + date).",
              9.5, MUTED, False, True)])


# --------------------------------------------------------------------------- SLIDE 4: IRE PRIMER
PRECEDENCE = [
    ("SG-SCCM wins",             "OS & OS version, CPU / disk / RAM, manufacturer, asset & ownership specs", GREEN),
    ("ServiceNow Discovery wins","IP address, MAC address (network-sourced facts)",                          ACCENT),
    ("EA / LeanIX wins",         "Application Category, Life Cycle Stage (Business App)",                     AMBER),
    ("Manual / Data Cert wins",  "Owners, support & approval groups, SOX flag (NOT discovered)",              RED),
]


def add_ire_slide(prs, blank):
    s = prs.slides.add_slide(blank)
    header(s, "Education  -  IRE in 60 Seconds",
           "Why some attributes are 'discovery-sourced' - and why that decides how we audit them")

    # two concept cards
    cx, cw = 0.55, 5.95
    rect(s, Inches(cx), Inches(1.40), Inches(cw), Inches(1.75), ACCENT_TINT, line=ACCENT)
    textbox(s, Inches(cx + 0.2), Inches(1.50), Inches(cw - 0.4), Inches(0.4),
            [("1.  Identification", 14, ACCENT, True, False)])
    textbox(s, Inches(cx + 0.2), Inches(1.92), Inches(cw - 0.4), Inches(1.15),
            [("How ServiceNow decides two records are the SAME CI - via identifier "
              "rules (e.g. Server = Serial Number; Computer = Name + Serial). ", 10.5, TEXT, False, False),
             ("Independent = identified on its own;  Dependent = only in the context "
              "of a host (e.g. a Database on its Server).", 10, MUTED, False, True)])

    x2 = cx + cw + 0.20
    rect(s, Inches(x2), Inches(1.40), Inches(cw), Inches(1.75), RGBColor(0xE7, 0xF2, 0xEC), line=GREEN)
    textbox(s, Inches(x2 + 0.2), Inches(1.50), Inches(cw - 0.4), Inches(0.4),
            [("2.  Reconciliation", 14, GREEN, True, False)])
    textbox(s, Inches(x2 + 0.2), Inches(1.92), Inches(cw - 0.4), Inches(1.15),
            [("When MORE THAN ONE source writes the same attribute, precedence decides "
              "which source wins - so the value doesn't 'flap'.", 10.5, TEXT, False, False),
             ("The winning source is authoritative; other sources are ignored for that "
              "attribute.", 10, MUTED, False, True)])

    # precedence cheat-sheet
    textbox(s, Inches(0.55), Inches(3.35), Inches(12.2), Inches(0.3),
            [("PPL source precedence (as-built, 2026-07-20)", 13, DARK, True, False)])
    t = s.shapes.add_table(len(PRECEDENCE), 2, Inches(0.55), Inches(3.72),
                           Inches(12.2), Inches(0.52 * len(PRECEDENCE))).table
    t.columns[0].width = Inches(3.05); t.columns[1].width = Inches(9.15)
    for i, (src, attrs, col) in enumerate(PRECEDENCE):
        cell(t, i, 0, src,   11, WHITE, True, PP_ALIGN.LEFT, col, MSO_ANCHOR.MIDDLE)
        cell(t, i, 1, attrs, 10.5, TEXT, False, PP_ALIGN.LEFT, LIGHT if i % 2 else WHITE, MSO_ANCHOR.MIDDLE)
    for rr in range(len(PRECEDENCE)):
        t.rows[rr].height = Inches(0.52)

    # why it matters
    rect(s, Inches(0.55), Inches(6.02), Inches(12.2), Inches(0.92), ACCENT_TINT, line=ACCENT)
    rect(s, Inches(0.55), Inches(6.02), Inches(0.08), Inches(0.92), ACCENT)
    textbox(s, Inches(0.75), Inches(6.05), Inches(1.4), Inches(0.9),
            [("Why it", 11, DARK, True, False), ("matters", 11, DARK, True, False)], MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(2.15), Inches(6.08), Inches(10.4), Inches(0.82),
            [("If a source is authoritative for an attribute, 'Audited = Y' means we measure the "
              "DISCOVERED value - we don't ask humans to maintain it.", 10.5, TEXT, False, False),
             ("The ownership / governance fields (owners, groups, SOX) are NOT discovered - those "
              "are the ones Data Certification must fill, so they carry the Mandatory decisions.",
              10, MUTED, False, True)])


# --------------------------------------------------------------------------- CLASS VOTE SLIDES
def add_class_slide(prs, blank, c):
    s = prs.slides.add_slide(blank)
    header(s, "Vote  -  %s" % c["name"],
           "Class Manager: %s   -   Table: %s   -   %d decision(s)"
           % (c["owner"], c["table"], len(c["items"])))

    # owner chip (red if unassigned)
    unassigned = c["owner"] == "none assigned"
    chip = RED if unassigned else ACCENT
    rect(s, Inches(10.75), Inches(0.30), Inches(2.02), Inches(0.44), chip)
    textbox(s, Inches(10.75), Inches(0.30), Inches(2.02), Inches(0.44),
            [(("ASSIGN OWNER" if unassigned else "PROPOSED - PENDING CCB"), 9, WHITE, True, False)],
            MSO_ANCHOR.MIDDLE, PP_ALIGN.CENTER)

    top = 1.35
    # optional warning band
    if c.get("warn"):
        rect(s, Inches(0.55), Inches(top), Inches(12.2), Inches(0.62), RED_TINT, line=RED)
        rect(s, Inches(0.55), Inches(top), Inches(0.08), Inches(0.62), RED)
        textbox(s, Inches(0.75), Inches(top + 0.02), Inches(12.0), Inches(0.58),
                [("!  " + c["warn"], 10, RED, True, False)], MSO_ANCHOR.MIDDLE)
        top += 0.78

    items = c["items"]
    rows = len(items) + 1
    widths = [Inches(0.42), Inches(1.85), Inches(1.05), Inches(3.83), Inches(3.20), Inches(1.85)]
    tbl_h = 0.34 + 0.58 * len(items)
    t = s.shapes.add_table(rows, 6, Inches(0.55), Inches(top), Inches(12.2), Inches(tbl_h)).table
    for j, cw in enumerate(widths):
        t.columns[j].width = cw
    for j, h in enumerate(["#", "Attribute", "Decision", "Context", "Team suggestion", "CCB ruling"]):
        cell(t, 0, j, h, 10, WHITE, True,
             PP_ALIGN.CENTER if j in (0, 2, 5) else PP_ALIGN.LEFT, DARK, MSO_ANCHOR.MIDDLE)
    for i, (num, attr, dec, ctx, sug) in enumerate(items, start=1):
        fill = LIGHT if i % 2 else WHITE
        cell(t, i, 0, str(num), 9,   MUTED, True,  PP_ALIGN.CENTER, fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 1, attr,     9.5, DARK,  True,  PP_ALIGN.LEFT,   fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 2, dec,      9,   DEC_COLOR.get(dec, TEXT), True, PP_ALIGN.CENTER, fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 3, ctx,      8.5, TEXT,  False, PP_ALIGN.LEFT,   fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 4, sug,      8.5, TEXT,  False, PP_ALIGN.LEFT,   fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 5, "Accept / Amend / Defer", 8, MUTED, False, PP_ALIGN.CENTER,
             RGBColor(0xFC, 0xFD, 0xFE), MSO_ANCHOR.MIDDLE, italic=True)
    t.rows[0].height = Inches(0.34)
    for rr in range(1, rows):
        t.rows[rr].height = Inches(0.58)

    textbox(s, Inches(0.55), Inches(6.95), Inches(12.2), Inches(0.28),
            [(CAPTION, 8, MUTED, False, True)])


# --------------------------------------------------------------------------- SLIDE: AGENDA
AGENDA_ITEMS = [
    ("Welcome & Purpose",         "5 min",   ACCENT,
     "Josh Sterling (CCB Chair) opens — why we are here and what ratification means"),
    ("Education Block",           "15 min",  GREEN,
     "Data Dictionary · CMP context · CCB role · Decision types · IRE primer"),
    ("Live Vote — 38 decisions",  "40 min",  DARK,
     "Class-by-class: Business App → Service Instance → Computer → Server family → Database family\n"
     "Each item: Accept / Amend / Defer  –  Chair records the ruling live"),
    ("Wrap-up & Next Steps",      "5 min",   AMBER,
     "Confirm deferrals, assign owners, set target dates  –  updated Data Dictionary distributed within 5 business days"),
]


def add_agenda_slide(prs, blank):
    s = prs.slides.add_slide(blank)
    header(s, "Meeting Agenda  -  CMDB CCB  |  July 21, 2026",
           "Chair: Josh Sterling    Presenter: Manuel Vazquez    Duration: ~65 min")
    y = 1.38; gap = 0.12
    bh = 1.14
    for label, duration, col, body in AGENDA_ITEMS:
        rect(s, Inches(0.55), Inches(y), Inches(12.2), Inches(bh), LIGHT, line=col)
        rect(s, Inches(0.55), Inches(y), Inches(0.10), Inches(bh), col)
        # duration chip
        rect(s, Inches(11.65), Inches(y + 0.12), Inches(0.95), Inches(0.38), col)
        textbox(s, Inches(11.65), Inches(y + 0.12), Inches(0.95), Inches(0.38),
                [(duration, 10, WHITE, True, False)], MSO_ANCHOR.MIDDLE, PP_ALIGN.CENTER)
        textbox(s, Inches(0.80), Inches(y + 0.10), Inches(10.75), Inches(0.38),
                [(label, 13, col, True, False)])
        textbox(s, Inches(0.80), Inches(y + 0.52), Inches(10.75), Inches(0.55),
                [(body, 9.5, TEXT, False, False)])
        y += bh + gap


# --------------------------------------------------------------------------- SLIDE: WHAT IS A DATA DICTIONARY
def add_what_is_slide(prs, blank):
    s = prs.slides.add_slide(blank)
    header(s, "What Is a Data Dictionary?",
           "Education  |  The formal register that governs every attribute stored in the CMDB")

    # left card — definition
    lx, ly, lw, lh = 0.55, 1.40, 5.95, 4.90
    rect(s, Inches(lx), Inches(ly), Inches(lw), Inches(lh), ACCENT_TINT, line=ACCENT)
    rect(s, Inches(lx), Inches(ly), Inches(0.10), Inches(lh), ACCENT)
    textbox(s, Inches(lx + 0.25), Inches(ly + 0.15), Inches(lw - 0.35), Inches(0.42),
            [("Definition", 15, ACCENT, True, False)])
    DEFN = (
        "A Data Dictionary is a formal register that defines every attribute stored "
        "in the CMDB for each CI class. For each attribute it records:"
    )
    textbox(s, Inches(lx + 0.25), Inches(ly + 0.65), Inches(lw - 0.35), Inches(0.55),
            [(DEFN, 10, TEXT, False, False)])
    ROWS = [
        ("Attribute name",    "What it is called — display name + technical field name"),
        ("Data type",         "String, Choice, Reference, True/False, etc."),
        ("Mandatory?",        "Is a record incomplete without it?"),
        ("Allowed values",    "Permitted choices or reference table"),
        ("Source of truth",   "Which system is authoritative — Discovery, SCCM, EA, Manual"),
        ("Audited?",          "Does it count toward the 90% completeness/correctness target?"),
        ("Certification owner", "Who is accountable for data quality in this field?"),
    ]
    t = s.shapes.add_table(len(ROWS), 2, Inches(lx + 0.25), Inches(ly + 1.28),
                           Inches(lw - 0.35), Inches(0.48 * len(ROWS))).table
    t.columns[0].width = Inches(1.80); t.columns[1].width = Inches(3.80)
    for i, (name, desc) in enumerate(ROWS):
        fill = LIGHT if i % 2 else WHITE
        cell(t, i, 0, name, 9, ACCENT, True,  PP_ALIGN.LEFT, fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 1, desc, 9, TEXT,   False, PP_ALIGN.LEFT, fill, MSO_ANCHOR.MIDDLE)
        t.rows[i].height = Inches(0.48)

    # right card — why it matters
    rx, ry, rw, rh = 6.70, 1.40, 6.08, 4.90
    rect(s, Inches(rx), Inches(ry), Inches(rw), Inches(rh), RGBColor(0xE7, 0xF2, 0xEC), line=GREEN)
    rect(s, Inches(rx), Inches(ry), Inches(0.10), Inches(rh), GREEN)
    textbox(s, Inches(rx + 0.25), Inches(ry + 0.15), Inches(rw - 0.35), Inches(0.42),
            [("Why it matters", 15, GREEN, True, False)])
    WHY = [
        ("Without agreed definitions, three things break down:", 10.5, DARK,  True,  False),
        ("",                                                      5,    TEXT,  False, False),
        ("Audits can't be scoped  —  auditors measure different things than the "
         "CMDB team tracks. Numbers diverge and no one trusts the data.",
         10, TEXT, False, False),
        ("",                                                      5,    TEXT,  False, False),
        ("Discovery conflicts can't be resolved  —  when SCCM and ServiceNow "
         "disagree on an attribute value, there is no agreed rule for which source wins.",
         10, TEXT, False, False),
        ("",                                                      5,    TEXT,  False, False),
        ("Data quality can't be measured  —  a completeness percentage means nothing "
         "if the team hasn't agreed which attributes are required.",
         10, TEXT, False, False),
        ("",                                                      7,    TEXT,  False, False),
        ("The Data Dictionary fixes all three:  it is the single agreed-upon record "
         "of what the CMDB should contain and who is accountable for each piece.",
         10.5, DARK, True, False),
    ]
    textbox(s, Inches(rx + 0.25), Inches(ry + 0.68), Inches(rw - 0.35), Inches(4.05), WHY)


# --------------------------------------------------------------------------- SLIDE: HOW IT FITS IN THE CMP
CMP_STAGES = [
    ("Stage 1",  "Class & Attribute\nData Definitions",
     "The Data Dictionary voted on today.\nDefines every attribute, its source of truth, "
     "mandatory/audited flags, and certification owner for each CI class.",
     True, ACCENT),
    ("Stage 2",  "CI Lifecycle &\nHealth Targets",
     "Defines what 'healthy' looks like — completeness / correctness thresholds, "
     "staleness rules, and the lifecycle states a CI moves through.",
     False, GREEN),
    ("Stage 3",  "Governance Process\n& Controls",
     "Formalizes the ongoing CCB cadence, change approval flows, audit cadence, "
     "and escalation paths. Gated by CO6 Governance milestone (Sep 30, 2026).",
     False, AMBER),
]


def add_cmp_slide(prs, blank):
    s = prs.slides.add_slide(blank)
    header(s, "How the Data Dictionary Fits in the CMP",
           "Education  |  The Configuration Management Plan is the governance umbrella  —  Stage 1 is today's vote")

    # intro sentence
    textbox(s, Inches(0.55), Inches(1.35), Inches(12.2), Inches(0.42),
            [("The Configuration Management Plan (CMP) defines how PPL's CMDB is built, maintained, and governed. "
              "It has three stages. Today ratifies Stage 1 — the foundation everything else depends on.",
              11, TEXT, False, False)])

    # stage cards
    for i, (stage, title, body, active, col) in enumerate(CMP_STAGES):
        x = 0.55 + i * 4.22; cw = 4.02; cy = 1.95; ch = 3.90
        tint = ACCENT_TINT if col == ACCENT else (RGBColor(0xE7, 0xF2, 0xEC) if col == GREEN else AMBER_TINT)
        rect(s, Inches(x), Inches(cy), Inches(cw), Inches(ch), tint, line=col)
        rect(s, Inches(x), Inches(cy), Inches(cw), Inches(0.10), col)
        # active badge
        if active:
            rect(s, Inches(x + cw - 1.35), Inches(cy + 0.16), Inches(1.22), Inches(0.34), col)
            textbox(s, Inches(x + cw - 1.35), Inches(cy + 0.16), Inches(1.22), Inches(0.34),
                    [("TODAY", 9, WHITE, True, False)], MSO_ANCHOR.MIDDLE, PP_ALIGN.CENTER)
        textbox(s, Inches(x + 0.20), Inches(cy + 0.18), Inches(cw - 1.55 if active else cw - 0.40), Inches(0.34),
                [(stage, 13, col, True, False)])
        textbox(s, Inches(x + 0.20), Inches(cy + 0.60), Inches(cw - 0.40), Inches(0.62),
                [(title, 12, DARK, True, False)])
        textbox(s, Inches(x + 0.20), Inches(cy + 1.28), Inches(cw - 0.40), Inches(2.50),
                [(body, 9.5, TEXT, False, False)])

    # bottom callout — why stage 1 is the foundation
    rect(s, Inches(0.55), Inches(6.00), Inches(12.2), Inches(0.84), ACCENT_TINT, line=ACCENT)
    rect(s, Inches(0.55), Inches(6.00), Inches(0.10), Inches(0.84), ACCENT)
    textbox(s, Inches(0.80), Inches(6.04), Inches(1.25), Inches(0.76),
            [("Why Stage 1", 10, DARK, True, False), ("first", 10, DARK, True, False)], MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(2.10), Inches(6.04), Inches(10.45), Inches(0.76),
            [("Audit targets, certification scope, and remediation priorities are all defined in terms of "
              "the Data Dictionary.  Stages 2 and 3 cannot be written until Stage 1 is ratified — "
              "the CCB's rulings today become the source of truth for all downstream governance work.",
              9.5, TEXT, False, False)], MSO_ANCHOR.MIDDLE)


# --------------------------------------------------------------------------- SLIDE: ROLE OF THE CCB
CCB_ROLES = [
    ("Ratify attribute governance",
     "For each attribute the CCB decides whether it is Mandatory, whether it is in Audit scope, "
     "what values are allowed, and who is the Certification Owner. These rulings are binding — "
     "they become the baseline for all future data quality measurement.",
     ACCENT),
    ("Assign Class Managers",
     "Each CI class must have a named Class Manager who is accountable for data quality "
     "after ratification. The CCB confirms or assigns the Class Manager for each class today. "
     "Service Instance currently has no assigned manager — this must be resolved in the meeting.",
     RED),
    ("Record and enforce decisions",
     "The Chair (Josh Sterling) captures each ruling live. Accept = adopt as proposed. "
     "Amend = record the change. Defer = name an owner and a target date. "
     "Deferred items are tracked to closure; the CCB does not re-vote without new information.",
     GREEN),
    ("Govern ongoing changes",
     "Post-ratification, any change to a mandatory attribute, audit flag, allowed-value list, "
     "or class ownership requires a CCB agenda item. The CCB is the standing governance body "
     "for the CMDB data model — not a one-time event.",
     AMBER),
]


def add_ccb_role_slide(prs, blank):
    s = prs.slides.add_slide(blank)
    header(s, "Role of the CCB",
           "Education  |  The CCB is the governance authority — it ratifies, assigns accountability, and enforces")

    y = 1.38; bh = 1.18; gap = 0.10
    for label, body, col in CCB_ROLES:
        tint = (ACCENT_TINT if col == ACCENT
                else RED_TINT if col == RED
                else RGBColor(0xE7, 0xF2, 0xEC) if col == GREEN
                else AMBER_TINT)
        rect(s, Inches(0.55), Inches(y), Inches(12.2), Inches(bh), tint, line=col)
        rect(s, Inches(0.55), Inches(y), Inches(0.10), Inches(bh), col)
        textbox(s, Inches(0.80), Inches(y + 0.10), Inches(12.0), Inches(0.38),
                [(label, 12, col, True, False)])
        textbox(s, Inches(0.80), Inches(y + 0.52), Inches(12.0), Inches(0.58),
                [(body, 9.5, TEXT, False, False)])
        y += bh + gap

    # CCB membership note
    textbox(s, Inches(0.55), Inches(6.72), Inches(12.2), Inches(0.52),
            [("CCB Chair: Josh Sterling (Configuration Management Process Owner)  |  "
              "Class Managers: Todd Dierksheide, Ray Reuter, Monica Green, Paul Becker, "
              "Mark Smith, Adam Gross, Jason Finn, Cheryl Guia  |  "
              "Consulting: Cybersecurity, EA, ServiceNow Platform, ITSM Process",
              8.5, MUTED, False, True)])


# --------------------------------------------------------------------------- BUILD
def build_deck():
    prs = Presentation(); prs.slide_width = EMU_W; prs.slide_height = EMU_H
    blank = prs.slide_layouts[6]
    add_title_slide(prs, blank)
    add_agenda_slide(prs, blank)
    add_what_is_slide(prs, blank)
    add_cmp_slide(prs, blank)
    add_ccb_role_slide(prs, blank)
    add_glance_slide(prs, blank)
    add_decoder_slide(prs, blank)
    add_ire_slide(prs, blank)
    for c in CLASSES:
        add_class_slide(prs, blank, c)
    return prs


def save_safe(prs, path):
    try:
        prs.save(path); print("Saved", os.path.basename(path), "-", len(prs.slides._sldIdLst), "slide(s)")
    except PermissionError:
        base, ext = os.path.splitext(path)
        alt = base + "-new" + ext
        prs.save(alt); print("Target locked - saved to", os.path.basename(alt))


here = os.path.dirname(os.path.abspath(__file__))
save_safe(build_deck(), os.path.join(here, "CMDB-Data-Dictionary-CCB-Vote-" + EDITION + ".pptx"))
