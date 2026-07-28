"""
SUPERSEDED for class content (reconciliation note, 2026-07-27).
The deck actually DELIVERED to PPL - "CMDB-Data-Dictionary-CCB-2026 (July 21 2026)
- DELIVERED to PPL.pptx" - kept THIS deck's title / agenda / education / wrap-up
framing, but REPLACED the 4 class "decision-table" slides with the CI Class
Reference "Managed Attributes" slides (Required/Recommended/Mandatory designations).
In the delivered Server slide, `environment` is Mandatory - which is NOT present in
this script's Server-parent rows below. Do not treat this script's class rows as the
delivered dictionary. Authoritative content mirror: cmdb-health-completeness-correctness-stories.md.

CMDB Data Dictionary - CCB Review & Approval Deck.

Built for the 7/21/2026 CMDB CCB. This is THE meeting deck (it replaces the
earlier granular 38-decision vote deck). It follows the agreed agenda and is
scoped to the 4 core PARENT classes only:

  1  Title
  2  Agenda & Purpose        (purpose + 15-decision scope + out-of-scope callout)
  3  CMP & CCB               (what each is + how they work together)
  4  What is a Data Dictionary (definition + how it fits in the CMP = Stage 1)
  5  Review - Business Application
  6  Review - Computer
  7  Review - Server (parent)
  8  Review - Database (parent)
  9  Wrap-up & Next Steps

SCOPE: 4 parent classes = 15 of the 38 Data Dictionary decisions
(BA 4, Computer 3, Server-parent 5, Database-parent 3). Service Instance and the
child classes (Windows/Linux Server, Oracle/SQL Database) are NOT covered today.

Sources of truth:
  configuration-management-plan-stage1.md  - full attribute-level data dictionary
  data-dictionary-ccb-decision-sheet.md    - the open [CCB confirm] decisions

"Team suggestion" is the delivery team's proposed default (NOT a decision);
the CCB ruling is authoritative and captured live on 7/21 (Accept / Amend / Defer).

Theme mirrors the CMDB-CSDM deck family so it reads as a set.

Run:     python build_ccb_review_deck.py
Output:  CMDB-Data-Dictionary-CCB-<EDITION>.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

EDITION   = "2026-07-21"
TITLE     = "CMDB Data Dictionary - CCB Review & Approval"
SUBTITLE  = "CMP Stage 1  -  4 core CI classes  -  15 attribute-governance decisions"
AUTHOR    = "Manuel Vazquez  -  Scrum Master, CMDB-CSDM"
FOOTLABEL = "PPL CMDB-CSDM  |  Data Dictionary CCB Review & Approval (CMP Stage 1)  |  Jul 21, 2026"
CAPTION   = ("Amber = CCB ruling needed today.  'Team suggestion' is the delivery team's proposed default "
             "(NOT a decision).  Chair records Accept / Amend / Defer per item.")

# --------------------------------------------------------------------------- THEME
DARK   = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0x2E, 0x6D, 0xB4)
LIGHT  = RGBColor(0xF4, 0xF6, 0xF9)
TEXT   = RGBColor(0x2D, 0x2D, 0x2D)
MUTED  = RGBColor(0x6B, 0x72, 0x80)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
AMBER  = RGBColor(0xE0, 0x8A, 0x00)
RED    = RGBColor(0xC0, 0x39, 0x2B)
RULE   = RGBColor(0xDD, 0xE3, 0xEA)
ACCENT_TINT = RGBColor(0xE9, 0xF0, 0xF8)
ACCENT_FILL = RGBColor(0xDC, 0xE8, 0xF6)   # slightly deeper than ACCENT_TINT, for the active-stage row
AMBER_TINT  = RGBColor(0xFB, 0xF1, 0xDD)
RED_TINT    = RGBColor(0xF9, 0xE7, 0xE5)
GREEN_TINT  = RGBColor(0xE7, 0xF2, 0xEC)
EMU_W, EMU_H = Inches(13.333), Inches(7.5)
FONT = "Segoe UI"

# --------------------------------------------------------------------------- DATA
# One dict per PARENT class. rows = (attr, dtype, source, mandatory, audited, decision, suggestion)
#   mandatory / audited: "Y" | "N" | "CCB?"   (CCB? = open decision on this cell)
#   decision: None | "Mandatory?" | "Audited?" | "Values?"
#   suggestion: team's proposed default (blank for settled rows)
CLASSES = [
    dict(
        name="Business Application", table="cmdb_ci_business_app",
        owner="Todd Dierksheide  (CCB Class Manager)",
        rows=[
            ("Name",                 "String(255)",    "EA (LeanIX) / Manual",           "Y",    "Y",    None,         ""),
            ("Application Category", "Choice",         "EA integration (authoritative)", "CCB?", "Y",    "Mandatory?", "Y  -  core classification"),
            ("Life Cycle Stage",     "Choice",         "EA integration (authoritative)", "CCB?", "Y",    "Mandatory?", "Y  -  lifecycle reporting"),
            ("IT Application Owner", "Ref (sys_user)", "Manual / Data Certification",    "CCB?", "Y",    "Mandatory?", "Y  -  ownership / audit"),
            ("Business Owner",       "Ref (sys_user)", "Manual / Data Certification",    "CCB?", "Y",    "Mandatory?", "Y  -  Dec audit remediation"),
        ],
        note=("Name is the identity / IRE key (mandatory + audited - settled).  "
              "Plus 3 PPL custom attributes added 7/20 (Access Control Source System, Approval Group, "
              "Recovery Tier) appear in the full CMP but are not in today's decision scope."),
    ),
    dict(
        name="Computer", table="cmdb_ci_computer",
        owner="Monica Green (Physical) / Paul Becker (Virtual)  (CCB Class Managers)",
        rows=[
            ("Name",             "String(255)",    "SCCM/Intune SGC / ACC",           "Y",    "Y",    None,         ""),
            ("Serial Number",    "String",         "Asset Mgmt / SCCM (authoritative)","Y",   "Y",    None,         ""),
            ("Assigned To",      "Ref (sys_user)", "Asset Mgmt / Manual",             "CCB?", "Y",    "Mandatory?", "PPL call"),
            ("Operating System", "String",         "SCCM/Intune (authoritative)",     "N",    "CCB?", "Audited?",   "Align to Dec audit scope"),
            ("MAC Address",      "String",         "Endpoint connector / Discovery",  "N",    "CCB?", "Audited?",   "Align to Dec audit scope"),
        ],
        note=("Name + Serial Number are the identity / IRE keys (mandatory + audited - settled).  "
              "Plus 2 PPL custom attributes (Requested Item, Off Network) appear in the full CMP, "
              "not in today's scope."),
    ),
    dict(
        name="Server  (parent class)", table="cmdb_ci_server",
        owner="Ray Reuter  (CCB Class Manager - Servers & Storage)",
        rows=[
            ("Name",               "String(255)",             "ServiceNow Discovery",             "Y", "Y",    None,       ""),
            ("Serial Number",      "String",                  "ServiceNow Discovery",             "Y", "Y",    None,       ""),
            ("Operational Status", "Choice - confirm values", "ServiceNow Discovery / Manual",    "Y", "Y",    "Values?",  "Confirm allowed list"),
            ("IP Address",         "String",                  "ServiceNow Discovery (auth)",      "N", "CCB?", "Audited?", "Align to Dec audit scope"),
            ("MAC Address",        "String",                  "ServiceNow Discovery (auth)",      "N", "CCB?", "Audited?", "Align to Dec audit scope"),
            ("CPU Count",          "Integer",                 "vCenter override (else Discovery)","N", "CCB?", "Audited?", "Align to Dec audit scope"),
            ("RAM (MB)",           "Integer",                 "vCenter override (else Discovery)","N", "CCB?", "Audited?", "Align to Dec audit scope"),
        ],
        note=("Operational Status (Values?): confirm allowed list - Build / Operational / Non-Operational / "
              "Retired (vs. OOB choices).  Parent Server class only; Windows/Linux child classes and the "
              "Hardware Type custom attribute are not covered today."),
    ),
    dict(
        name="Database  (parent class)", table="cmdb_ci_database",
        owner="Ray Reuter  (CCB Class Manager - Database CIs)",
        rows=[
            ("Name",               "String(255)",             "ServiceNow Discovery (App Probes)", "Y", "Y",    None,       ""),
            ("Version",            "String",                  "ServiceNow Discovery (auth)",       "N", "CCB?", "Audited?", "Align to Dec audit scope"),
            ("TCP Port",           "Integer",                 "ServiceNow Discovery (auth)",       "N", "CCB?", "Audited?", "Align to Dec audit scope"),
            ("Operational Status", "Choice - confirm values", "ServiceNow Discovery / Manual",     "Y", "Y",    "Values?",  "Confirm allowed list"),
        ],
        note=("Operational Status (Values?): confirm allowed list - Operational / Retired / ... (confirm full "
              "list).  Parent Database class only; Oracle/SQL child classes are not covered today."),
    ),
]

# --------------------------------------------------------------------------- HELPERS
def _set(run, size, color, bold=False, italic=False, name=FONT):
    run.font.size = Pt(size); run.font.color.rgb = color
    run.font.bold = bold; run.font.italic = italic; run.font.name = name


def rect(slide, x, y, w, h, color, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def textbox(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(4); tf.margin_top = tf.margin_bottom = Pt(2)
    for i, spec in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        text, size, color, bold, italic = spec
        r = p.add_run(); r.text = text; _set(r, size, color, bold, italic)
    return tb


def header(slide, title, subtitle=None):
    rect(slide, 0, 0, EMU_W, Inches(1.05), DARK)
    rect(slide, 0, Inches(1.05), EMU_W, Inches(0.06), ACCENT)
    textbox(slide, Inches(0.55), Inches(0.16), Inches(12.2), Inches(0.55),
            [(title, 23, WHITE, True, False)], MSO_ANCHOR.MIDDLE)
    if subtitle:
        textbox(slide, Inches(0.57), Inches(0.66), Inches(12.2), Inches(0.32),
                [(subtitle, 12, RGBColor(0xC9, 0xD6, 0xE6), False, False)], MSO_ANCHOR.MIDDLE)
    textbox(slide, Inches(0.55), Inches(7.08), Inches(10), Inches(0.32),
            [(FOOTLABEL, 9, MUTED, False, False)])


def cell(t, i, j, lines, size, color, bold=False, align=PP_ALIGN.LEFT,
         fill=None, anchor=None, name=FONT, italic=False):
    c = t.cell(i, j)
    if fill is not None:
        c.fill.solid(); c.fill.fore_color.rgb = fill
    if anchor is not None:
        c.vertical_anchor = anchor
    c.margin_left = c.margin_right = Pt(5); c.margin_top = c.margin_bottom = Pt(2)
    tf = c.text_frame; tf.word_wrap = True
    if isinstance(lines, str):
        lines = [lines]
    for k, ln in enumerate(lines):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.alignment = align; p.space_before = Pt(0); p.space_after = Pt(1)
        r = p.add_run(); r.text = ln; _set(r, size, color, bold, italic, name=name)


# --------------------------------------------------------------------------- SLIDE 1: TITLE
def add_title_slide(prs, blank):
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, EMU_W, EMU_H, DARK)
    rect(s, 0, Inches(4.55), EMU_W, Inches(0.06), ACCENT)
    textbox(s, Inches(0.9), Inches(2.35), Inches(11.6), Inches(1.9),
            [(TITLE, 38, WHITE, True, False)], MSO_ANCHOR.BOTTOM)
    textbox(s, Inches(0.92), Inches(4.75), Inches(11.5), Inches(0.5),
            [(SUBTITLE, 16, RGBColor(0xC9, 0xD6, 0xE6), False, False)])
    textbox(s, Inches(0.92), Inches(5.35), Inches(11.5), Inches(0.4),
            [("CMDB Change Control Board  -  Chair: Josh Sterling  -  July 21, 2026",
              12.5, RGBColor(0x9F, 0xB3, 0xCC), False, False)])
    textbox(s, Inches(0.92), Inches(6.7), Inches(11.5), Inches(0.4),
            [(AUTHOR, 11, RGBColor(0x9F, 0xB3, 0xCC), False, False)])


# --------------------------------------------------------------------------- SLIDE 2: AGENDA & PURPOSE
AGENDA = [
    ("1", "Purpose & scope",                                    "this slide"),
    ("2", "CMP & CCB - how they work together",                 "education"),
    ("3", "What a Data Dictionary is & how it fits in the CMP", "education"),
    ("4", "Review - Business Application data dictionary",      "4 decisions"),
    ("5", "Review - Computer data dictionary",                  "3 decisions"),
    ("6", "Review - Server data dictionary",                    "5 decisions"),
    ("7", "Review - Database data dictionary",                  "3 decisions"),
    ("8", "Wrap-up & next steps",                               "close"),
]


def add_agenda_slide(prs, blank):
    s = prs.slides.add_slide(blank)
    header(s, "Agenda & Purpose  -  CMDB CCB  |  July 21, 2026",
           "Chair: Josh Sterling    Presenter: Manuel Vazquez    Duration: ~60 min")

    # purpose band
    rect(s, Inches(0.55), Inches(1.30), Inches(12.2), Inches(0.82), ACCENT_TINT, line=ACCENT)
    rect(s, Inches(0.55), Inches(1.30), Inches(0.10), Inches(0.82), ACCENT)
    textbox(s, Inches(0.80), Inches(1.30), Inches(1.35), Inches(0.82),
            [("Purpose", 12, ACCENT, True, False)], MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(2.15), Inches(1.34), Inches(10.4), Inches(0.74),
            [("Ratify CMP Stage 1 - the CMDB Data Dictionary - for the 4 core CI classes: "
              "Business Application, Computer, Server, and Database.  For each attribute the CCB confirms "
              "whether it is Mandatory, in Audit scope, and (where relevant) its allowed values.",
              11, TEXT, False, False)], MSO_ANCHOR.MIDDLE)

    # left column: agenda list
    y = 2.40; rh = 0.46; gap = 0.055
    for num, label, tag in AGENDA:
        rect(s, Inches(0.55), Inches(y), Inches(7.75), Inches(rh), LIGHT, line=RULE)
        rect(s, Inches(0.55), Inches(y), Inches(0.46), Inches(rh), ACCENT)
        textbox(s, Inches(0.55), Inches(y), Inches(0.46), Inches(rh),
                [(num, 12, WHITE, True, False)], MSO_ANCHOR.MIDDLE, PP_ALIGN.CENTER)
        textbox(s, Inches(1.15), Inches(y), Inches(5.55), Inches(rh),
                [(label, 10.5, DARK, False, False)], MSO_ANCHOR.MIDDLE)
        textbox(s, Inches(6.55), Inches(y), Inches(1.65), Inches(rh),
                [(tag, 9, MUTED, False, True)], MSO_ANCHOR.MIDDLE, PP_ALIGN.RIGHT)
        y += rh + gap

    # right column: scope callout
    rx, rw = 8.55, 4.20
    rect(s, Inches(rx), Inches(2.40), Inches(rw), Inches(0.40), GREEN)
    textbox(s, Inches(rx + 0.12), Inches(2.40), Inches(rw - 0.2), Inches(0.40),
            [("In scope today", 11.5, WHITE, True, False)], MSO_ANCHOR.MIDDLE)
    rect(s, Inches(rx), Inches(2.86), Inches(rw), Inches(1.50), GREEN_TINT, line=GREEN)
    textbox(s, Inches(rx + 0.18), Inches(2.92), Inches(rw - 0.32), Inches(1.42),
            [("4 parent classes", 11, DARK, True, False),
             ("Business Application, Computer, Server, Database", 9.5, TEXT, False, False),
             ("", 5, TEXT, False, False),
             ("15 attribute-governance decisions", 11, DARK, True, False),
             ("of the 38 in the full Data Dictionary", 9.5, TEXT, False, False)])

    rect(s, Inches(rx), Inches(4.55), Inches(rw), Inches(0.40), MUTED)
    textbox(s, Inches(rx + 0.12), Inches(4.55), Inches(rw - 0.2), Inches(0.40),
            [("Not covered today", 11.5, WHITE, True, False)], MSO_ANCHOR.MIDDLE)
    rect(s, Inches(rx), Inches(5.01), Inches(rw), Inches(1.55), LIGHT, line=RULE)
    textbox(s, Inches(rx + 0.18), Inches(5.07), Inches(rw - 0.32), Inches(1.47),
            [("Service Instance", 10, DARK, True, False),
             ("needs a Class Manager assigned first", 9, MUTED, False, False),
             ("", 4, TEXT, False, False),
             ("Child classes", 10, DARK, True, False),
             ("Windows/Linux Server, Oracle/SQL Database", 9, MUTED, False, False),
             ("- to be scheduled separately", 9, MUTED, False, True)])


# --------------------------------------------------------------------------- SLIDE 3: CMP & CCB
CMP_STAGES = [
    ("Stage 1", "Class & Attribute Data Definitions (Data Dictionary)", True),
    ("Stage 2", "CI Lifecycle & Health Targets", False),
    ("Stage 3", "Governance Process & Controls  (CO6 gate, Sep 30)", False),
]
CCB_DUTIES = [
    "Ratifies attribute governance - mandatory, audit scope, allowed values, ownership",
    "Confirms / assigns the Class Manager accountable for each class's data quality",
    "Records each ruling (Accept / Amend / Defer) and enforces it",
    "Standing body - governs all future changes to the CMDB data model",
]


def add_cmp_ccb_slide(prs, blank):
    s = prs.slides.add_slide(blank)
    header(s, "The CMP & the CCB - How They Work Together",
           "Education  |  The CMP defines what is governed; the CCB is the authority that ratifies and enforces it")

    cardy, cardh = 1.35, 3.55
    # left card - CMP
    lx, lw = 0.55, 5.95
    rect(s, Inches(lx), Inches(cardy), Inches(lw), Inches(cardh), ACCENT_TINT, line=ACCENT)
    rect(s, Inches(lx), Inches(cardy), Inches(0.10), Inches(cardh), ACCENT)
    textbox(s, Inches(lx + 0.25), Inches(cardy + 0.14), Inches(lw - 0.4), Inches(0.4),
            [("The CMP - Configuration Management Plan", 14, ACCENT, True, False)])
    textbox(s, Inches(lx + 0.25), Inches(cardy + 0.62), Inches(lw - 0.45), Inches(0.75),
            [("PPL's governance document for how the CMDB is built, maintained, and governed. "
              "Delivered in three stages:", 10.5, TEXT, False, False)])
    sy = cardy + 1.42
    for stage, desc, active in CMP_STAGES:
        col = ACCENT if active else MUTED
        fill = ACCENT_FILL if active else WHITE
        rect(s, Inches(lx + 0.25), Inches(sy), Inches(lw - 0.5), Inches(0.58), fill, line=col)
        rect(s, Inches(lx + 0.25), Inches(sy), Inches(0.08), Inches(0.58), col)
        textbox(s, Inches(lx + 0.42), Inches(sy + 0.02), Inches(1.05), Inches(0.54),
                [(stage, 10.5, col, True, False)], MSO_ANCHOR.MIDDLE)
        desc_w = (lw - 2.95) if active else (lw - 1.80)  # narrower when the TODAY chip is present
        textbox(s, Inches(lx + 1.50), Inches(sy + 0.02), Inches(desc_w), Inches(0.54),
                [(desc, 9.5, DARK if active else TEXT, active, False)], MSO_ANCHOR.MIDDLE)
        if active:
            rect(s, Inches(lx + lw - 1.30), Inches(sy + 0.11), Inches(0.95), Inches(0.34), ACCENT)
            textbox(s, Inches(lx + lw - 1.30), Inches(sy + 0.11), Inches(0.95), Inches(0.34),
                    [("TODAY", 9, WHITE, True, False)], MSO_ANCHOR.MIDDLE, PP_ALIGN.CENTER)
        sy += 0.66

    # right card - CCB
    rx, rw = 6.70, 6.08
    rect(s, Inches(rx), Inches(cardy), Inches(rw), Inches(cardh), GREEN_TINT, line=GREEN)
    rect(s, Inches(rx), Inches(cardy), Inches(0.10), Inches(cardh), GREEN)
    textbox(s, Inches(rx + 0.25), Inches(cardy + 0.14), Inches(rw - 0.4), Inches(0.4),
            [("The CCB - Change Control Board", 14, GREEN, True, False)])
    textbox(s, Inches(rx + 0.25), Inches(cardy + 0.62), Inches(rw - 0.45), Inches(0.55),
            [("The governance authority for the CMDB data model.  "
              "Chair: Josh Sterling (Configuration Management Process Owner).", 10.5, TEXT, False, False)])
    dy = cardy + 1.35
    for duty in CCB_DUTIES:
        rect(s, Inches(rx + 0.28), Inches(dy + 0.05), Inches(0.12), Inches(0.12), GREEN)
        textbox(s, Inches(rx + 0.52), Inches(dy - 0.05), Inches(rw - 0.80), Inches(0.55),
                [(duty, 10, TEXT, False, False)])
        dy += 0.53

    # bottom band - how they work together
    by = 5.20
    rect(s, Inches(0.55), Inches(by), Inches(12.2), Inches(1.55), DARK)
    rect(s, Inches(0.55), Inches(by), Inches(0.10), Inches(1.55), ACCENT)
    textbox(s, Inches(0.80), Inches(by + 0.12), Inches(12.0), Inches(0.4),
            [("How they work together", 13, WHITE, True, False)])
    textbox(s, Inches(0.80), Inches(by + 0.58), Inches(11.9), Inches(0.90),
            [("The CMP defines WHAT is governed - the classes, attributes, and rules.  "
              "The CCB is WHO governs - the authority that ratifies and enforces.  "
              "A CMP stage is only authoritative once the CCB ratifies it.", 11, WHITE, False, False),
             ("Today the CCB ratifies CMP Stage 1 - the Data Dictionary - for the 4 core classes.",
              11, RGBColor(0x9F, 0xB3, 0xCC), True, False)])


# --------------------------------------------------------------------------- SLIDE 4: WHAT IS A DATA DICTIONARY
DD_RECORDS = [
    ("Attribute name",      "Display name + technical field name"),
    ("Data type",           "String, Choice, Reference, True/False, ..."),
    ("Mandatory?",          "Is a record incomplete without it?"),
    ("Allowed values",      "Permitted choices or reference table"),
    ("Source of truth",     "Authoritative system - Discovery, SCCM, EA, Manual"),
    ("Audited?",            "Counts toward the 90% completeness / correctness target?"),
    ("Certification owner", "Who is accountable for data quality in this field"),
]


def add_datadict_slide(prs, blank):
    s = prs.slides.add_slide(blank)
    header(s, "What Is a Data Dictionary - and How It Fits in the CMP",
           "Education  |  The formal register of every attribute in the CMDB - and it IS CMP Stage 1")

    cardy, cardh = 1.35, 4.05
    # left card - definition
    lx, lw = 0.55, 6.30
    rect(s, Inches(lx), Inches(cardy), Inches(lw), Inches(cardh), ACCENT_TINT, line=ACCENT)
    rect(s, Inches(lx), Inches(cardy), Inches(0.10), Inches(cardh), ACCENT)
    textbox(s, Inches(lx + 0.25), Inches(cardy + 0.14), Inches(lw - 0.4), Inches(0.4),
            [("What it is", 14, ACCENT, True, False)])
    textbox(s, Inches(lx + 0.25), Inches(cardy + 0.60), Inches(lw - 0.45), Inches(0.55),
            [("A formal register that defines every attribute stored in the CMDB for each CI class. "
              "For each attribute it records:", 10, TEXT, False, False)])
    t = s.shapes.add_table(len(DD_RECORDS), 2, Inches(lx + 0.25), Inches(cardy + 1.22),
                           Inches(lw - 0.50), Inches(0.37 * len(DD_RECORDS))).table
    t.columns[0].width = Inches(1.85); t.columns[1].width = Inches(4.20)
    for i, (nm, desc) in enumerate(DD_RECORDS):
        fill = LIGHT if i % 2 else WHITE
        cell(t, i, 0, nm,   9, ACCENT, True,  PP_ALIGN.LEFT, fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 1, desc, 9, TEXT,   False, PP_ALIGN.LEFT, fill, MSO_ANCHOR.MIDDLE)
        t.rows[i].height = Inches(0.37)

    # right card - how it fits
    rx, rw = 7.05, 5.73
    rect(s, Inches(rx), Inches(cardy), Inches(rw), Inches(cardh), GREEN_TINT, line=GREEN)
    rect(s, Inches(rx), Inches(cardy), Inches(0.10), Inches(cardh), GREEN)
    textbox(s, Inches(rx + 0.25), Inches(cardy + 0.14), Inches(rw - 0.4), Inches(0.4),
            [("How it fits in the CMP", 14, GREEN, True, False)])
    textbox(s, Inches(rx + 0.25), Inches(cardy + 0.62), Inches(rw - 0.45), Inches(3.30),
            [("The Data Dictionary IS CMP Stage 1 - the foundation the rest of the plan is built on.",
              11, DARK, True, False),
             ("", 6, TEXT, False, False),
             ("Stages 2 and 3 are written in its terms: audit targets, certification scope, and "
              "remediation priorities all reference the attributes defined here.", 10, TEXT, False, False),
             ("", 6, TEXT, False, False),
             ("Ratifying it today turns the delivery team's draft into the agreed baseline - and "
              "unlocks the downstream governance work.", 10, TEXT, False, False),
             ("", 6, TEXT, False, False),
             ("Why it matters - without agreed definitions:", 10, DARK, True, False),
             ("- Audits can't be scoped   - Discovery conflicts can't be resolved   "
              "- Data quality can't be measured", 9.5, TEXT, False, False)])


# --------------------------------------------------------------------------- CLASS REVIEW SLIDES
def add_class_slide(prs, blank, c):
    s = prs.slides.add_slide(blank)
    n_dec = sum(1 for r in c["rows"] if r[5] is not None)
    header(s, "Review - %s Data Dictionary" % c["name"].split("  (")[0],
           "Class Manager: %s   -   Table: %s" % (c["owner"], c["table"]))

    # decisions chip
    rect(s, Inches(10.95), Inches(0.30), Inches(1.82), Inches(0.44), AMBER)
    textbox(s, Inches(10.95), Inches(0.30), Inches(1.82), Inches(0.44),
            [("%d decisions" % n_dec, 11, WHITE, True, False)], MSO_ANCHOR.MIDDLE, PP_ALIGN.CENTER)

    top = 1.45
    rows = c["rows"]
    nrows = len(rows) + 1
    widths = [Inches(2.15), Inches(1.75), Inches(3.30), Inches(1.05), Inches(0.95), Inches(3.00)]
    rh = 0.52
    t = s.shapes.add_table(nrows, 6, Inches(0.55), Inches(top), Inches(12.2),
                           Inches(0.34 + rh * len(rows))).table
    for j, cw in enumerate(widths):
        t.columns[j].width = cw
    for j, h in enumerate(["Attribute", "Type", "Source of Truth", "Mandatory", "Audited", "Team suggestion"]):
        cell(t, 0, j, h, 10, WHITE, True,
             PP_ALIGN.CENTER if j in (3, 4) else PP_ALIGN.LEFT, DARK, MSO_ANCHOR.MIDDLE)
    for i, (attr, dtype, src, mand, aud, dec, sug) in enumerate(rows, start=1):
        base = LIGHT if i % 2 else WHITE
        is_dec = dec is not None
        # attribute
        cell(t, i, 0, attr, 9.5, DARK, True, PP_ALIGN.LEFT, base, MSO_ANCHOR.MIDDLE)
        # type (amber when the decision is a Values? confirmation)
        type_open = (dec == "Values?")
        cell(t, i, 1, dtype, 8.5, AMBER if type_open else TEXT, type_open, PP_ALIGN.LEFT,
             AMBER_TINT if type_open else base, MSO_ANCHOR.MIDDLE)
        # source
        cell(t, i, 2, src, 8.5, TEXT, False, PP_ALIGN.LEFT, base, MSO_ANCHOR.MIDDLE)
        # mandatory
        m_open = (mand == "CCB?")
        cell(t, i, 3, ("CCB?" if m_open else mand), 9,
             AMBER if m_open else (DARK if mand == "Y" else MUTED), m_open or mand == "Y",
             PP_ALIGN.CENTER, AMBER_TINT if m_open else base, MSO_ANCHOR.MIDDLE)
        # audited
        a_open = (aud == "CCB?")
        cell(t, i, 4, ("CCB?" if a_open else aud), 9,
             AMBER if a_open else (DARK if aud == "Y" else MUTED), a_open or aud == "Y",
             PP_ALIGN.CENTER, AMBER_TINT if a_open else base, MSO_ANCHOR.MIDDLE)
        # team suggestion
        cell(t, i, 5, sug, 8.5, AMBER if is_dec else MUTED, is_dec, PP_ALIGN.LEFT, base, MSO_ANCHOR.MIDDLE)
    t.rows[0].height = Inches(0.34)
    for rr in range(1, nrows):
        t.rows[rr].height = Inches(rh)

    # class-specific note
    ny = top + 0.34 + rh * len(rows) + 0.15
    rect(s, Inches(0.55), Inches(ny), Inches(12.2), Inches(0.78), LIGHT, line=RULE)
    rect(s, Inches(0.55), Inches(ny), Inches(0.08), Inches(0.78), ACCENT)
    textbox(s, Inches(0.75), Inches(ny + 0.04), Inches(11.9), Inches(0.72),
            [(c["note"], 9, TEXT, False, False)], MSO_ANCHOR.MIDDLE)

    # caption
    textbox(s, Inches(0.55), Inches(6.95), Inches(12.2), Inches(0.28),
            [(CAPTION, 8, MUTED, False, True)])


# --------------------------------------------------------------------------- SLIDE 9: WRAP-UP
def add_wrapup_slide(prs, blank):
    s = prs.slides.add_slide(blank)
    header(s, "Wrap-up & Next Steps",
           "How we capture approval today - and what ratification unlocks")

    # left - how approval is captured
    lx, lw = 0.55, 6.05
    rect(s, Inches(lx), Inches(1.40), Inches(lw), Inches(3.05), ACCENT_TINT, line=ACCENT)
    rect(s, Inches(lx), Inches(1.40), Inches(0.10), Inches(3.05), ACCENT)
    textbox(s, Inches(lx + 0.25), Inches(1.52), Inches(lw - 0.4), Inches(0.4),
            [("How we capture approval", 13, ACCENT, True, False)])
    textbox(s, Inches(lx + 0.25), Inches(2.00), Inches(lw - 0.45), Inches(2.35),
            [("Accept  -  adopt the attribute as the team suggested", 10.5, GREEN, True, False),
             ("Amend  -  the chair records the change in the ruling", 10.5, AMBER, True, False),
             ("Defer  -  name an owner and a target date", 10.5, RED, True, False),
             ("", 8, TEXT, False, False),
             ("A class is accepted once all its items are ruled.  The chair rolls the per-class "
              "results into the Data Dictionary acceptance in the CCB minutes.", 10, TEXT, False, False),
             ("", 6, TEXT, False, False),
             ("Deferred items are tracked to closure; the CCB does not re-vote without new information.",
              10, TEXT, False, False)])

    # right - next steps
    rx, rw = 6.90, 5.88
    rect(s, Inches(rx), Inches(1.40), Inches(rw), Inches(3.05), GREEN_TINT, line=GREEN)
    rect(s, Inches(rx), Inches(1.40), Inches(0.10), Inches(3.05), GREEN)
    textbox(s, Inches(rx + 0.25), Inches(1.52), Inches(rw - 0.4), Inches(0.4),
            [("Next steps", 13, GREEN, True, False)])
    steps = [
        "Update the Data Dictionary with today's rulings; distribute the ratified Stage 1 within 5 business days.",
        "Ratified Stage 1 becomes the baseline for audit targets and certification scope.",
        "Unlocks CMP Stage 2 (Lifecycle & Health Targets) and Stage 3 (Governance Process & Controls).",
        "Schedule the out-of-scope classes: assign a Service Instance Class Manager; agenda the child classes (Windows/Linux, Oracle/SQL).",
    ]
    dy = 2.02
    for st in steps:
        rect(s, Inches(rx + 0.28), Inches(dy + 0.05), Inches(0.12), Inches(0.12), GREEN)
        textbox(s, Inches(rx + 0.52), Inches(dy - 0.05), Inches(rw - 0.80), Inches(0.70),
                [(st, 9.5, TEXT, False, False)])
        dy += 0.62

    # bottom band - what today ratifies
    rect(s, Inches(0.55), Inches(4.70), Inches(12.2), Inches(1.35), DARK)
    rect(s, Inches(0.55), Inches(4.70), Inches(0.10), Inches(1.35), ACCENT)
    textbox(s, Inches(0.80), Inches(4.82), Inches(12.0), Inches(0.4),
            [("What today ratifies", 13, WHITE, True, False)])
    textbox(s, Inches(0.80), Inches(5.28), Inches(11.9), Inches(0.70),
            [("15 attribute-governance decisions across 4 core classes - Business Application (4), "
              "Computer (3), Server (5), Database (3) - establishing the agreed source of truth, "
              "mandatory rules, and audit scope for the CMDB.", 11, WHITE, False, False)])


# --------------------------------------------------------------------------- BUILD
def build_deck():
    prs = Presentation(); prs.slide_width = EMU_W; prs.slide_height = EMU_H
    blank = prs.slide_layouts[6]
    add_title_slide(prs, blank)
    add_agenda_slide(prs, blank)
    add_cmp_ccb_slide(prs, blank)
    add_datadict_slide(prs, blank)
    for c in CLASSES:
        add_class_slide(prs, blank, c)
    add_wrapup_slide(prs, blank)
    return prs


def save_safe(prs, path):
    try:
        prs.save(path); print("Saved", os.path.basename(path), "-", len(prs.slides._sldIdLst), "slide(s)")
    except PermissionError:
        base, ext = os.path.splitext(path)
        alt = base + "-new" + ext
        prs.save(alt); print("Target locked - saved to", os.path.basename(alt))


here = os.path.dirname(os.path.abspath(__file__))
save_safe(build_deck(), os.path.join(here, "CMDB-Data-Dictionary-CCB-" + EDITION + ".pptx"))
