"""
CMDB Class reference slide - TEMPLATE.

One slide per CI class showing:
  - Class metadata header (label, table, Class Owner = CCB Class Manager, parent, CI count)
  - Managed Attributes table: Label | Field name | Type | Requirement | Audit (T/F)
  - IRE: short prose summary + Identification rule/entries + Reconciliation precedence

Theme mirrors the CMDB-CSDM deck family (PI-3 / CO6 decks) so it reads as a set.

To add a class: append another dict to CLASSES (same shape).  To use as a blank
template: copy a dict and clear the values.

Each class is populated from its audit-dashboard governance attribute set.
Field names/types are sourced where confirmable; Requirement and
Audit are PROPOSED (pending CCB) and marked with a trailing "*".  "(?)" = confirm
the technical name in ServiceNow.  IRE + metadata reflect the as-built config
captured 2026-07-20 (ci-class-ire-reconciliation-config.md + Stage 1 Data Dictionary).

Run:     python build_cmdb_class_template_deck.py
Output:  CMDB-Class-Reference-<EDITION>.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

EDITION   = "2026-07-20"
FOOTLABEL = "PPL CMDB-CSDM  |  CI Class Reference (CMP Stage 1)  |  DRAFT - Jul 20, 2026"
CAPTION   = ("* = proposed, pending CCB;  values without * are sourced (Stage 1 dictionary / as-built config).   "
             "(?) = confirm technical field name in ServiceNow.   IRE reflects the as-built configuration, 2026-07-20.")

# --------------------------------------------------------------------------- DATA
# One dict per class.
# attributes = (Label, Field name, Type, Requirement, req_proposed, Audit bool, audit_proposed)
#   *_proposed=True renders a trailing "*" (value proposed, pending CCB).  "(?)" in a name/type = confirm in SN.
CLASSES = [
    dict(
        label="Server",
        table="cmdb_ci_server",
        owner="Ray Reuter",
        owner_role="CCB Class Manager - Servers & Storage CIs",
        parent="Computer  >  Hardware  >  CI",
        ci_count="15,218 active",
        source_note="Source: Server audit dashboard",
        proposed=True,
        # (Label, Field name, Type, Requirement, req_proposed, Audit, audit_proposed)
        attributes=[
            ("CI Owner",              "owned_by",                "Reference -> sys_user",       "Required",    True, True,  True),
            ("Technical Owner Group", "u_technical_owner_group", "Reference -> sys_user_group", "Required",    True, True,  True),
            ("Support Group",         "support_group",           "Reference -> sys_user_group", "Required",    True, True,  False),
            ("Asset Tag",             "asset_tag",               "String",                      "Recommended", True, False, True),
            ("Value Stream",          "u_value_stream (?)",      "Reference (?)",               "Required",    True, True,  True),
            ("SOX Type",              "u_sox_type",              "Choice",                      "Recommended", True, True,  True),
            ("Classification",        "u_classification (?)",    "Choice (?)",                  "Required",    True, True,  True),
            ("Location",              "location",                "Reference -> cmn_location",   "Required",    True, True,  True),
            ("IP Address",            "ip_address",              "String",                      "Recommended", True, True,  True),
        ],
        ire_ident=dict(
            rule="PPL Server", type="Independent",
            entries=[
                ("100", "Serial number",                "Standard"),
                ("200", "Serial number (lookup table)", "Lookup"),
                ("300", "Name",                         "Standard"),
            ],
        ),
        ire_recon=dict(
            summary="Created: 0 on class  ·  Derived: 5 (inherited from Computer / Hardware)",
            rows=[
                ("80",  "ServiceNow", "IP Address, MAC Address, OS Version"),
                ("90",  "SG-SCCM",    "CPU count, Disk space, Manufacturer, +more"),
                ("100", "ServiceNow", "All  (Hardware)"),
                ("100", "SG-SCCM",    "Approval group, Asset, Assigned, +more"),
                ("200", "ServiceNow", "Approval group, Asset, Asset tag, +more"),
            ],
        ),
        ire_summary=(
            "Servers are identified independently - Serial Number (Priority 100), a serial-number "
            "lookup (200), then Name (300) - never via a parent CI. No reconciliation rules are defined "
            "on the class itself; source precedence is inherited from the Computer / Hardware hierarchy, "
            "where SG-SCCM wins for hardware, asset and ownership attributes and ServiceNow Discovery "
            "wins for IP, MAC and OS."
        ),
    ),
    dict(
        label="Computer",
        table="cmdb_ci_computer",
        owner="Monica Green  /  Paul Becker",
        owner_role="CCB Class Managers - Computers (Physical / Virtual)",
        parent="Hardware  >  CI",
        ci_count="28,203 active",
        source_note="Source: Computer audit dashboard",
        proposed=True,
        # (Label, Field name, Type, Requirement, req_proposed, Audit, audit_proposed)
        attributes=[
            ("Assigned To",           "assigned_to",             "Reference -> sys_user",       "Required",    True,  True,  False),
            ("CI Owner",              "owned_by",                "Reference -> sys_user",       "Required",    True,  True,  True),
            ("Technical Owner Group", "u_technical_owner_group", "Reference -> sys_user_group", "Required",    True,  True,  True),
            ("Support Group",         "support_group",           "Reference -> sys_user_group", "Required",    True,  True,  True),
            ("Asset Tag",             "asset_tag",               "String",                      "Recommended", True,  False, True),
            ("OS",                    "os",                      "String",                      "Recommended", True,  True,  True),
            ("Location",              "location",                "Reference -> cmn_location",   "Required",    True,  True,  True),
            ("IP Address",            "ip_address",              "String",                      "Recommended", True,  True,  True),
            ("Serial Number",         "serial_number",           "String",                      "Required",    False, True,  False),
        ],
        ire_ident=dict(
            rule="PPL Computers", type="Independent",
            entries=[
                ("100", "Name + Serial number", "Standard"),
            ],
        ),
        ire_recon=dict(
            summary="Created: 4 on class  ·  Derived: 1 (inherited from Hardware)",
            rows=[
                ("80",  "ServiceNow", "IP Address, MAC Address, OS Version"),
                ("90",  "SG-SCCM",    "CPU count, Disk space, Manufacturer, +more"),
                ("100", "SG-SCCM",    "Approval group, Asset, Assigned, Attestation, +more"),
                ("300", "ServiceNow", "Approval group, Asset, Asset tag, Assigned, +more"),
                ("100", "ServiceNow", "All  (Hardware, derived)"),
            ],
        ),
        ire_summary=(
            "Computers are identified independently by a single rule (Priority 100) matching on "
            "Name + Serial Number (the 'Duplicate Workstation' identifier). Unlike Server, reconciliation "
            "is defined directly on the class - four created rules (ServiceNow[80] to ServiceNow[300]) plus "
            "one allow-all rule inherited from Hardware. SG-SCCM is authoritative for OS, OS Version and "
            "hardware specifications; ServiceNow Discovery populates IP and MAC address."
        ),
    ),
    dict(
        label="Business Application",
        table="cmdb_ci_business_app",
        owner="Todd Dierksheide",
        owner_role="CCB Class Manager - Business Application",
        domain="Design & Planning",
        ci_count="~2,150 records",
        source_note="Source: Business Application audit dashboard",
        proposed=True,
        # (Label, Field name, Type, Requirement, req_proposed, Audit, audit_proposed)
        attributes=[
            ("CI Owner",              "owned_by",                "Reference -> sys_user",       "Required",    True, True,  True),
            ("Technical Owner Group", "u_technical_owner_group", "Reference -> sys_user_group", "Required",    True, True,  True),
            ("Business Owner",        "business_owner",          "Reference -> sys_user",       "Required",    True, True,  False),
            ("Approval Group",        "u_approval_group",        "Reference -> sys_user_group", "Required",    True, True,  True),
            ("Support Group",         "support_group",           "Reference -> sys_user_group", "Required",    True, True,  True),
            ("Asset Tag",             "asset_tag",               "String",                      "Optional",    True, False, True),
            ("Value Stream",          "u_value_stream (?)",      "Reference (?)",               "Required",    True, True,  True),
            ("SOX Type",              "u_sox_type",              "Choice",                      "Recommended", True, True,  True),
            ("Classification",        "u_classification (?)",    "Choice (?)",                  "Required",    True, True,  True),
        ],
        ire_ident=dict(
            rule="not captured", type="by Name (Stage 1)",
            entries=[
                ("-", "Name", "-"),
            ],
        ),
        ire_recon=dict(
            summary="No CI Class Manager rules captured - source of truth per Stage 1 dictionary:",
            rows=[
                ("-", "EA (LeanIX)",        "Application Category, Life Cycle Stage - authoritative"),
                ("-", "Manual / Data Cert", "Ownership & governance fields - via certification"),
            ],
        ),
        ire_summary=(
            "Business Applications are the abstract application model, identified by Name - not a "
            "discovered CI. No rule-level IRE or reconciliation was captured for this class in CI Class "
            "Manager. Per the Stage 1 dictionary, Enterprise Architecture integration (LeanIX) is "
            "authoritative for Application Category and Life Cycle Stage, while ownership and governance "
            "fields are maintained manually through the Data Certification workflow."
        ),
    ),
    dict(
        label="Database Instance",
        table="cmdb_ci_database",
        owner="Ray Reuter",
        owner_role="CCB Class Manager - Database CIs",
        domain="Service Delivery",
        ci_count="3,113 active",
        source_note="Source: Database audit dashboard",
        proposed=True,
        # (Label, Field name, Type, Requirement, req_proposed, Audit, audit_proposed)
        attributes=[
            ("CI Owner",              "owned_by",                "Reference -> sys_user",       "Required",    True, True, True),
            ("Technical Owner Group", "u_technical_owner_group", "Reference -> sys_user_group", "Required",    True, True, True),
            ("Support Group",         "support_group",           "Reference -> sys_user_group", "Required",    True, True, False),
            ("Environment",           "environment (?)",         "Choice",                      "Recommended", True, True, True),
            ("Location",              "location",                "Reference -> cmn_location",   "Recommended", True, True, True),
        ],
        ire_ident=dict(
            rule="Database instance rule", type="Dependent (host Server)",
            entries=[
                ("100", "Serial number",                "Standard"),
                ("200", "Edition + Name + TCP port(s)", "Standard"),
            ],
        ),
        ire_recon=dict(
            summary="No reconciliation rules captured on the parent class - source of truth per Stage 1:",
            rows=[
                ("-", "ServiceNow (App Probes)", "Active running version + TCP port - authoritative"),
            ],
        ),
        ire_summary=(
            "A Database Instance is identified dependently - always in the context of its host Server - "
            "first by Serial Number (Priority 100), then by Edition + Name + TCP port(s) (Priority 200). "
            "No reconciliation rules were captured on the parent class; per Stage 1, ServiceNow Discovery "
            "application probes are authoritative for the active running version and port details. Child "
            "classes (Oracle, MS SQL) add engine-specific discovery and reconciliation."
        ),
    ),
]

# --------------------------------------------------------------------------- THEME
DARK   = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0x2E, 0x6D, 0xB4)
LIGHT  = RGBColor(0xF4, 0xF6, 0xF9)
TEXT   = RGBColor(0x2D, 0x2D, 0x2D)
MUTED  = RGBColor(0x6B, 0x72, 0x80)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
AMBER  = RGBColor(0xE0, 0x8A, 0x00)
RED    = RGBColor(0xC0, 0x39, 0x2B)
RULE   = RGBColor(0xDD, 0xE3, 0xEA)
ACCENT_TINT = RGBColor(0xE9, 0xF0, 0xF8)
EMU_W, EMU_H = Inches(13.333), Inches(7.5)
FONT = "Segoe UI"
MONO = "Consolas"

# Requirement value -> (text color, bold)
REQ_STYLE = {
    "Required":    (RED,   True),
    "Recommended": (AMBER, True),
    "Optional":    (MUTED, False),
}


def _set(run, size, color, bold=False, italic=False, name=FONT):
    run.font.size = Pt(size); run.font.color.rgb = color
    run.font.bold = bold; run.font.italic = italic; run.font.name = name


def rect(slide, x, y, w, h, color, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def textbox(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(4); tf.margin_top = tf.margin_bottom = Pt(2)
    for i, spec in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        text, size, color, bold, italic = spec
        r = p.add_run(); r.text = text; _set(r, size, color, bold, italic)
    return tb


def header(slide, title, subtitle=None):
    rect(slide, 0, 0, EMU_W, Inches(1.05), DARK)
    rect(slide, 0, Inches(1.05), EMU_W, Inches(0.06), ACCENT)
    textbox(slide, Inches(0.55), Inches(0.16), Inches(12.2), Inches(0.55),
            [(title, 23, WHITE, True, False)], MSO_ANCHOR.MIDDLE)
    if subtitle:
        textbox(slide, Inches(0.57), Inches(0.66), Inches(12.2), Inches(0.32),
                [(subtitle, 12, RGBColor(0xC9, 0xD6, 0xE6), False, False)], MSO_ANCHOR.MIDDLE)
    textbox(slide, Inches(0.55), Inches(7.08), Inches(10), Inches(0.32),
            [(FOOTLABEL, 9, MUTED, False, False)])


def cell(t, i, j, lines, size, color, bold=False, align=PP_ALIGN.LEFT,
         fill=None, anchor=None, name=FONT):
    c = t.cell(i, j)
    if fill is not None:
        c.fill.solid(); c.fill.fore_color.rgb = fill
    if anchor is not None:
        c.vertical_anchor = anchor
    c.margin_left = c.margin_right = Pt(5); c.margin_top = c.margin_bottom = Pt(2)
    tf = c.text_frame; tf.word_wrap = True
    if isinstance(lines, str):
        lines = [lines]
    for k, ln in enumerate(lines):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.alignment = align; p.space_before = Pt(0); p.space_after = Pt(1)
        r = p.add_run(); r.text = ln; _set(r, size, color, bold, name=name)


def meta_field(slide, x, w, label, value, sub=None):
    lines = [(label.upper(), 8.5, MUTED, True, False), (value, 12.5, DARK, True, False)]
    if sub:
        lines.append((sub, 8.5, MUTED, False, True))
    textbox(slide, Inches(x), Inches(1.28), Inches(w), Inches(0.86), lines, MSO_ANCHOR.MIDDLE)


def attributes_table(slide, x, y, w, attrs):
    rows = len(attrs) + 1
    t = slide.shapes.add_table(rows, 5, Inches(x), Inches(y), Inches(w), Inches(0.30 * rows)).table
    for j, cw in enumerate([Inches(1.55), Inches(1.95), Inches(1.80), Inches(1.30), Inches(0.85)]):
        t.columns[j].width = cw
    for j, h in enumerate(["Label", "Field name", "Type", "Requirement", "Audit"]):
        cell(t, 0, j, h, 10, WHITE, True, PP_ALIGN.CENTER if j == 4 else PP_ALIGN.LEFT, DARK)
    for i, (lab, field, typ, req, req_p, audit, audit_p) in enumerate(attrs, start=1):
        fill = LIGHT if i % 2 else WHITE
        rc, rb = REQ_STYLE.get(req, (TEXT, False))
        req_txt = req + (" *" if req_p else "")
        audit_txt = ("T" if audit else "F") + ("*" if audit_p else "")
        cell(t, i, 0, lab,     9,   DARK,   True,  PP_ALIGN.LEFT, fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 1, field,   8.5, ACCENT, False, PP_ALIGN.LEFT, fill, MSO_ANCHOR.MIDDLE, name=MONO)
        cell(t, i, 2, typ,     8.5, TEXT,   False, PP_ALIGN.LEFT, fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 3, req_txt, 8.5, rc,     rb,    PP_ALIGN.LEFT, fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 4, audit_txt, 10, GREEN if audit else MUTED, True,
             PP_ALIGN.CENTER, fill, MSO_ANCHOR.MIDDLE)
    for r in range(rows):
        t.rows[r].height = Inches(0.30)


def ire_panel(slide, x, y, w, ident, recon):
    # panel title
    rect(slide, Inches(x), Inches(y), Inches(w), Inches(0.32), ACCENT)
    textbox(slide, Inches(x + 0.10), Inches(y), Inches(w - 0.2), Inches(0.32),
            [("Identification & Reconciliation (IRE)", 11, WHITE, True, False)], MSO_ANCHOR.MIDDLE)

    # Identification
    iy = y + 0.40
    textbox(slide, Inches(x), Inches(iy), Inches(w), Inches(0.22),
            [("IDENTIFICATION", 9, DARK, True, False)], MSO_ANCHOR.MIDDLE)
    textbox(slide, Inches(x), Inches(iy + 0.22), Inches(w), Inches(0.22),
            [("Rule: %s   ·   Type: %s" % (ident["rule"], ident["type"]), 9, TEXT, False, False)],
            MSO_ANCHOR.MIDDLE)
    ie = ident["entries"]
    t = slide.shapes.add_table(len(ie) + 1, 3, Inches(x), Inches(iy + 0.46),
                               Inches(w), Inches(0.24 * (len(ie) + 1))).table
    for j, cw in enumerate([Inches(0.65), Inches(2.75), Inches(1.05)]):
        t.columns[j].width = cw
    for j, h in enumerate(["Pri", "Criterion attributes", "Match"]):
        cell(t, 0, j, h, 8.5, WHITE, True, PP_ALIGN.CENTER if j != 1 else PP_ALIGN.LEFT, DARK)
    for r, (pri, crit, match) in enumerate(ie, start=1):
        fill = LIGHT if r % 2 else WHITE
        cell(t, r, 0, pri,   8.5, ACCENT, True,  PP_ALIGN.CENTER, fill, MSO_ANCHOR.MIDDLE)
        cell(t, r, 1, crit,  8.5, TEXT,   False, PP_ALIGN.LEFT,   fill, MSO_ANCHOR.MIDDLE)
        cell(t, r, 2, match, 8.5, MUTED,  False, PP_ALIGN.CENTER, fill, MSO_ANCHOR.MIDDLE)
    for r in range(len(ie) + 1):
        t.rows[r].height = Inches(0.24)

    # Reconciliation
    ry = iy + 0.46 + 0.24 * (len(ie) + 1) + 0.12
    textbox(slide, Inches(x), Inches(ry), Inches(w), Inches(0.22),
            [("RECONCILIATION", 9, DARK, True, False)], MSO_ANCHOR.MIDDLE)
    textbox(slide, Inches(x), Inches(ry + 0.22), Inches(w), Inches(0.22),
            [(recon["summary"], 8.5, TEXT, False, False)], MSO_ANCHOR.MIDDLE)
    rr = recon["rows"]
    t = slide.shapes.add_table(len(rr) + 1, 3, Inches(x), Inches(ry + 0.46),
                               Inches(w), Inches(0.24 * (len(rr) + 1))).table
    for j, cw in enumerate([Inches(0.65), Inches(1.35), Inches(2.45)]):
        t.columns[j].width = cw
    for j, h in enumerate(["Pri", "Source", "Attributes (winning source)"]):
        cell(t, 0, j, h, 8.5, WHITE, True, PP_ALIGN.CENTER if j == 0 else PP_ALIGN.LEFT, DARK)
    for r, (pri, src, at) in enumerate(rr, start=1):
        fill = LIGHT if r % 2 else WHITE
        cell(t, r, 0, pri, 8.5, ACCENT, True,  PP_ALIGN.CENTER, fill, MSO_ANCHOR.MIDDLE)
        cell(t, r, 1, src, 8.5, DARK,   True,  PP_ALIGN.LEFT,   fill, MSO_ANCHOR.MIDDLE)
        cell(t, r, 2, at,  8,   TEXT,   False, PP_ALIGN.LEFT,   fill, MSO_ANCHOR.MIDDLE)
    for r in range(len(rr) + 1):
        t.rows[r].height = Inches(0.24)


def add_class_slide(prs, blank, c):
    s = prs.slides.add_slide(blank)
    sub = "Table: %s" % c["table"]
    if c.get("source_note"):
        sub += "   ·   " + c["source_note"]
    header(s, "CMDB Class  -  %s" % c["label"], sub)
    if c.get("proposed"):
        rect(s, Inches(10.75), Inches(0.30), Inches(2.02), Inches(0.44), RED)
        textbox(s, Inches(10.75), Inches(0.30), Inches(2.02), Inches(0.44),
                [("PROPOSED DRAFT", 10, WHITE, True, False)], MSO_ANCHOR.MIDDLE, PP_ALIGN.CENTER)

    # metadata strip
    rect(s, Inches(0.55), Inches(1.22), Inches(12.2), Inches(0.98), ACCENT_TINT, line=ACCENT)
    audited = sum(1 for a in c["attributes"] if a[5])
    meta_field(s, 0.75, 4.30, "Class Owner", c["owner"], c["owner_role"])
    if c.get("parent"):
        meta_field(s, 5.30, 3.10, "Parent Class", c["parent"])
    elif c.get("domain"):
        meta_field(s, 5.30, 3.10, "CSDM Domain", c["domain"])
    meta_field(s, 8.55, 2.10, "Active CIs", c["ci_count"])
    meta_field(s, 10.75, 1.90, "Audited (proposed)", "%d of %d" % (audited, len(c["attributes"])))

    # section label
    textbox(s, Inches(0.55), Inches(2.24), Inches(7.45), Inches(0.28),
            [("Managed Attributes", 13, DARK, True, False)], MSO_ANCHOR.MIDDLE)

    attributes_table(s, 0.55, 2.52, 7.45, c["attributes"])
    ire_panel(s, 8.30, 2.24, 4.45, c["ire_ident"], c["ire_recon"])

    # bottom IRE summary band (full width)
    rect(s, Inches(0.55), Inches(6.20), Inches(12.2), Inches(0.80), ACCENT_TINT, line=ACCENT)
    rect(s, Inches(0.55), Inches(6.20), Inches(0.08), Inches(0.80), ACCENT)
    textbox(s, Inches(0.72), Inches(6.20), Inches(1.20), Inches(0.80),
            [("IRE", 11, DARK, True, False), ("Summary", 11, DARK, True, False)], MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(2.00), Inches(6.24), Inches(10.6), Inches(0.44),
            [(c["ire_summary"], 9.5, TEXT, False, False)], MSO_ANCHOR.TOP)
    textbox(s, Inches(2.00), Inches(6.70), Inches(10.6), Inches(0.28),
            [(CAPTION, 8, MUTED, False, True)], MSO_ANCHOR.TOP)


def build_deck():
    prs = Presentation(); prs.slide_width = EMU_W; prs.slide_height = EMU_H
    blank = prs.slide_layouts[6]
    for c in CLASSES:
        add_class_slide(prs, blank, c)
    return prs


def save_safe(prs, path):
    try:
        prs.save(path); print("Saved", os.path.basename(path), "-", len(prs.slides._sldIdLst), "slide(s)")
    except PermissionError:
        base, ext = os.path.splitext(path)
        alt = base + "-new" + ext
        prs.save(alt); print("Target locked - saved to", os.path.basename(alt))


here = os.path.dirname(os.path.abspath(__file__))
save_safe(build_deck(), os.path.join(here, "CMDB-Class-Reference-" + EDITION + ".pptx"))
