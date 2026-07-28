"""
PI-3 Plan deck - CMDB / CSDM team ADO objectives & features (team share).

Two requested views:
  * Objective View - features grouped under each of the 4 objectives
  * Sprint View    - features grouped by iteration (3.1 - 3.6 IP)
  * PI Program Board - visual swimlane of features across the 6 sprints

Source: PI-3/pi3-cmdb-csdm-ado-tracking.md (Azure DevOps, area A-INFOPS\\FY26\\PI3),
as of 2026-07-22. CMDB/CSDM focus area is a SUBSET of CO6 (4 of 7 CO6 objectives).
Full 7 CO6 objectives live in the PI3-Objectives decks. Theme mirrors those decks (a set).

Run:  python build_pi3_features_deck.py
Output:  PI3-Objectives-Features-<EDITION>.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --------------------------------------------------------------------------- CONTENT
EDITION   = "2026-07-22"
TITLE     = "PI-3 Plan - CMDB / CSDM"
SUBTITLE  = "Objectives & Features  ·  Objective View + Sprint View"
AUTHOR    = "Manuel Vazquez  -  Scrum Master, CMDB-CSDM"
FOOTLABEL = "PPL CMDB-CSDM  |  PI-3 Objectives & Features  |  ADO A-INFOPS\\FY26\\PI3  |  Jul 22, 2026"
SOURCE_NOTE = ("Source: Azure DevOps (A-INFOPS\\FY26\\PI3), as of Jul 22, 2026.  "
               "CMDB/CSDM focus area - a subset of CO6 (maps to 4 of the 7 CO6 objectives).")

# Objectives (ADO order).  id, name, workstream key, CO6 map, owner, one-line descriptor
OBJECTIVES = [
    dict(id="1516954", name="Automate Network Device Discovery in the CMDB", ws="NET", co6="CO6 Obj 4",
         gates="Aug 31 · Sep 30 · Oct 30", owner="M. Vazquez",
         desc="Automated credentialed discovery of all network device types, mandatory attribute population, and ≥90% owner-validated coverage."),
    dict(id="1516955", name="Build & Validate Service Maps for Priority Business Apps", ws="SM", co6="CO6 Obj 5",
         gates="Aug 31 · Sep 30 · Oct 30", owner="A. Griffis",
         desc="Owner-validated Silver- and Gold-tier service maps, with ≥75% of the business-service inventory represented."),
    dict(id="1516956", name="Certify CMDB Data Across the Four Major CI Classes", ws="CERT", co6="CO6 Obj 3",
         gates="Sep 30 · Oct 30", owner="M. Vazquez",
         desc="End-to-end data certification - process, PROD build, dashboards, training - across Computers, Servers, Databases, Network Devices."),
    dict(id="1516957", name="Activate the ServiceNow → Qualys Attribute Sync in Production", ws="QLY", co6="CO6 Obj 2",
         gates="Sep 30", owner="M. Vazquez",
         desc="One-way ServiceNow → Qualys sync of owner, support group, and SOX flag, live in production on a defined schedule."),
]

# Features.  id, parent objective, workstream, sprint (ADO iteration), start, CO6 delivery date, owner, name, short label
# co6 = the CO6 acceptance-criteria completion date this feature delivers against (the delivery commitment).
FEATURES = [
    dict(id="1516993", obj="1516954", ws="NET",  sprint="3.1", start="Aug 5",  co6="Aug 31, 2026", owner="M. Vazquez",
         name="Automated Credentialed Discovery Live — All Network Device Types", short="Discovery Live — all device types"),
    dict(id="1517005", obj="1516954", ws="NET",  sprint="3.3", start="Sep 2",  co6="Sep 30, 2026", owner="M. Vazquez",
         name="Mandatory Attribute Population via Discovery — Network Devices", short="Mandatory Attribute Population"),
    dict(id="1356646", obj="1516954", ws="NET",  sprint="3.5", start="Sep 30", co6="Oct 30, 2026", owner="M. Vazquez",
         name="Network Inventory Coverage Reconciliation & Owner Validation (≥90%)", short="Coverage Reconciliation ≥90%"),
    dict(id="1517650", obj="1516955", ws="SM",   sprint="3.1", start="Aug 5",  co6="Aug 31, 2026", owner="A. Griffis",
         name="Silver-Tier Service Maps (App → Infrastructure) + ≥75% Business-Service Inventory", short="Silver-Tier Maps + ≥75% Inventory"),
    dict(id="1517655", obj="1516955", ws="SM",   sprint="3.3", start="Sep 2",  co6="Sep 30, 2026", owner="A. Griffis",
         name="Gold-Tier Service Maps (Service → App) + Silver Validation Pass", short="Gold-Tier Maps + Silver Validation"),
    dict(id="1518000", obj="1516955", ws="SM",   sprint="3.5", start="Sep 30", co6="Oct 30, 2026", owner="A. Griffis",
         name="Remaining Contractor-Managed Silver-Tier Maps + All-Silver Service → App Layer", short="Remaining Silver + Svc→App"),
    dict(id="1517029", obj="1516956", ws="CERT", sprint="3.2", start="Aug 19", co6="Sep 30, 2026", owner="M. Vazquez",
         name="Certify Computers — CMDB Data Certification", short="Certify Computers"),
    dict(id="1517032", obj="1516956", ws="CERT", sprint="3.2", start="Aug 19", co6="Sep 30, 2026", owner="M. Vazquez",
         name="Certify Servers — CMDB Data Certification", short="Certify Servers"),
    dict(id="1517037", obj="1516956", ws="CERT", sprint="3.4", start="Sep 16", co6="Oct 30, 2026", owner="M. Vazquez",
         name="Certify Databases — CMDB Data Certification", short="Certify Databases"),
    dict(id="1517040", obj="1516956", ws="CERT", sprint="3.4", start="Sep 16", co6="Oct 30, 2026", owner="M. Vazquez",
         name="Certify Network Devices — CMDB Data Certification", short="Certify Network Devices"),
    dict(id="1517065", obj="1516957", ws="QLY",  sprint="3.1", start="Aug 5",  co6="Sep 30, 2026", owner="S. Tomberg",
         name="Activate ServiceNow → Qualys Attribute Sync (PROD)", short="Activate SN→Qualys Sync (PROD)"),
]

# CO6 delivery gates (chronological): date, count timing note relative to the PI-3 window
GATE_NOTE = {
    "Aug 31, 2026": "delivered in Sprint 3.1 (ends Aug 18)  ·  ahead of the commitment",
    "Sep 30, 2026": "delivered in Sprints 3.1 - 3.3 (end by Sep 15)  ·  ahead of the commitment",
    "Oct 30, 2026": "delivered in Sprints 3.4 - 3.5 (end by Oct 13)  ·  ahead of the Oct 30 date (after PI-3 close, Oct 27)",
}
GATE_ORDER = ["Aug 31, 2026", "Sep 30, 2026", "Oct 30, 2026"]

# Sprints (iteration cadence).  code, dates
SPRINTS = [
    ("3.1",    "Aug 5 - Aug 18"),
    ("3.2",    "Aug 19 - Sep 1"),
    ("3.3",    "Sep 2 - Sep 15"),
    ("3.4",    "Sep 16 - Sep 29"),
    ("3.5",    "Sep 30 - Oct 13"),
    ("3.6 IP", "Oct 14 - Oct 27"),
]

# --------------------------------------------------------------------------- THEME
DARK   = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0x2E, 0x6D, 0xB4)
LIGHT  = RGBColor(0xF4, 0xF6, 0xF9)
TEXT   = RGBColor(0x2D, 0x2D, 0x2D)
MUTED  = RGBColor(0x6B, 0x72, 0x80)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
RED    = RGBColor(0xC0, 0x39, 0x2B)
RULE   = RGBColor(0xDD, 0xE3, 0xEA)
ACCENT_TINT = RGBColor(0xE9, 0xF0, 0xF8)
AMBER_FILL  = RGBColor(0xFB, 0xEE, 0xDB)
AMBER_BODY  = RGBColor(0x5A, 0x4A, 0x2A)

# Workstream (one per objective): fill (vivid, for chips) and text (darker, for on-white)
BLUE   = RGBColor(0x2E, 0x6D, 0xB4); BLUE_T   = RGBColor(0x2E, 0x6D, 0xB4)
TEAL   = RGBColor(0x2A, 0x9D, 0x8F); TEAL_T   = RGBColor(0x1F, 0x7A, 0x6E)
INDIGO = RGBColor(0x5B, 0x4B, 0x8A); INDIGO_T = RGBColor(0x4A, 0x3D, 0x73)
AMBER  = RGBColor(0xE0, 0x8A, 0x00); AMBER_T  = RGBColor(0xB5, 0x6E, 0x00)
WS = {
    "NET":  ("Network Discovery",  BLUE,   BLUE_T),
    "SM":   ("Service Maps",       TEAL,   TEAL_T),
    "CERT": ("Data Certification", INDIGO, INDIGO_T),
    "QLY":  ("Qualys Sync",        AMBER,  AMBER_T),
}

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
FONT = "Segoe UI"


def _set(run, size, color, bold=False, italic=False):
    run.font.size = Pt(size); run.font.color.rgb = color
    run.font.bold = bold; run.font.italic = italic; run.font.name = FONT


def rect(slide, x, y, w, h, color, line=None, shape=MSO_SHAPE.RECTANGLE):
    shp = slide.shapes.add_shape(shape, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def textbox(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(4); tf.margin_top = tf.margin_bottom = Pt(2)
    for i, spec in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        text, size, color, bold, italic = spec
        r = p.add_run(); r.text = text; _set(r, size, color, bold, italic)
    return tb


def header(slide, title, subtitle=None):
    rect(slide, 0, 0, EMU_W, Inches(1.05), DARK)
    rect(slide, 0, Inches(1.05), EMU_W, Inches(0.06), ACCENT)
    textbox(slide, Inches(0.55), Inches(0.16), Inches(12.2), Inches(0.55), [(title, 23, WHITE, True, False)], MSO_ANCHOR.MIDDLE)
    if subtitle:
        textbox(slide, Inches(0.57), Inches(0.66), Inches(12.2), Inches(0.32), [(subtitle, 12, RGBColor(0xC9, 0xD6, 0xE6), False, False)], MSO_ANCHOR.MIDDLE)
    textbox(slide, Inches(0.55), Inches(7.08), Inches(11.5), Inches(0.32), [(FOOTLABEL, 9, MUTED, False, False)])


def cell(t, i, j, lines, size, color, bold=False, align=PP_ALIGN.LEFT, fill=None, anchor=None, italic=False):
    c = t.cell(i, j)
    if fill is not None:
        c.fill.solid(); c.fill.fore_color.rgb = fill
    if anchor is not None:
        c.vertical_anchor = anchor
    tf = c.text_frame; tf.word_wrap = True
    if isinstance(lines, str):
        lines = [lines]
    for k, ln in enumerate(lines):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.alignment = align; p.space_before = Pt(0); p.space_after = Pt(2)
        r = p.add_run(); r.text = ln; _set(r, size, color, bold, italic)


def cell_multi(t, i, j, items, size, align=PP_ALIGN.LEFT, fill=None, anchor=MSO_ANCHOR.MIDDLE):
    """Fill one cell with multiple lines, each its own (text, color, bold)."""
    c = t.cell(i, j)
    if fill is not None:
        c.fill.solid(); c.fill.fore_color.rgb = fill
    c.vertical_anchor = anchor
    tf = c.text_frame; tf.word_wrap = True
    for k, (text, color, bold) in enumerate(items):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.alignment = align; p.space_before = Pt(0); p.space_after = Pt(3)
        r = p.add_run(); r.text = text; _set(r, size, color, bold)


def co6_color(date):
    """Delivery-commitment color: red once it lands after PI-3 close, green if comfortably early."""
    if date.startswith("Oct 30"):
        return RED
    if date.startswith("Aug 31"):
        return GREEN
    return ACCENT


def co6_short(date):
    return date.split(",")[0]


# --------------------------------------------------------------------------- SLIDES
def add_title(prs, blank):
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, EMU_W, EMU_H, DARK)
    rect(s, 0, Inches(4.55), EMU_W, Inches(0.08), ACCENT)
    textbox(s, Inches(0.8), Inches(1.65), Inches(11.7), Inches(1.1), [(TITLE, 40, WHITE, True, False)])
    textbox(s, Inches(0.82), Inches(2.72), Inches(11.7), Inches(0.6), [(SUBTITLE, 19, RGBColor(0xC9, 0xD6, 0xE6), False, False)])
    textbox(s, Inches(0.82), Inches(4.75), Inches(11.7), Inches(1.5), [
        ("4 objectives  ·  11 features  ·  6 sprints (Aug 5 - Oct 27, 2026)", 16, WHITE, True, False),
        ("CMDB / CSDM team focus area  -  a subset of CO6", 14, RGBColor(0xC9, 0xD6, 0xE6), False, False),
        (AUTHOR, 12, RGBColor(0x9F, 0xB2, 0xCC), False, False),
    ])
    textbox(s, Inches(0.82), Inches(6.5), Inches(11.7), Inches(0.4), [(SOURCE_NOTE, 11, RGBColor(0x7F, 0x8C, 0xA6), False, True)])


def add_overview(prs, blank):
    s = prs.slides.add_slide(blank)
    header(s, "PI-3 Objectives - Overview", "4 CMDB/CSDM objectives  ·  11 features  ·  each mapped to CO6")
    cols = ["#", "Objective", "ADO ID", "Owner", "Feat", "CO6"]
    widths = [Inches(0.45), Inches(6.75), Inches(1.15), Inches(1.7), Inches(0.65), Inches(1.5)]
    t = s.shapes.add_table(len(OBJECTIVES) + 1, len(cols), Inches(0.55), Inches(1.5), Inches(12.2), Inches(2.7)).table
    for j, wdt in enumerate(widths):
        t.columns[j].width = wdt
    for j, h in enumerate(cols):
        cell(t, 0, j, h, 11, WHITE, True, PP_ALIGN.LEFT if j == 1 else PP_ALIGN.CENTER, DARK)
    for i, o in enumerate(OBJECTIVES, start=1):
        fill = LIGHT if i % 2 else WHITE
        vivid = WS[o["ws"]][1]
        n_feat = sum(1 for f in FEATURES if f["obj"] == o["id"])
        cell(t, i, 0, str(i), 11, WHITE, True, PP_ALIGN.CENTER, vivid, MSO_ANCHOR.MIDDLE)
        cell(t, i, 1, [o["name"], o["desc"]], 10.5, DARK, True, PP_ALIGN.LEFT, fill, MSO_ANCHOR.MIDDLE)
        # dim the descriptor line
        t.cell(i, 1).text_frame.paragraphs[1].runs[0].font.bold = False
        t.cell(i, 1).text_frame.paragraphs[1].runs[0].font.size = Pt(8.5)
        t.cell(i, 1).text_frame.paragraphs[1].runs[0].font.color.rgb = MUTED
        cell(t, i, 2, o["id"], 10, ACCENT, True, PP_ALIGN.CENTER, fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 3, o["owner"], 10, TEXT, False, PP_ALIGN.CENTER, fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 4, str(n_feat), 11, DARK, True, PP_ALIGN.CENTER, fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 5, [o["co6"], o["gates"]], 9.5, WS[o["ws"]][2], True, PP_ALIGN.CENTER, fill, MSO_ANCHOR.MIDDLE)
        g = t.cell(i, 5).text_frame.paragraphs[1].runs[0].font
        g.bold = False; g.size = Pt(8); g.color.rgb = MUTED
    cols_hdr = t.cell(0, 5).text_frame.paragraphs[0].runs[0]
    cols_hdr.text = "CO6  ·  delivery gates"
    rect(s, Inches(0.55), Inches(4.55), Inches(12.2), Inches(0.85), ACCENT_TINT, line=ACCENT)
    textbox(s, Inches(0.72), Inches(4.6), Inches(11.9), Inches(0.75),
            [(SOURCE_NOTE, 10.5, DARK, False, False)], MSO_ANCHOR.MIDDLE)


def add_objective_view(prs, blank):
    s = prs.slides.add_slide(blank)
    header(s, "Objective View - Features by Objective", "11 features grouped under their 4 parent objectives  ·  with CO6 delivery commitment")
    # build row plan
    rowplan = [("head", None)]
    for o in OBJECTIVES:
        rowplan.append(("group", o))
        for f in FEATURES:
            if f["obj"] == o["id"]:
                rowplan.append(("feat", f))
    t = s.shapes.add_table(len(rowplan), 5, Inches(0.55), Inches(1.32), Inches(12.2), Inches(5.5)).table
    for j, wdt in enumerate([Inches(1.05), Inches(7.35), Inches(0.9), Inches(1.05), Inches(1.85)]):
        t.columns[j].width = wdt
    for r, (kind, data) in enumerate(rowplan):
        if kind == "head":
            for j, h in enumerate(["ADO ID", "Feature", "Sprint", "Start", "Delivery (CO6)"]):
                cell(t, r, j, h, 10.5, WHITE, True, PP_ALIGN.LEFT if j == 1 else PP_ALIGN.CENTER, DARK)
        elif kind == "group":
            t.cell(r, 0).merge(t.cell(r, 4))
            vivid = WS[data["ws"]][1]
            label = ("%s   ·   ADO %s   ·   %s   ·   Owner: %s"
                     % (data["name"], data["id"], data["co6"], data["owner"]))
            cell(t, r, 0, label, 10.5, WHITE, True, PP_ALIGN.LEFT, vivid, MSO_ANCHOR.MIDDLE)
        else:
            f = data
            fill = LIGHT
            cell(t, r, 0, f["id"], 9.5, ACCENT, True, PP_ALIGN.CENTER, fill, MSO_ANCHOR.MIDDLE)
            cell(t, r, 1, f["name"], 9.5, TEXT, False, PP_ALIGN.LEFT, fill, MSO_ANCHOR.MIDDLE)
            cell(t, r, 2, f["sprint"], 9.5, DARK, True, PP_ALIGN.CENTER, fill, MSO_ANCHOR.MIDDLE)
            cell(t, r, 3, f["start"], 9.5, TEXT, False, PP_ALIGN.CENTER, fill, MSO_ANCHOR.MIDDLE)
            cell(t, r, 4, f["co6"], 9.5, co6_color(f["co6"]), True, PP_ALIGN.CENTER, fill, MSO_ANCHOR.MIDDLE)


def add_sprint_view(prs, blank):
    s = prs.slides.add_slide(blank)
    header(s, "Sprint View - Features by Iteration", "Which features land in which PI-3 sprint (by ADO Iteration Path)")
    cols = ["Sprint", "Dates", "#", "Features  (ADO ID  ·  CO6 delivery date  ·  color = objective)"]
    widths = [Inches(1.2), Inches(2.15), Inches(0.55), Inches(8.3)]
    t = s.shapes.add_table(len(SPRINTS) + 1, len(cols), Inches(0.55), Inches(1.4), Inches(12.2), Inches(5.3)).table
    for j, wdt in enumerate(widths):
        t.columns[j].width = wdt
    for j, h in enumerate(cols):
        cell(t, 0, j, h, 10.5, WHITE, True, PP_ALIGN.LEFT if j == 3 else PP_ALIGN.CENTER, DARK)
    for i, (code, dates) in enumerate(SPRINTS, start=1):
        feats = [f for f in FEATURES if f["sprint"] == code]
        fill = LIGHT if i % 2 else WHITE
        is_ip = "IP" in code
        cell(t, i, 0, "Sprint " + code, 10.5, DARK if not is_ip else MUTED, True, PP_ALIGN.CENTER, fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 1, dates, 9.5, TEXT, False, PP_ALIGN.CENTER, fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 2, str(len(feats)), 10.5, DARK if feats else MUTED, True, PP_ALIGN.CENTER, fill, MSO_ANCHOR.MIDDLE)
        if feats:
            items = [("%s  ·  %s   —  CO6 due %s" % (f["id"], f["name"], co6_short(f["co6"])), WS[f["ws"]][2], True) for f in feats]
            cell_multi(t, i, 3, items, 9.5, PP_ALIGN.LEFT, fill)
        else:
            cell(t, i, 3, "— no features planned —", 9.5, MUTED, False, PP_ALIGN.LEFT, fill, MSO_ANCHOR.MIDDLE, italic=True)
    textbox(s, Inches(0.55), Inches(6.75), Inches(12.2), Inches(0.3),
            [("Features span Sprints 3.1 - 3.5 (2-3 each); the 3.6 IP iteration holds none. Every feature completes in a sprint ending before its CO6 delivery date.",
              9.5, MUTED, False, True)])


def add_commitments(prs, blank):
    s = prs.slides.add_slide(blank)
    header(s, "Delivery Commitments - CO6 Gates",
           "Every feature delivers a CO6 acceptance criterion; its completion date is the delivery commitment")
    rowplan = [("head", None)]
    for gate in GATE_ORDER:
        gfeats = [f for f in FEATURES if f["co6"] == gate]
        rowplan.append(("group", (gate, gfeats)))
        for f in gfeats:
            rowplan.append(("feat", f))
    t = s.shapes.add_table(len(rowplan), 4, Inches(0.55), Inches(1.3), Inches(12.2), Inches(5.0)).table
    for j, wdt in enumerate([Inches(1.05), Inches(7.15), Inches(2.1), Inches(1.9)]):
        t.columns[j].width = wdt
    for r, (kind, data) in enumerate(rowplan):
        if kind == "head":
            for j, h in enumerate(["ADO ID", "Feature", "Objective", "Scheduled sprint"]):
                cell(t, r, j, h, 10.5, WHITE, True, PP_ALIGN.LEFT if j == 1 else PP_ALIGN.CENTER, DARK)
        elif kind == "group":
            gate, gfeats = data
            t.cell(r, 0).merge(t.cell(r, 3))
            label = ("%s   ·   %d feature%s   ·   %s"
                     % (gate, len(gfeats), "" if len(gfeats) == 1 else "s", GATE_NOTE[gate]))
            cell(t, r, 0, label, 10.5, WHITE, True, PP_ALIGN.LEFT, co6_color(gate), MSO_ANCHOR.MIDDLE)
        else:
            f = data
            cell(t, r, 0, f["id"], 9.5, ACCENT, True, PP_ALIGN.CENTER, LIGHT, MSO_ANCHOR.MIDDLE)
            cell(t, r, 1, f["name"], 9.5, TEXT, False, PP_ALIGN.LEFT, LIGHT, MSO_ANCHOR.MIDDLE)
            cell(t, r, 2, WS[f["ws"]][0], 9.5, WS[f["ws"]][2], True, PP_ALIGN.LEFT, LIGHT, MSO_ANCHOR.MIDDLE)
            cell(t, r, 3, "Sprint " + f["sprint"], 9.5, DARK, True, PP_ALIGN.CENTER, LIGHT, MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.5), [
        ("CO6 check: every feature is scheduled in a sprint that ends before its delivery date - the ADO plan fits within CO6.", 10, GREEN, True, False),
        ("Oct 30 gates fall after PI-3 close (Oct 27); those features complete by Sprint 3.5 (Oct 13), leaving the 3.6 IP iteration as validation buffer.", 9.5, MUTED, False, True),
    ])


def add_board(prs, blank):
    s = prs.slides.add_slide(blank)
    header(s, "PI-3 Program Board", "Features across the 6 sprints  ·  color-coded by objective")
    x0 = Inches(0.55); usable = 12.2; gap = 0.12
    colw = (usable - 5 * gap) / 6.0
    top = 1.35
    hdr_h = 0.62
    chip_top = top + hdr_h + 0.12
    chip_h = 0.78; chip_gap = 0.08
    for ci, (code, dates) in enumerate(SPRINTS):
        cx = 0.55 + ci * (colw + gap)
        is_ip = "IP" in code
        # column header
        rect(s, Inches(cx), Inches(top), Inches(colw), Inches(hdr_h), MUTED if is_ip else DARK)
        textbox(s, Inches(cx), Inches(top + 0.04), Inches(colw), Inches(hdr_h - 0.06),
                [("Sprint " + code, 12, WHITE, True, False), (dates, 8, RGBColor(0xC9, 0xD6, 0xE6), False, False)],
                MSO_ANCHOR.MIDDLE, PP_ALIGN.CENTER)
        feats = [f for f in FEATURES if f["sprint"] == code]
        if not feats:
            rect(s, Inches(cx), Inches(chip_top), Inches(colw), Inches(chip_h), LIGHT, line=RULE)
            textbox(s, Inches(cx), Inches(chip_top), Inches(colw), Inches(chip_h),
                    [("— none —", 9, MUTED, False, True)], MSO_ANCHOR.MIDDLE, PP_ALIGN.CENTER)
            continue
        for fi, f in enumerate(feats):
            cy = chip_top + fi * (chip_h + chip_gap)
            vivid = WS[f["ws"]][1]
            rect(s, Inches(cx), Inches(cy), Inches(colw), Inches(chip_h), vivid, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            textbox(s, Inches(cx + 0.06), Inches(cy + 0.03), Inches(colw - 0.12), Inches(chip_h - 0.06),
                    [(f["short"], 8, WHITE, True, False),
                     ("ADO " + f["id"], 7, RGBColor(0xE4, 0xEC, 0xF6), False, False),
                     ("CO6 due " + co6_short(f["co6"]), 7.5, WHITE, True, False)],
                    MSO_ANCHOR.MIDDLE, PP_ALIGN.LEFT)
    # legend
    ly = 6.62
    lx = 0.55
    textbox(s, Inches(lx), Inches(ly - 0.02), Inches(1.0), Inches(0.3), [("Objective:", 9.5, TEXT, True, False)], MSO_ANCHOR.MIDDLE)
    lx += 1.05
    for key in ["NET", "SM", "CERT", "QLY"]:
        name, vivid, _ = WS[key]
        rect(s, Inches(lx), Inches(ly), Inches(0.24), Inches(0.24), vivid)
        textbox(s, Inches(lx + 0.3), Inches(ly - 0.03), Inches(2.6), Inches(0.3), [(name, 9.5, TEXT, False, False)], MSO_ANCHOR.MIDDLE)
        lx += 2.85


def build_deck():
    prs = Presentation(); prs.slide_width = EMU_W; prs.slide_height = EMU_H
    blank = prs.slide_layouts[6]
    add_title(prs, blank)
    add_overview(prs, blank)
    add_objective_view(prs, blank)
    add_sprint_view(prs, blank)
    add_commitments(prs, blank)
    add_board(prs, blank)
    return prs


def save_safe(prs, path):
    try:
        prs.save(path); print("Saved", os.path.basename(path), "-", len(prs.slides._sldIdLst), "slides")
    except PermissionError:
        base, ext = os.path.splitext(path)
        alt = base + "-new" + ext
        prs.save(alt); print("Target locked - saved to", os.path.basename(alt))


here = os.path.dirname(os.path.abspath(__file__))
save_safe(build_deck(), os.path.join(here, "PI3-Objectives-Features-" + EDITION + ".pptx"))
