"""
PI-3 Objectives deck - FOUR proposed objectives (Network Gear, Service Mapping,
CMDB Governance, Qualys Integration).

Each objective is SMART, mapped to its CO6 v3 deliverable, and shows the v3
acceptance criteria (verbatim, with target dates).

Source: authoritative CO6 v3 ("AP00105-6 ... Release 2 and 3 Deployment v3.docx").
NO contract-commercial data.  Theme mirrors the CO6 / PI-3 decks (a set).

Run:  python build_pi3_objectives_4_deck.py
Output:  PI3-Objectives-4-deliverables-<EDITION>.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

EDITION   = "2026-07-12"
TITLE     = "PI-3 Objectives - CMDB / CSDM"
SUBTITLE  = "Four proposed objectives  ·  SMART, mapped to CO6 v3 deliverables"
AUTHOR    = "Manuel Vazquez  -  Scrum Master, CMDB-CSDM"
FOOTLABEL = "PPL CMDB-CSDM  |  PI-3 Objectives (proposed)  |  CO6 v3  -  Jul 12, 2026"
PLANNING_NOTE = "Proposed for PI Planning - Business Value set by Business Owners.  Grounded in the authoritative CO6 v3."

# Overview: (num, name, deliverable, measurable, gate, commit_word, commit_color)
OVERVIEW = [
    ("1", "Automate Network Device Discovery", "v3 #1  Network Gear", "0 failed auth · >=90% coverage", "Oct 30", "Mixed", "AMBER"),
    ("2", "Build & Validate Service Maps", "v3 #2  Service Mapping", "15 Silver + Gold · >=75% inventory", "Oct 30", "Stretch", "AMBER"),
    ("3", "Certify CMDB Data (4 CI classes)", "v3 #4  CMDB Governance", "4 CI classes certified", "Oct 30", "Committed", "GREEN"),
    ("4", "Activate ServiceNow->Qualys Sync", "v3 #3  Qualys Integration", "Live in PROD · scheduled sync", "Sep 30", "Committed*", "GREEN"),
]
OVERVIEW_FOOT = ("* Obj 4 committable only if the Qualys plug-in is confirmed.  Mixed = later gate is a stretch (Obj 1 90% coverage).  "
                 + PLANNING_NOTE)

# Objectives. num, name, maps(subtitle), statement, smart{}, a_risk, acs=[(bullets[],date)], dependency
OBJECTIVES = [
    dict(num="1", name="Automate Network Device Discovery in the CMDB",
         maps="Maps to CO6 v3 Deliverable #1 (Network Gear Discovery)   ·   Committed to Sep 30; 90% = Stretch",
         statement="Enable automated discovery of all in-scope network device types with validated credentials (0 failed auth) and active schedules by Aug 31, 100% of mandatory attributes by Sep 30, and >=90% of the business-owner-validated device inventory in the CMDB by Oct 30.",
         smart=dict(
             s="Named device types (routers, switches, firewalls, load balancers, WAPs, controllers); creds, schedules, attributes, coverage.",
             m="0 failed authentications; 0 blank mandatory attributes; >=90% vs an owner-validated inventory.",
             a="At risk - credential distribution is an unsolved CO6 risk; 90% by Oct 30 is aggressive.",
             r="CMDB completeness and incident impact analysis.",
             t="Aug 31 -> Sep 30 -> Oct 30 (Oct 30 lands just after PI-3 ends Oct 27)."),
         a_risk=True,
         acs=[
             (["Device types discovered via validated credentials - no failed authentications",
               "Discovery schedules active on defined intervals; CMDB auto-updated each cycle"], "Aug 31, 2026"),
             (["All mandatory CMDB attributes populated via discovery - none blank"], "Sep 30, 2026"),
             ([">=90% of expected inventory in the CMDB, validated by business owners"], "Oct 30, 2026"),
         ],
         dependency="Prove credential distribution on a pilot subnet before committing 90%. Define the coverage denominator (authoritative device list)."),
    dict(num="2", name="Build & Validate Service Maps for Priority Business Apps",
         maps="Maps to CO6 v3 Deliverable #2 (Service Mapping)   ·   Stretch (app-owner dependency)",
         statement="Deliver owner-validated, ServiceNow-consumable service maps - 15 Silver-tier apps (app->infra) and >=75% of the business-service inventory by Aug 31, Gold-tier (service->app) maps by Sep 30, and the remaining Contractor-managed Silver-tier maps by Oct 30.",
         smart=dict(
             s="Tiered targets; two mapping layers (app->infra, service->app); owner validation; consumable in ServiceNow.",
             m="15 Silver-tier maps; Gold-tier maps; >=75% business-service inventory.",
             a="Highest risk - depends on app-owner availability and business services being defined before start.",
             r="Faster incident triage and change-impact assessment; lower MTTR.",
             t="Aug 31 / Sep 30 / Oct 30."),
         a_risk=True,
         acs=[
             (["15 Silver-tier maps (business app -> infrastructure), owner-validated, teams enabled",
               ">=75% business-service inventory represented, owner-validated  (*assumes services defined prior to start)"], "Aug 31, 2026"),
             (["15 Silver-tier maps: additional validation pass",
               "Gold-tier maps (business service -> business application), owner-validated, teams enabled"], "Sep 30, 2026"),
             (["Remaining Contractor-managed Silver-tier maps (app -> infra)",
               "Silver-tier maps (business service -> business application)"], "Oct 30, 2026"),
         ],
         dependency="Secure the app-owner list before Aug 5. Consider committing 8-10 Silver as the first gate; full 15 + Gold as stretch."),
    dict(num="3", name="Certify CMDB Data Across the Four Major CI Classes",
         maps="Maps to CO6 v3 Deliverable #4 (CMDB Governance)   ·   Committed",
         statement="Stand up an end-to-end data certification process - documented flow, technical build in production, tracking dashboards, and delivered training - for Computers and Servers by Sep 30 and Databases and Network Devices by Oct 30.",
         smart=dict(
             s="4 CI classes; each = process flow + PROD build + dashboards + training.",
             m="4 classes fully certified (4 artifacts each).",
             a="Depends on PI-2 carryover: Data Dictionary CCB approval (Jul 21) + audit-dashboard spikes.",
             r="Audit-readiness and a trustworthy CMDB for change / incident / regulatory decisions.",
             t="Computers & Servers Sep 30; Databases & Network Devices Oct 30."),
         a_risk=True,
         acs=[
             (["Data certification for Computers - process flow, PROD build, dashboards, training",
               "Data certification for Servers - same"], "Sep 30, 2026"),
             (["Data certification for Databases - same",
               "Data certification for Network Devices - same"], "Oct 30, 2026"),
         ],
         dependency="Close the PI-2 carryover gates before PI-3 starts, or Sprint-1 is blocked."),
    dict(num="4", name="Activate the ServiceNow -> Qualys Attribute Sync in Production",
         maps="Maps to CO6 v3 Deliverable #3 (Qualys Integration)   ·   Committed (contingent on plug-in)",
         statement="Configure, test, and deploy to production a one-way ServiceNow -> Qualys integration that automatically synchronizes CI owner, support group, and SOX flag on a defined schedule with zero manual handoffs, by Sep 30.",
         smart=dict(
             s="One integration; direction ServiceNow -> Qualys; three fields (owner, support group, SOX flag).",
             m="Live in production; scheduled sync; 0 manual handoffs.",
             a="Contingent on the Qualys plug-in; stories 1428703/1428704 are blocked and scoped for the old direction.",
             r="Ties vulnerability data to the right owners and compliance context; audit readiness.",
             t="Sep 30 - single mid-PI gate."),
         a_risk=True,
         acs=[
             (["One-way ServiceNow -> Qualys integration configured, tested, live in PROD; owner, support group & SOX flag auto-synchronized on a defined schedule, no manual handoffs  (*assumes the required Qualys plug-in is available)"], "Sep 30, 2026"),
         ],
         dependency="Confirm the Qualys plug-in is available; re-scope stories 1428703 / 1428704 to the ServiceNow -> Qualys direction."),
]

# Gate timeline: (when, what's due, objectives)
GATES = [
    ("Aug 31", "Network Gear (creds + schedules)  ·  Service Mapping (15 Silver + >=75% inventory)", "Obj 1, 2"),
    ("Sep 30", "Qualys live in PROD  ·  Data cert (Computers, Servers)  ·  Network attributes  ·  Service Mapping (Gold + Silver pass)", "Obj 1, 2, 3, 4"),
    ("Oct 27", "PI-3 ends", "-"),
    ("Oct 30", "Network 90%  ·  Service Mapping (remaining Silver + service->app)  ·  Data cert (Databases, Network)", "Obj 1, 2, 3"),
]
GATES_NOTE = "The Oct 30 gates land ~3 days after PI-3 ends (Oct 27) - plan them into the final sprint / IP iteration."

# Dependencies: (item, threatens, owner)
DEPS = [
    ("Qualys plug-in confirmed available", "Obj 4 (stories 1428703/04 blocked)", "Stan / Rich Santillo"),
    ("Re-scope Qualys stories to ServiceNow -> Qualys direction", "Obj 4", "Stan"),
    ("App-owner list secured (Silver / Gold tiers)", "Obj 2", "Tanzeel / Joe"),
    ("PI-2 carryover: Data Dictionary CCB approval (Jul 21) + audit spikes", "Obj 3", "Manuel / Josh Sterling"),
    ("Credential-distribution approach proven on a pilot", "Obj 1 (90% coverage)", "Stan"),
]

# --------------------------------------------------------------------------- THEME
DARK   = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0x2E, 0x6D, 0xB4)
LIGHT  = RGBColor(0xF4, 0xF6, 0xF9)
TEXT   = RGBColor(0x2D, 0x2D, 0x2D)
MUTED  = RGBColor(0x6B, 0x72, 0x80)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
AMBER  = RGBColor(0xE0, 0x8A, 0x00)
RED    = RGBColor(0xC0, 0x39, 0x2B)
RULE   = RGBColor(0xDD, 0xE3, 0xEA)
ACCENT_TINT = RGBColor(0xE9, 0xF0, 0xF8)
AMBER_FILL  = RGBColor(0xFB, 0xEE, 0xDB)
AMBER_BODY  = RGBColor(0x5A, 0x4A, 0x2A)
RAG = {"GREEN": GREEN, "AMBER": AMBER, "RED": RED}
EMU_W, EMU_H = Inches(13.333), Inches(7.5)
FONT = "Segoe UI"


def _set(run, size, color, bold=False, italic=False):
    run.font.size = Pt(size); run.font.color.rgb = color
    run.font.bold = bold; run.font.italic = italic; run.font.name = FONT


def rect(slide, x, y, w, h, color, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def textbox(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(4); tf.margin_top = tf.margin_bottom = Pt(2)
    for i, spec in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        text, size, color, bold, italic = spec
        r = p.add_run(); r.text = text; _set(r, size, color, bold, italic)
    return tb


def header(slide, title, subtitle=None):
    rect(slide, 0, 0, EMU_W, Inches(1.05), DARK)
    rect(slide, 0, Inches(1.05), EMU_W, Inches(0.06), ACCENT)
    textbox(slide, Inches(0.55), Inches(0.16), Inches(12.2), Inches(0.55), [(title, 23, WHITE, True, False)], MSO_ANCHOR.MIDDLE)
    if subtitle:
        textbox(slide, Inches(0.57), Inches(0.66), Inches(12.2), Inches(0.32), [(subtitle, 12, RGBColor(0xC9, 0xD6, 0xE6), False, False)], MSO_ANCHOR.MIDDLE)
    textbox(slide, Inches(0.55), Inches(7.08), Inches(10), Inches(0.32), [(FOOTLABEL, 9, MUTED, False, False)])


def cell(t, i, j, lines, size, color, bold=False, align=PP_ALIGN.LEFT, fill=None, anchor=None):
    c = t.cell(i, j)
    if fill is not None:
        c.fill.solid(); c.fill.fore_color.rgb = fill
    if anchor is not None:
        c.vertical_anchor = anchor
    tf = c.text_frame; tf.word_wrap = True
    if isinstance(lines, str):
        lines = [lines]
    for k, ln in enumerate(lines):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.alignment = align; p.space_before = Pt(0); p.space_after = Pt(2)
        r = p.add_run(); r.text = ln; _set(r, size, color, bold)


def ac_table(slide, x, y, w, h, acs):
    t = slide.shapes.add_table(len(acs) + 1, 2, Inches(x), Inches(y), Inches(w), Inches(h)).table
    t.columns[0].width = Inches(w - 2.6); t.columns[1].width = Inches(2.6)
    cell(t, 0, 0, "Acceptance criteria (CO6 v3)", 10.5, WHITE, True, PP_ALIGN.LEFT, DARK)
    cell(t, 0, 1, "Target date", 10.5, WHITE, True, PP_ALIGN.CENTER, DARK)
    for i, (bullets, date) in enumerate(acs, start=1):
        fill = LIGHT if i % 2 else WHITE
        disp = bullets if len(bullets) == 1 else ["•  " + b for b in bullets]
        cell(t, i, 0, disp, 9, TEXT, False, PP_ALIGN.LEFT, fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 1, date, 9.5, ACCENT, True, PP_ALIGN.CENTER, fill, MSO_ANCHOR.MIDDLE)


# --------------------------------------------------------------------------- SLIDES
def add_title(prs, blank):
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, EMU_W, EMU_H, DARK)
    rect(s, 0, Inches(4.55), EMU_W, Inches(0.08), ACCENT)
    textbox(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(1.1), [(TITLE, 40, WHITE, True, False)])
    textbox(s, Inches(0.82), Inches(2.75), Inches(11.7), Inches(0.6), [(SUBTITLE, 19, RGBColor(0xC9, 0xD6, 0xE6), False, False)])
    textbox(s, Inches(0.82), Inches(4.75), Inches(11.7), Inches(1.5), [
        ("Network Gear  ·  Service Mapping  ·  CMDB Governance  ·  Qualys Integration", 15, WHITE, True, False),
        ("PI-3 window: Aug 5 - Oct 27, 2026   ·   for PI Planning (Jul 22 - Aug 4)", 14, RGBColor(0xC9, 0xD6, 0xE6), False, False),
        (AUTHOR, 12, RGBColor(0x9F, 0xB2, 0xCC), False, False),
    ])
    textbox(s, Inches(0.82), Inches(6.5), Inches(11.7), Inches(0.4), [(PLANNING_NOTE, 11, RGBColor(0x7F, 0x8C, 0xA6), False, True)])


def add_overview(prs, blank):
    s = prs.slides.add_slide(blank)
    header(s, "PI-3 Objectives - Overview", "Four objectives, each mapped to a CO6 v3 deliverable")
    rows = len(OVERVIEW) + 1
    t = s.shapes.add_table(rows, 6, Inches(0.55), Inches(1.5), Inches(12.2), Inches(2.9)).table
    for j, wdt in enumerate([Inches(0.5), Inches(3.5), Inches(2.7), Inches(3.1), Inches(1.0), Inches(1.4)]):
        t.columns[j].width = wdt
    for j, h in enumerate(["#", "Objective", "CO6 v3 deliverable", "Measurable target", "Gate", "Commit"]):
        cell(t, 0, j, h, 11, WHITE, True, PP_ALIGN.LEFT if j in (1, 2, 3) else PP_ALIGN.CENTER, DARK)
    for i, (num, name, deliv, meas, gate, cw, cc) in enumerate(OVERVIEW, start=1):
        fill = LIGHT if i % 2 else WHITE
        cell(t, i, 0, num, 10.5, MUTED, True, PP_ALIGN.CENTER, fill)
        cell(t, i, 1, name, 10.5, DARK, True, PP_ALIGN.LEFT, fill)
        cell(t, i, 2, deliv, 9.5, ACCENT, True, PP_ALIGN.LEFT, fill)
        cell(t, i, 3, meas, 9.5, TEXT, False, PP_ALIGN.LEFT, fill)
        cell(t, i, 4, gate, 10, RED if gate == "Oct 30" else ACCENT, True, PP_ALIGN.CENTER, fill)
        cell(t, i, 5, cw, 9.5, WHITE, True, PP_ALIGN.CENTER, RAG[cc])
    textbox(s, Inches(0.55), Inches(4.7), Inches(12.2), Inches(1.0), [(OVERVIEW_FOOT, 9.5, TEXT, False, False)])


def add_objective(prs, blank, o):
    s = prs.slides.add_slide(blank)
    header(s, "Objective " + o["num"] + "  -  " + o["name"], o["maps"])
    # statement strip
    rect(s, Inches(0.55), Inches(1.2), Inches(12.2), Inches(0.6), ACCENT_TINT, line=ACCENT)
    textbox(s, Inches(0.72), Inches(1.23), Inches(11.9), Inches(0.54), [(o["statement"], 11, DARK, False, True)], MSO_ANCHOR.MIDDLE)
    # SMART table
    sm = o["smart"]
    rowspec = [("S  ·  Specific", sm["s"]), ("M  ·  Measurable", sm["m"]), ("A  ·  Achievable", sm["a"]),
               ("R  ·  Relevant", sm["r"]), ("T  ·  Time-bound", sm["t"])]
    t = s.shapes.add_table(5, 2, Inches(0.55), Inches(1.92), Inches(12.2), Inches(1.85)).table
    t.columns[0].width = Inches(1.95); t.columns[1].width = Inches(10.25)
    for i, (lab, val) in enumerate(rowspec):
        fill = LIGHT if i % 2 else WHITE
        is_a = lab.startswith("A")
        cell(t, i, 0, lab, 10, DARK, True, PP_ALIGN.LEFT, fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 1, val, 9.5, RED if (is_a and o["a_risk"]) else TEXT, is_a and o["a_risk"], PP_ALIGN.LEFT, fill, MSO_ANCHOR.MIDDLE)
    # acceptance criteria table
    ac_table(s, 0.55, 3.86, 12.2, 0.4 * (len(o["acs"]) + 1), o["acs"])
    # dependency strip
    rect(s, Inches(0.55), Inches(5.95), Inches(12.2), Inches(0.92), AMBER_FILL, line=AMBER)
    textbox(s, Inches(0.72), Inches(5.99), Inches(11.9), Inches(0.84),
            [("Lock before planning", 10.5, AMBER_BODY, True, False), (o["dependency"], 10, AMBER_BODY, False, False)], MSO_ANCHOR.MIDDLE)


def add_timeline(prs, blank):
    s = prs.slides.add_slide(blank)
    header(s, "Gate Timeline - PI-3", "When each objective's acceptance lands")
    rows = len(GATES) + 1
    t = s.shapes.add_table(rows, 3, Inches(0.55), Inches(1.5), Inches(12.2), Inches(2.7)).table
    for j, wdt in enumerate([Inches(1.5), Inches(8.6), Inches(2.1)]):
        t.columns[j].width = wdt
    for j, h in enumerate(["When", "What's due", "Objectives"]):
        cell(t, 0, j, h, 11.5, WHITE, True, PP_ALIGN.LEFT if j == 1 else PP_ALIGN.CENTER, DARK)
    for i, (when, what, objs) in enumerate(GATES, start=1):
        end = when == "Oct 27"
        fill = AMBER_FILL if end else (LIGHT if i % 2 else WHITE)
        cell(t, i, 0, when, 10.5, RED if end else DARK, True, PP_ALIGN.CENTER, fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 1, what, 9.5, AMBER_BODY if end else TEXT, end, PP_ALIGN.LEFT, fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 2, objs, 9.5, DARK, True, PP_ALIGN.CENTER, fill, MSO_ANCHOR.MIDDLE)
    rect(s, Inches(0.55), Inches(4.5), Inches(12.2), Inches(0.8), AMBER_FILL, line=AMBER)
    textbox(s, Inches(0.72), Inches(4.55), Inches(11.9), Inches(0.7), [(GATES_NOTE, 11, AMBER_BODY, False, False)], MSO_ANCHOR.MIDDLE)


def add_deps(prs, blank):
    s = prs.slides.add_slide(blank)
    header(s, "Lock These Before PI Planning", "Each turns an objective from aspirational to genuinely Achievable")
    rows = len(DEPS) + 1
    t = s.shapes.add_table(rows, 3, Inches(0.55), Inches(1.5), Inches(12.2), Inches(3.4)).table
    for j, wdt in enumerate([Inches(6.2), Inches(3.2), Inches(2.8)]):
        t.columns[j].width = wdt
    for j, h in enumerate(["Lock before Aug 5 / PI Planning", "Threatens", "Owner"]):
        cell(t, 0, j, h, 11.5, WHITE, True, PP_ALIGN.LEFT, DARK)
    for i, (item, threat, owner) in enumerate(DEPS, start=1):
        fill = LIGHT if i % 2 else WHITE
        cell(t, i, 0, item, 10.5, DARK, False, PP_ALIGN.LEFT, fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 1, threat, 10, RED, True, PP_ALIGN.LEFT, fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 2, owner, 10, TEXT, False, PP_ALIGN.LEFT, fill, MSO_ANCHOR.MIDDLE)
    rect(s, Inches(0.55), Inches(5.5), Inches(12.2), Inches(0.95), ACCENT_TINT, line=ACCENT)
    textbox(s, Inches(0.72), Inches(5.56), Inches(11.9), Inches(0.85),
            [("A commitment isn't real until its dependency is closed - resolve these before committing at PI Planning.", 10.5, DARK, False, False)], MSO_ANCHOR.MIDDLE)


def build_deck():
    prs = Presentation(); prs.slide_width = EMU_W; prs.slide_height = EMU_H
    blank = prs.slide_layouts[6]
    add_title(prs, blank)
    add_overview(prs, blank)
    for o in OBJECTIVES:
        add_objective(prs, blank, o)
    add_timeline(prs, blank)
    add_deps(prs, blank)
    return prs


def save_safe(prs, path):
    try:
        prs.save(path); print("Saved", os.path.basename(path), "-", len(prs.slides._sldIdLst), "slides")
    except PermissionError:
        base, ext = os.path.splitext(path)
        alt = base + "-new" + ext
        prs.save(alt); print("Target locked - saved to", os.path.basename(alt))


here = os.path.dirname(os.path.abspath(__file__))
save_safe(build_deck(), os.path.join(here, "PI3-Objectives-4-deliverables-" + EDITION + ".pptx"))
