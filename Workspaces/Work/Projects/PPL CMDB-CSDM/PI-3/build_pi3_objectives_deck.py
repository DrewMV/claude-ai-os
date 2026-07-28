"""
PI-3 Objectives deck - proposed SMART objectives for PI-3 planning.

Source: PI-3/pi3-objectives.md (grounded in the authoritative CO6 v3,
"AP00105-6 to ServiceNow Release 2 and 3 Deployment v3.docx").
Each objective slide gives a SMART breakdown (Specific / Measurable / Achievable /
Relevant / Time-bound), acceptance gates with dates, committed-vs-stretch, and the
dependency to lock before PI planning.

NO contract-commercial data (fees excluded - not relevant to objectives).

Theme mirrors the CO6 deck so the decks look like a set.
Single source of truth. Edit CONTENT, then run:
    python build_pi3_objectives_deck.py
Output:  PI3-Objectives-<EDITION>.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# CONTENT
# ---------------------------------------------------------------------------
EDITION   = "2026-07-11"
TITLE     = "PI-3 Objectives - CMDB / CSDM"
SUBTITLE  = "Proposed objectives for PI Planning  ·  SMART, CO6-aligned"
AUTHOR    = "Manuel Vazquez  -  Scrum Master, CMDB-CSDM"
FOOTLABEL = "PPL CMDB-CSDM  |  PI-3 Objectives (proposed)  |  Jul 11, 2026"
PLANNING_NOTE = ("Proposed for planning discussion - Business Value set by Business Owners; "
                 "committed vs. stretch confirmed at PI Planning.  Grounded in CO6 v3 (authoritative).")

# Overview: (#, name, measurable, gate, commit_word, commit_color, priority)
OVERVIEW = [
    ("1", "Sustain ServiceNow Platform Operations", "40 pts/member/sprint · ATF plan", "Oct 30", "Committed", "GREEN", "P0"),
    ("2", "Activate Qualys Integration in Production", "Live in PROD · scheduled sync", "Sep 30", "Committed*", "GREEN", "P1"),
    ("3", "Expand CI Data Certification (all classes)", "4 CI classes certified", "Oct 30", "Committed", "GREEN", "P1"),
    ("4", "Automate Network Device Discovery", "0 failed auth · >=90% coverage", "Oct 30", "Mixed", "AMBER", "P1"),
    ("5", "Build Service Maps (priority apps)", "15 Silver + Gold · >=75% inventory", "Oct 30", "Stretch", "AMBER", "P2"),
    ("6", "Deliver NERC-CIP Platform Strategy", "Evaluation + exec package", "Oct 30", "Committed", "GREEN", "P2"),
    ("7", "Define & Begin Legacy Rationalization", "4-platform analysis · plan · Oct exec", "Oct 30", "Mixed", "AMBER", "P3"),
]
OVERVIEW_FOOT = ("* Obj 2 committable only if the Qualys plugin is confirmed available.  "
                 "Mixed = the later gate is a stretch (Obj 4 90% coverage; Obj 7 Oct execution).  " + PLANNING_NOTE)

# Objectives. num, name, sub(priority+commit), statement, smart{}, a_risk, gates, dependency
OBJECTIVES = [
    dict(num="1", name="Sustain ServiceNow Platform Operations",
         sub="Priority P0   ·   Committed",
         statement="Sustain BAU delivery at 40 completed story points per team member per sprint across every PI-3 sprint, publish monthly sprint-allocation reports, and publish an ATF rollout plan by Oct 30.",
         smart=dict(
             s="BAU support + ITSM PO coverage + ATF plan, scoped to production ITSM / CMDB modules.",
             m="40 pts/member/sprint (PTO-adjusted); monthly allocation reports; 1 ATF plan (all prod capabilities in use as of Jul 31).",
             a="Achievable if the Platform Support team is staffed at PI start with a PI-2 velocity baseline.",
             r="The delivery engine for the whole PI - if capacity drops, every objective slips. (CO6 Platform Support + ITSM PO + ATF.)",
             t="Per sprint through PI-3; ATF plan by Oct 30."),
         a_risk=False,
         gates="Per sprint - 40 pts/member  ·  Monthly - allocation reports  ·  Oct 30 - ATF rollout plan",
         dependency="Pin '40 pts' to your actual estimation unit; lock the team roster + PI-2 velocity baseline before planning."),
    dict(num="2", name="Activate Qualys Integration in Production",
         sub="Priority P1   ·   Committed (contingent on Qualys plugin)",
         statement="Configure, test, and deploy the one-way ServiceNow -> Qualys integration to production by Sep 30, syncing CI owner, support group, and SOX flag on a defined schedule with zero manual handoffs.",
         smart=dict(
             s="One integration, explicit direction (ServiceNow -> Qualys) and fields (owner, support group, SOX flag).",
             m="Live in PROD; scheduled sync; 0 manual handoffs; 0 data-loss / transformation errors.",
             a="AT RISK - stories 1428703 / 1428704 are blocked on the Qualys vendor plugin.",
             r="Highest contractual visibility; carries CO5 -> CO6; security / vulnerability posture.",
             t="Sep 30 - single mid-PI gate, leaving recovery room."),
         a_risk=True,
         gates="Sep 30 - integration live in PROD, scheduled sync, 0 manual handoffs",
         dependency="Confirm the Qualys plugin is available (Stan / Rich Santillo) before planning. If unconfirmed, split: non-prod build (commit) vs. PROD go-live (stretch). Reconcile direction/data against CO6 v3."),
    dict(num="3", name="Expand CI Data Certification Across All Major CI Classes",
         sub="Priority P1   ·   Committed",
         statement="Stand up data certification (process, PROD build, dashboards, training) for Computers and Servers by Sep 30 and Databases and Network Devices by Oct 30, all against the approved data dictionary.",
         smart=dict(
             s="4 named CI classes; each = process definition + PROD build + dashboard + training.",
             m="4 classes certified; per-class 4 artifacts; coverage evidenced by the audit dashboards.",
             a="Achievable only if PI-2 carryover closes first (CCB approval, audit spikes, 90% stories).",
             r="Re-baselines CO5 D1; audit-readiness; core CMDB health.",
             t="Computers & Servers Sep 30; Databases & Network Devices Oct 30."),
         a_risk=True,
         gates="Sep 30 - Computers & Servers  ·  Oct 30 - Databases & Network Devices",
         dependency="Entry criteria: Data Dictionary CCB approval (Jul 21) + 3 audit-dashboard spikes + 90%-coverage stories. If CCB slips, Sprint-1 is blocked."),
    dict(num="4", name="Automate Network Device Discovery in CMDB",
         sub="Priority P1   ·   Committed to Sep 30; 90% coverage is Stretch",
         statement="Discover all in-scope network device types with validated credentials (0 failed auth) and active schedules by Aug 31, populate 100% of mandatory attributes by Sep 30, and reach >=90% coverage validated by business owners by Oct 30.",
         smart=dict(
             s="In-scope device types (routers, switches, firewalls, load balancers, WAPs, controllers).",
             m="0 failed authentications; 0 blank mandatory attributes; >=90% vs an owner-validated inventory.",
             a="Credential distribution is an unresolved CO6 RISK - 90% by Oct 30 is aggressive.",
             r="CMDB completeness + incident response; foundational.",
             t="Aug 31 (creds/schedules) -> Sep 30 (attrs) -> Oct 30 (90%)."),
         a_risk=True,
         gates="Aug 31 - creds validated + schedules active  ·  Sep 30 - mandatory attrs 100%  ·  Oct 30 - >=90% coverage",
         dependency="Prove credential distribution on a pilot subnet before committing 90%. Define the coverage denominator (whose authoritative device list?)."),
    dict(num="5", name="Build Service Maps for Priority Business Applications",
         sub="Priority P2   ·   Stretch (app-owner dependency)",
         statement="Build and owner-validate service maps for 15 Silver-tier apps and >=75% of the business-service inventory by Aug 31, Gold-tier apps by Sep 30, and remaining Accenture-managed Silver-tier apps by Oct 30 - all consumable in ServiceNow.",
         smart=dict(
             s="Tiered app targets; dependencies down to infrastructure CIs; owner validation.",
             m="15 Silver + >=75% service inventory (Aug 31); Gold-tier (Sep 30); remaining Silver (Oct 30).",
             a="HIGHEST RISK - CO6 assumes business services defined pre-start; app-owner availability is the gate.",
             r="End-to-end service visibility; impact / incident analysis.",
             t="Aug 31 / Sep 30 / Oct 30."),
         a_risk=True,
         gates="Aug 31 - 15 Silver + >=75% inventory  ·  Sep 30 - Gold-tier  ·  Oct 30 - remaining Silver",
         dependency="Secure the app-owner list before Aug 5 (hard entry criterion). Consider committing 8-10 Silver as the first gate; full 15 + Gold as stretch."),
    dict(num="6", name="Deliver NERC-CIP ServiceNow Platform Strategy",
         sub="Priority P2   ·   Committed",
         statement="Deliver a completed dual-instance NERC-CIP evaluation (considerations, pros/cons, risks, dependencies) by Sep 30 and an executive strategy package accepted by PPL IT leadership by Oct 30.",
         smart=dict(
             s="2 artifacts - dual-instance evaluation + executive strategy package.",
             m="Evaluation delivered Sep 30; exec package delivered and accepted Oct 30.",
             a="Highly achievable - analysis / strategy work, low technical / infra dependency.",
             r="Regulatory (NERC-CIP); executive visibility.",
             t="Sep 30 (evaluation) / Oct 30 (package accepted)."),
         a_risk=False,
         gates="Sep 30 - dual-instance evaluation  ·  Oct 30 - exec strategy package accepted",
         dependency="Name the accepting stakeholder (PPL IT Leadership) now - 'accepted' needs a named approver. Scope is in CO6 v3 but not the 3/27 draft - confirm with Joe / Jordan."),
    dict(num="7", name="Define and Begin Legacy Platform Rationalization",
         sub="Priority P3   ·   Committed to plan; Oct execution is Stretch",
         statement="Complete a gap analysis of iTeam, DISCO, Cherwell, and AIM vs. ServiceNow by Aug 31, deliver a prioritized migration plan + roadmap by Sep 30, and execute the October-scope items defined in that plan by Oct 30.",
         smart=dict(
             s="4 platforms (iTeam, DISCO, Cherwell, AIM); analysis -> plan -> limited execution.",
             m="4-platform gap analysis; 1 prioritized plan; defined Oct-scope items executed.",
             a="Analysis achievable; execution is greenfield (only iTeam import 1452028 today), intentionally limited.",
             r="Retire legacy tooling into ServiceNow.",
             t="Aug 31 (analysis) / Sep 30 (plan) / Oct 30 (Oct-scope execution)."),
         a_risk=False,
         gates="Aug 31 - 4-platform gap analysis  ·  Sep 30 - migration plan + roadmap  ·  Oct 30 - Oct-scope execution",
         dependency="Oct execution scope is defined only by the Sep 30 plan -> keep it stretch. Decide Cherwell / AIM in-or-out at planning."),
]

# Gate timeline: (when, what's due, objectives)
GATES = [
    ("Aug 31", "Network Gear (creds + schedules)  ·  Service Mapping (15 Silver + 75% inventory)  ·  Legacy (gap analysis)", "Obj 4, 5, 7"),
    ("Sep 30", "Qualys live in PROD  ·  Data Cert (Computers, Servers)  ·  Network attrs  ·  Service Mapping (Gold)  ·  Legacy (plan)", "Obj 2, 3, 4, 5, 7"),
    ("Oct 27", "PI-3 ends", "-"),
    ("Oct 30", "ATF plan  ·  Data Cert (Databases, Network)  ·  Network 90%  ·  Service Mapping (remaining Silver)  ·  NERC-CIP package  ·  Legacy (Oct exec)", "Obj 1, 3, 4, 5, 6, 7"),
]
GATES_NOTE = "The Oct 30 CO6 gates land ~3 days after PI-3 ends (Oct 27) - plan them into the final sprint / IP iteration."

# Dependencies to lock: (item, threatens, owner)
DEPS = [
    ("CO6 executed (signature)", "All objectives", "Christian / Aaron Simeon"),
    ("Qualys vendor plugin confirmed available", "Obj 2 (blocked stories 1428703/04)", "Stan / Rich Santillo"),
    ("App-owner list secured (Silver / Gold tiers)", "Obj 5", "Tanzeel / Joe"),
    ("Data Dictionary CCB approval (Jul 21) + audit spikes + 90% stories", "Obj 3", "Manuel / Josh Sterling"),
    ("Credential-distribution approach proven on a pilot", "Obj 4 (90% coverage)", "Stan"),
    ("Named accepting stakeholder for NERC-CIP", "Obj 6", "Joe"),
]

# ---------------------------------------------------------------------------
# THEME  (mirrors the CO6 deck)
# ---------------------------------------------------------------------------
DARK   = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0x2E, 0x6D, 0xB4)
LIGHT  = RGBColor(0xF4, 0xF6, 0xF9)
TEXT   = RGBColor(0x2D, 0x2D, 0x2D)
MUTED  = RGBColor(0x6B, 0x72, 0x80)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
AMBER  = RGBColor(0xE0, 0x8A, 0x00)
RED    = RGBColor(0xC0, 0x39, 0x2B)
PALE   = RGBColor(0xEC, 0xEE, 0xF1)
RULE   = RGBColor(0xDD, 0xE3, 0xEA)
ACCENT_TINT = RGBColor(0xE9, 0xF0, 0xF8)
AMBER_FILL  = RGBColor(0xFB, 0xEE, 0xDB)
AMBER_BODY  = RGBColor(0x5A, 0x4A, 0x2A)
RAG = {"GREEN": GREEN, "AMBER": AMBER, "RED": RED}

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
FONT = "Segoe UI"


def _set(run, size, color, bold=False, italic=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = FONT


def rect(slide, x, y, w, h, color, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def textbox(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, spec in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        text, size, color, bold, italic = spec
        r = p.add_run(); r.text = text
        _set(r, size, color, bold, italic)
    return tb


def header(slide, title, subtitle=None):
    rect(slide, 0, 0, EMU_W, Inches(1.05), DARK)
    rect(slide, 0, Inches(1.05), EMU_W, Inches(0.06), ACCENT)
    textbox(slide, Inches(0.55), Inches(0.16), Inches(12.2), Inches(0.55),
            [(title, 23, WHITE, True, False)], MSO_ANCHOR.MIDDLE)
    if subtitle:
        textbox(slide, Inches(0.57), Inches(0.66), Inches(12.2), Inches(0.32),
                [(subtitle, 12, RGBColor(0xC9, 0xD6, 0xE6), False, False)], MSO_ANCHOR.MIDDLE)
    textbox(slide, Inches(0.55), Inches(7.08), Inches(10), Inches(0.32),
            [(FOOTLABEL, 9, MUTED, False, False)])


def cell(t, i, j, lines, size, color, bold=False, align=PP_ALIGN.LEFT, fill=None, anchor=None):
    c = t.cell(i, j)
    if fill is not None:
        c.fill.solid(); c.fill.fore_color.rgb = fill
    if anchor is not None:
        c.vertical_anchor = anchor
    tf = c.text_frame; tf.word_wrap = True
    if isinstance(lines, str):
        lines = [lines]
    for k, ln in enumerate(lines):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0); p.space_after = Pt(2)
        r = p.add_run(); r.text = ln
        _set(r, size, color, bold)


# ---------------------------------------------------------------------------
# SLIDES
# ---------------------------------------------------------------------------
def add_title(prs, blank):
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, EMU_W, EMU_H, DARK)
    rect(s, 0, Inches(4.55), EMU_W, Inches(0.08), ACCENT)
    textbox(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(1.1), [(TITLE, 40, WHITE, True, False)])
    textbox(s, Inches(0.82), Inches(2.75), Inches(11.7), Inches(0.6), [(SUBTITLE, 20, RGBColor(0xC9, 0xD6, 0xE6), False, False)])
    textbox(s, Inches(0.82), Inches(4.75), Inches(11.7), Inches(1.5), [
        ("7 proposed objectives  ·  SMART  ·  committed & stretch", 16, WHITE, True, False),
        ("PI-3 window: Aug 5 - Oct 27, 2026   ·   for PI Planning (Jul 22 - Aug 4)", 14, RGBColor(0xC9, 0xD6, 0xE6), False, False),
        (AUTHOR, 12, RGBColor(0x9F, 0xB2, 0xCC), False, False),
    ])
    textbox(s, Inches(0.82), Inches(6.45), Inches(11.7), Inches(0.5), [(PLANNING_NOTE, 11, RGBColor(0x7F, 0x8C, 0xA6), False, True)])


def add_overview(prs, blank):
    s = prs.slides.add_slide(blank)
    header(s, "PI-3 Objectives - Overview", "7 objectives  ·  ordered by priority  ·  committed & stretch")
    rows = len(OVERVIEW) + 1
    t = s.shapes.add_table(rows, 6, Inches(0.55), Inches(1.4), Inches(12.2), Inches(4.3)).table
    for j, wdt in enumerate([Inches(0.5), Inches(3.7), Inches(3.9), Inches(1.2), Inches(1.5), Inches(1.4)]):
        t.columns[j].width = wdt
    for j, h in enumerate(["#", "Objective", "Measurable target", "Gate", "Commit", "Priority"]):
        cell(t, 0, j, h, 11, WHITE, True, PP_ALIGN.LEFT if j == 1 else PP_ALIGN.CENTER, DARK)
    for i, (num, name, meas, gate, cw, cc, pri) in enumerate(OVERVIEW, start=1):
        fill = LIGHT if i % 2 else WHITE
        cell(t, i, 0, num, 10.5, MUTED, True, PP_ALIGN.CENTER, fill)
        cell(t, i, 1, name, 10.5, DARK, True, PP_ALIGN.LEFT, fill)
        cell(t, i, 2, meas, 9.5, TEXT, False, PP_ALIGN.LEFT, fill)
        cell(t, i, 3, gate, 10, RED if gate == "Oct 30" else ACCENT, True, PP_ALIGN.CENTER, fill)
        cell(t, i, 4, cw, 9.5, WHITE, True, PP_ALIGN.CENTER, RAG[cc])
        cell(t, i, 5, pri, 10, DARK, True, PP_ALIGN.CENTER, fill)
    textbox(s, Inches(0.55), Inches(5.85), Inches(12.2), Inches(1.0), [(OVERVIEW_FOOT, 9.5, TEXT, False, False)])


def add_objective(prs, blank, o):
    s = prs.slides.add_slide(blank)
    header(s, "Objective " + o["num"] + "  -  " + o["name"], o["sub"])
    # statement strip
    rect(s, Inches(0.55), Inches(1.22), Inches(12.2), Inches(0.66), ACCENT_TINT, line=ACCENT)
    textbox(s, Inches(0.72), Inches(1.25), Inches(11.9), Inches(0.6),
            [(o["statement"], 11.5, DARK, False, True)], MSO_ANCHOR.MIDDLE)
    # SMART table
    sm = o["smart"]
    rowspec = [("S  ·  Specific", sm["s"]), ("M  ·  Measurable", sm["m"]), ("A  ·  Achievable", sm["a"]),
               ("R  ·  Relevant", sm["r"]), ("T  ·  Time-bound", sm["t"])]
    t = s.shapes.add_table(5, 2, Inches(0.55), Inches(2.02), Inches(12.2), Inches(2.35)).table
    t.columns[0].width = Inches(2.0)
    t.columns[1].width = Inches(10.2)
    for i, (lab, val) in enumerate(rowspec):
        fill = LIGHT if i % 2 else WHITE
        is_a = lab.startswith("A")
        cell(t, i, 0, lab, 10.5, DARK, True, PP_ALIGN.LEFT, fill, MSO_ANCHOR.MIDDLE)
        vcolor = RED if (is_a and o["a_risk"]) else TEXT
        vbold = is_a and o["a_risk"]
        cell(t, i, 1, val, 10, vcolor, vbold, PP_ALIGN.LEFT, fill, MSO_ANCHOR.MIDDLE)
    # gates strip
    rect(s, Inches(0.55), Inches(4.55), Inches(12.2), Inches(0.82), LIGHT, line=RULE)
    textbox(s, Inches(0.72), Inches(4.6), Inches(11.9), Inches(0.72),
            [("Acceptance gates", 10.5, DARK, True, False), (o["gates"], 10, TEXT, False, False)], MSO_ANCHOR.MIDDLE)
    # dependency (amber)
    rect(s, Inches(0.55), Inches(5.5), Inches(12.2), Inches(1.15), AMBER_FILL, line=AMBER)
    textbox(s, Inches(0.72), Inches(5.55), Inches(11.9), Inches(1.05),
            [("Lock before planning", 10.5, AMBER_BODY, True, False), (o["dependency"], 10, AMBER_BODY, False, False)], MSO_ANCHOR.MIDDLE)


def add_timeline(prs, blank):
    s = prs.slides.add_slide(blank)
    header(s, "Gate Timeline - PI-3", "When each objective's acceptance lands")
    rows = len(GATES) + 1
    t = s.shapes.add_table(rows, 3, Inches(0.55), Inches(1.45), Inches(12.2), Inches(3.4)).table
    for j, wdt in enumerate([Inches(1.5), Inches(8.6), Inches(2.1)]):
        t.columns[j].width = wdt
    for j, h in enumerate(["When", "What's due", "Objectives"]):
        cell(t, 0, j, h, 11.5, WHITE, True, PP_ALIGN.LEFT if j == 1 else PP_ALIGN.CENTER, DARK)
    for i, (when, what, objs) in enumerate(GATES, start=1):
        end = when == "Oct 27"
        fill = AMBER_FILL if end else (LIGHT if i % 2 else WHITE)
        cell(t, i, 0, when, 10.5, RED if end else DARK, True, PP_ALIGN.CENTER, fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 1, what, 9.5, AMBER_BODY if end else TEXT, end, PP_ALIGN.LEFT, fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 2, objs, 9.5, DARK, True, PP_ALIGN.CENTER, fill, MSO_ANCHOR.MIDDLE)
    rect(s, Inches(0.55), Inches(5.1), Inches(12.2), Inches(0.8), AMBER_FILL, line=AMBER)
    textbox(s, Inches(0.72), Inches(5.15), Inches(11.9), Inches(0.7),
            [(GATES_NOTE, 11, AMBER_BODY, False, False)], MSO_ANCHOR.MIDDLE)


def add_deps(prs, blank):
    s = prs.slides.add_slide(blank)
    header(s, "Lock These Before PI Planning", "Each turns an objective from aspirational to genuinely Achievable")
    rows = len(DEPS) + 1
    t = s.shapes.add_table(rows, 3, Inches(0.55), Inches(1.45), Inches(12.2), Inches(3.9)).table
    for j, wdt in enumerate([Inches(6.0), Inches(3.4), Inches(2.8)]):
        t.columns[j].width = wdt
    for j, h in enumerate(["Lock before Aug 5 / PI Planning", "Threatens", "Owner"]):
        cell(t, 0, j, h, 11.5, WHITE, True, PP_ALIGN.LEFT, DARK)
    for i, (item, threat, owner) in enumerate(DEPS, start=1):
        fill = LIGHT if i % 2 else WHITE
        cell(t, i, 0, item, 10.5, DARK, False, PP_ALIGN.LEFT, fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 1, threat, 10, RED, True, PP_ALIGN.LEFT, fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 2, owner, 10, TEXT, False, PP_ALIGN.LEFT, fill, MSO_ANCHOR.MIDDLE)
    rect(s, Inches(0.55), Inches(5.6), Inches(12.2), Inches(0.95), ACCENT_TINT, line=ACCENT)
    textbox(s, Inches(0.72), Inches(5.66), Inches(11.9), Inches(0.85),
            [("These are already in the PI-3 open-actions list - reframed here as entry criteria: a commitment isn't real until its dependency is closed.",
              10.5, DARK, False, False)], MSO_ANCHOR.MIDDLE)


def build_deck():
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    blank = prs.slide_layouts[6]
    add_title(prs, blank)
    add_overview(prs, blank)
    for o in OBJECTIVES:
        add_objective(prs, blank, o)
    add_timeline(prs, blank)
    add_deps(prs, blank)
    return prs


def save_safe(prs, path):
    try:
        prs.save(path); print("Saved", os.path.basename(path), "-", len(prs.slides._sldIdLst), "slides")
    except PermissionError:
        base, ext = os.path.splitext(path)
        alt = base + "-new" + ext
        prs.save(alt); print("Target locked - saved to", os.path.basename(alt))


here = os.path.dirname(os.path.abspath(__file__))
save_safe(build_deck(), os.path.join(here, "PI3-Objectives-" + EDITION + ".pptx"))
