"""
CMDB-CSDM Weekly Leadership Status Deck generator.

Single source of truth for the weekly status report. Edit the CONTENT
section below, then run:  python build_status_deck.py
Outputs (named by EDITION date):
  - <EDITION>.pptx   PowerPoint for customer leadership
  - <EDITION>.md     Marp markdown mirror for Obsidian

Weekly update workflow:
  1. Copy last week's values, bump EDITION / WEEK_LABEL / ITERATION.
  2. Move resolved items out of RISKS; refresh THIS_WEEK with the new deltas
     (this is what powers the "week over week" view).
  3. Re-run the script.

2026-07-01 edition: scorecard reframed from the 6 PI objectives to the 3 CO5
contractual deliverables (per the 6/29 CO5 Deliverable Alignment). Sprint 2.3
velocity + delivery metrics are now derived from the 6/30 ADO board, replacing
the prior "pending" placeholders (#3 iteration velocity, #5 delivery).
"""
import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------------------------------------------------------------------------
# CONTENT  (edit this block each week)
# ---------------------------------------------------------------------------
EDITION      = "2026-07-08"
WEEK_LABEL   = "Week of July 8, 2026"
PI           = "PI-2"
PI_WINDOW    = "May 13 - Aug 4, 2026"
ITERATION    = "Iteration 2.5  (Jul 8 - Jul 21)"
AUTHOR       = "Manuel Vazquez  -  Scrum Master, CMDB-CSDM"
EDITION_NOTE = ("Sprint 2.4 close-out (Jul 7). Board state: Jul 8 ADO snapshot. "
                "Sprint 2.5 scope TBD — work items to be added once planning is confirmed. "
                "CO6 pending execution as of Jul 8.")

OVERALL_STATUS = ("CO5 FOCUS - CO6 PENDING; DEV FREEZE ACTIVE",
                  "Sprint 2.4 closed Jul 7 with 9 items / 10 pts accepted. Delivery momentum is real - "
                  "SOX BA review done, Airlift exports complete, CI retirement executed - but the CO5 "
                  "acceptance gate (Data Dictionary CCB approval Jul 21) is 13 days out and CO6 is still unsigned. "
                  "Dev code freeze (Jul 4-18) is now active, constraining Sprint 2.5 delivery capacity. "
                  "The $533,775 holdback stands until CO6 is executed or CO5 is formally accepted.")

# Operational label (D1/D2/D3), CO5 deliverable, ADO feature, "Accept by", RAG, one-line status,
# then sub-item counts: scope (sub-items / acceptance stories), moving, stuck.
OBJECTIVES = [
    ("D1", "Governance - Data Dictionary, Data Cert, ESS-02, SOX, CCB", "1480087", "CCB 7/21", "AMBER",
     "SOX BA review (1480105) Closed - D1 acceptance evidence confirmed; 4 Data Dictionaries in Validation "
     "(CCB review Jul 7-20, approval Jul 21); ESS-02 (1480102) and Monthly CCB (1480107) Ready DoR; "
     "Data Cert pilot (1480099) Ready DoR; code freeze gates ESS-02/CCB until Jul 18", 8, 8, 0),
    ("D2", "Automated Data Ingestion - Computers / Servers / Databases", "-", "7/21", "AMBER",
     "Airlift exports (PA/KY VMware + PA Physical) Closed; imports in Validation; 3 audit-dashboard spikes "
     "(Servers/DB/Computer) carried from 2.4 Active - did not close; no formal 90%-measurement acceptance story", 3, 3, 0),
    ("D3", "Other Enhancement - Evaluate Qualys Integration", "1428703", "CO6 Sep 30", "AMBER",
     "Part 1 (1428703) still Active - did not close in 2.4; dev freeze further delays; "
     "Part 2 (1428704) status unconfirmed (G4); CO5 bar = evaluate + document; "
     "CO6 (Sep 30) delivers full one-way PROD integration (ServiceNow → Qualys)", 2, 1, 1),
]
# Portfolio totals (sub-item counts across the 3 CO5 deliverables)
OBJ_SCOPE  = sum(o[6] for o in OBJECTIVES)
OBJ_MOVING = sum(o[7] for o in OBJECTIVES)
OBJ_STUCK  = sum(o[8] for o in OBJECTIVES)
STORY_CAVEAT = ("Sprint 2.4 final board snapshot (Jul 8). Accepted/Done: 9 items / 10 pts. "
                "12 items / 18 pts in validation carrying into Sprint 2.5. "
                "CO6 unsigned as of Jul 8 — $533,775 holdback risk persists.")

# Sprint 2.4 velocity - DERIVED from the Jul 8 ADO board (sprint close-out snapshot).
# Done = Closed + Resolved. Board grew mid-sprint with retirement cluster + Airlift imports.
VELOCITY_2_4 = [
    ("On board",                 "~38 items / ~48 pts", "ACCENT"),
    ("Accepted / Done",          "9 items - 10 pts",    "GREEN"),
    ("Delivered, in validation", "12 items - 18 pts",   "AMBER"),
    ("Active / not started",     "~17 items",           "MUTED"),
]
VELOCITY_NOTE = ("Derived from the Jul 8 ADO board (Sprint 2.4 close-out). Board grew during sprint "
                 "with retirement cluster and Airlift imports added mid-sprint. "
                 "18 pts in validation carry into Sprint 2.5 pending acceptance.")

# ---------------------------------------------------------------------------
# WEEK-OVER-WEEK HISTORY  (#1: persist a snapshot each run, diff vs prior)
# ---------------------------------------------------------------------------
HISTORY_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "status_history.json")
RAG_RANK = {"RED": 1, "AMBER": 2, "GREEN": 3}
RAGWORD  = {"GREEN": "On track", "AMBER": "Watch", "RED": "Blocked"}


def load_history():
    if os.path.exists(HISTORY_PATH):
        with open(HISTORY_PATH, encoding="utf-8") as f:
            return json.load(f)
    return {}


def current_snapshot():
    return {
        "week_label": WEEK_LABEL,
        "iteration": ITERATION,
        "portfolio": {"scope": OBJ_SCOPE, "moving": OBJ_MOVING, "stuck": OBJ_STUCK},
        "objectives": {o[1]: {"pri": o[0], "rag": o[4], "scope": o[6],
                              "moving": o[7], "stuck": o[8]} for o in OBJECTIVES},
    }


HISTORY = load_history()
CURR = current_snapshot()
_prior_keys = sorted(k for k in HISTORY if k < EDITION)
PRIOR_EDITION = _prior_keys[-1] if _prior_keys else None
PRIOR = HISTORY.get(PRIOR_EDITION)
# This edition reframes the scorecard (6 PI objectives -> 3 CO5 deliverables), so the
# prior snapshot's objective names will not match; week-over-week resets to a CO5 baseline.
STRUCTURE_REFRAMED = False


def wow_rows():
    """Per-objective status transitions vs the prior snapshot.
    Empty list on a baseline edition (no prior snapshot exists yet)."""
    if not PRIOR:
        return []
    prev = PRIOR.get("objectives", {})
    out = []
    for o in OBJECTIVES:
        name, rag, scope, moving, stuck = o[1], o[4], o[6], o[7], o[8]
        p = prev.get(name)
        if not p:
            out.append({"name": name, "last": None, "this": rag,
                        "d_moving": None, "d_stuck": None, "trend": "NEW"})
            continue
        if RAG_RANK[rag] > RAG_RANK[p["rag"]]:
            trend = "UP"
        elif RAG_RANK[rag] < RAG_RANK[p["rag"]]:
            trend = "DOWN"
        else:
            trend = "SAME"
        out.append({"name": name, "last": p["rag"], "this": rag,
                    "d_moving": moving - p["moving"], "d_stuck": stuck - p["stuck"],
                    "trend": trend})
    return out


def wow_transitions():
    """(newly_blocked, newly_unblocked) objective-name lists vs prior snapshot."""
    blocked, unblocked = [], []
    if PRIOR:
        prev = PRIOR.get("objectives", {})
        for o in OBJECTIVES:
            name, rag = o[1], o[4]
            p = prev.get(name)
            if not p:
                continue
            if p["rag"] != "RED" and rag == "RED":
                blocked.append(name)
            elif p["rag"] == "RED" and rag != "RED":
                unblocked.append(name)
    return blocked, unblocked

# PI over PI
PI_PROGRESS = [
    ("PI-2", "May 13 - Aug 4, 2026", "ACTIVE", "Iteration 2.5 of 6", "Sprint 2.4 closed Jul 7 (9 items / 10 pts); CO5 acceptance gate = Data Dictionary CCB approval Jul 21"),
    ("PI-3", "Aug 5 - Oct 27, 2026", "PLANNING", "Planning Jul 22 - Aug 4", "CO6 gap closure (if signed), NowAssist rollout, Qualys full integration, CI Cert extended scope"),
]
PI_NOTE = "Only PI-2 has live data today. PI-over-PI comparison populates once PI-3 begins (planning during the PI-2 IP iteration, Jul 22 - Aug 4)."

# Iteration over Iteration: id, window, status, theme, outcomes[]
ITERATIONS = [
    ("2.3", "Jun 10 - 23",    "COMPLETED", "Deploy & scale",
        ["Data Certification dashboard -> PROD (6/23)",
         "NowAssist accepted: prerequisites, plugin activation, CI summarization",
         "11 pts accepted; 23 pts delivered into validation"]),
    ("2.4", "Jun 24 - Jul 7", "COMPLETED", "CO5 acceptance push",
        ["9 items / 10 pts accepted: SOX BA review, Airlift exports, CI retirement, NowAssist Duplicate CIs",
         "Data Dictionaries (Servers/Computers/Biz Apps/DB) delivered to Validation",
         "12 items / 18 pts in validation - largest carry-over pipeline of PI-2",
         "3 audit dashboard spikes carried to 2.5 Active (Servers/DB/Computer)"]),
    ("2.5", "Jul 8 - Jul 21",  "ACTIVE",    "CO5 acceptance gate + code freeze navigation",
        ["Scope TBD - 2.5 work items to be confirmed; planning in progress",
         "Dev code freeze Jul 4-18 constrains first 10 days of sprint",
         "Data Dictionary CCB approval target Jul 21 (sprint end)",
         "3 audit spikes (Servers/DB/Computer) carried from 2.4"]),
]

# Week over week
THIS_WEEK = [
    "Sprint 2.4 closed (Jul 7) - 9 items / 10 pts accepted: SOX BA review, Airlift exports (PA/KY VMware + PA Physical), bulk CI retirement, NowAssist Duplicate CIs skill, IRE behaviour review spike.",
    "SOX Business Apps (1480105) Closed - CO5 D1 acceptance evidence confirmed; Airlift inventory export cluster complete.",
    "Data Dictionaries (Servers/Computers/Biz Apps/DB) delivered to Validation - CCB review Jul 7-20, approval target Jul 21.",
    "12 items / 18 pts in validation entering Sprint 2.5 - largest carry-over pipeline of PI-2.",
]
WEEK_NOTE = "Week-over-week deltas resume vs. the 7/1 CO5 baseline (same 3-deliverable frame)."

# Deep dives: header, [(item, status, detail)]
DEEP_DIVE_1 = ("CO5 Deliverable 1 - Governance", [
    ("Data Dictionaries (Servers/Computers/Biz Apps/DB)", "AMBER", "All 4 dictionaries in Validation (1480088/090/097/098) - CCB customer review Jul 7-20, approval target Jul 21 (monthly CCB). Databases closes G1 at artifact level."),
    ("SOX + Data Certification", "GREEN", "SOX BA review (1480105) Closed - CO5 D1 acceptance evidence confirmed. Data Cert pilot (1480099) Ready DoR; dashboard in PROD since 6/23."),
    ("ESS-02 / Monthly CCB", "AMBER", "ESS-02 alignment (1480102) and Monthly CCB (1480107) remain Ready DoR - not started; dev-dependent work gated by code freeze (lifts Jul 18). CCB story closes G5."),
])
DEEP_DIVE_2 = ("CO5 Deliverables 2 & 3 - Data Ingestion + Qualys", [
    ("Data Ingestion - Airlift inventory + SCCM/Discovery precedence", "AMBER", "Airlift exports (PA/KY VMware + PA Physical) Closed; imports (1487867/868/869) in Validation. Retirement: execute Closed (1487878), validate Active (1487872). SCCM/Discovery precedence reported complete."),
    ("90% Coverage (Computers G2 / Servers G3 / DB)", "AMBER", "3 audit-dashboard spikes (1480112 Servers, 1480113 DB, 1480114 Computer) carried from 2.4 Active - did not close in sprint. BA spike (1480111) in Validation. No formal 90%-measurement acceptance story yet."),
    ("Qualys Integration (Other Enhancement)", "AMBER", "Part 1 (1428703) still Active - did not close in 2.4; dev freeze (Jul 4-18) further delays. Part 2 (1428704) status unconfirmed (G4). CO5 bar = evaluate + document. CO6 (pending signature) delivers full one-way PROD integration (ServiceNow → Qualys) by Sep 30."),
])

# Milestones
MILESTONES = [
    ("Jul 4-18",      "Dev code freeze — ACTIVE NOW; constrains Sprint 2.5 delivery through Jul 18"),
    ("Jul 7",         "Sprint 2.4 end ✓ — 9 items / 10 pts accepted; 12 items / 18 pts to validation"),
    ("Jul 8",         "Sprint 2.5 start — CO5 acceptance gate; Data Dictionary CCB push"),
    ("Jul 18",        "Dev code freeze lifts — ESS-02 / CCB / Qualys work can resume"),
    ("Jul 18-Aug 15", "Test code freeze — overlaps Sprint 2.5 tail and IP iteration"),
    ("Jul 21",        "Data Dictionary CCB approval (monthly CCB) — critical CO5 D1 gate; Sprint 2.5 end"),
    ("Jul 22-Aug 4",  "IP Iteration — PI-3 planning begins"),
    ("Aug 5",         "PI-3 begins"),
]

# Risks & asks: risk, severity, detail, ask
RISKS = [
    ("CO6 pending execution", "HIGH", "CO6 unsigned as of Jul 8; signature expected this week. Effective Jun 30 retroactively when signed. $2.7M, 9 workstreams through Oct 30: Network Discovery, Service Mapping, Qualys PROD (Sep 30), CMDB Governance (Data Cert expanded), Legacy Rationalization, ITSM PO, ATF Strategy, Platform Support, NERC-CIP Strategy. $533,775 June holdback stands until signed.", "Execute CO6 this week. Confirm with Christian / Aaron Simeon — effective date is already Jun 30; delay only extends holdback exposure."),
    ("Dev code freeze Jul 4-18", "HIGH", "Freeze is currently active - spans first 10 days of Sprint 2.5. Dev-dependent CO5 items (ESS-02, CCB story, Qualys Part 1) cannot land until Jul 18.", "Awareness; prioritize validation and documentation work in 2.5 through Jul 18."),
    ("3 audit dashboard spikes incomplete", "MEDIUM", "Servers (1480112), Database (1480113), Computer (1480114) audit spikes carried into 2.5 Active - did not close in 2.4. G2/G3 coverage measurement evidence incomplete.", "Prioritize spike completion in 2.5; required before 90%-coverage acceptance stories can be created."),
    ("90% coverage acceptance gap", "MEDIUM", "No formal 90%-measurement acceptance story for Computers (G2) or Servers (G3). Audit spikes must complete first; stories must then be created and accepted.", "Endorse creating coverage-measurement acceptance stories once spikes close."),
    ("Deprioritized streams stranded", "LOW", "Service Mapping, Network Discovery, and NowAssist phases 2-5 remain paused. Test freeze (Jul 18-Aug 15) further limits resumption window.", "Confirm resumption plan after CO5 acceptance; sequence around test freeze."),
]

# What's next
NEXT = [
    "Execute CO6 this week (effective Jun 30 retroactively) — $2.7M, 9 workstreams through Oct 30; releases $533,775 holdback.",
    "Secure Data Dictionary CCB approval Jul 21 — critical CO5 D1 gate; validation window Jul 7-20.",
    "Complete 3 audit dashboard spikes (Servers/DB/Computer) in 2.5 — prerequisite for G2/G3 acceptance stories.",
    "Progress ESS-02 (1480102) and Monthly CCB (1480107) off Ready DoR — code freeze lifts Jul 18.",
    "Confirm Qualys Part 1 (1428703) path to completion; CO6 delivers full PROD integration by Sep 30.",
    "Begin CO6 workstream onboarding: Network Discovery, Service Mapping, ITSM PO, Platform Support, Legacy Rationalization.",
]

# ---------------------------------------------------------------------------
# THEME
# ---------------------------------------------------------------------------
DARK   = RGBColor(0x1F, 0x3A, 0x5F)   # navy
ACCENT = RGBColor(0x2E, 0x6D, 0xB4)   # blue
LIGHT  = RGBColor(0xF4, 0xF6, 0xF9)   # near-white panel
TEXT   = RGBColor(0x2D, 0x2D, 0x2D)
MUTED  = RGBColor(0x6B, 0x72, 0x80)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
AMBER  = RGBColor(0xE0, 0x8A, 0x00)
RED    = RGBColor(0xC0, 0x39, 0x2B)
RAG = {"GREEN": GREEN, "AMBER": AMBER, "RED": RED,
       "HIGH": RED, "MEDIUM": AMBER, "LOW": GREEN}
# named colors for metric strips
COLORMAP = {"ACCENT": ACCENT, "GREEN": GREEN, "AMBER": AMBER, "MUTED": MUTED, "DARK": DARK}
# week-over-week trend styling: (cell color, label)
TREND = {"UP": (GREEN, "▲ Improved"), "DOWN": (RED, "▼ Slipped"),
         "SAME": (MUTED, "→ No change"), "NEW": (ACCENT, "＋ New")}

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
FONT = "Segoe UI"

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def _set(run, size, color, bold=False, italic=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = FONT


def rect(slide, x, y, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def textbox(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    """lines: list of (text, size, color, bold, italic) or list of such for paragraphs."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, spec in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        text, size, color, bold, italic = spec
        r = p.add_run(); r.text = text
        _set(r, size, color, bold, italic)
    return tb


def header(slide, title, subtitle=None):
    rect(slide, 0, 0, EMU_W, Inches(1.05), DARK)
    rect(slide, 0, Inches(1.05), EMU_W, Inches(0.06), ACCENT)
    textbox(slide, Inches(0.55), Inches(0.18), Inches(11.5), Inches(0.55),
            [(title, 26, WHITE, True, False)], MSO_ANCHOR.MIDDLE)
    if subtitle:
        textbox(slide, Inches(0.57), Inches(0.66), Inches(11.5), Inches(0.32),
                [(subtitle, 12, RGBColor(0xC9, 0xD6, 0xE6), False, False)], MSO_ANCHOR.MIDDLE)
    # footer
    textbox(slide, Inches(0.55), Inches(7.05), Inches(8), Inches(0.35),
            [(f"PPL CMDB-CSDM  |  {PI}  |  {WEEK_LABEL}", 9, MUTED, False, False)])


def chip(slide, x, y, label, color, w=Inches(1.0)):
    c = rect(slide, x, y, w, Inches(0.30), color)
    tf = c.text_frame; tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    _set(r, 10, WHITE, True)
    return c


def pending_box(slide, x, y, w, h, title, body):
    """Annotated placeholder for a metric that is intentionally not yet populated."""
    rect(slide, x, y, w, h, LIGHT)
    rect(slide, x, y, Inches(0.12), h, AMBER)
    chip(slide, x + Inches(0.28), y + Inches(0.16), "PENDING DATA", AMBER, Inches(1.7))
    textbox(slide, x + Inches(2.15), y + Inches(0.13), w - Inches(2.4), Inches(0.45),
            [(title, 12.5, DARK, True, False)], MSO_ANCHOR.MIDDLE)
    textbox(slide, x + Inches(0.28), y + Inches(0.62), w - Inches(0.55), h - Inches(0.72),
            [(body, 10.5, MUTED, False, True)])


def metric_strip(slide, x, y, cards, cwid=Inches(2.95), gap=Inches(0.12)):
    """Row of metric cards: cards = [(label, value, COLORKEY), ...]."""
    cx = x
    for label, value, ckey in cards:
        col = COLORMAP.get(ckey, ACCENT)
        rect(slide, cx, y, cwid, Inches(1.0), LIGHT)
        rect(slide, cx, y, cwid, Inches(0.08), col)
        textbox(slide, cx + Inches(0.15), y + Inches(0.16), cwid - Inches(0.3), Inches(0.35),
                [(label, 11, MUTED, True, False)])
        textbox(slide, cx + Inches(0.15), y + Inches(0.5), cwid - Inches(0.3), Inches(0.45),
                [(value, 15, DARK, True, False)])
        cx += cwid + gap


# ---- Slide 1: Title -------------------------------------------------------
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, EMU_W, EMU_H, DARK)
rect(s, 0, Inches(4.55), EMU_W, Inches(0.08), ACCENT)
textbox(s, Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.0),
        [("CMDB-CSDM Delivery Status", 44, WHITE, True, False)])
textbox(s, Inches(0.82), Inches(3.05), Inches(11.7), Inches(0.6),
        [("Weekly Leadership Report", 22, RGBColor(0xC9, 0xD6, 0xE6), False, False)])
textbox(s, Inches(0.82), Inches(4.75), Inches(11.7), Inches(1.4), [
    (f"{PI}  ·  {ITERATION}", 16, WHITE, True, False),
    (WEEK_LABEL, 14, RGBColor(0xC9, 0xD6, 0xE6), False, False),
    (AUTHOR, 12, RGBColor(0x9F, 0xB2, 0xCC), False, False),
])

# ---- Slide 2: How to read -------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "How to Read This Report", EDITION_NOTE)
textbox(s, Inches(0.55), Inches(1.4), Inches(12), Inches(0.4),
        [("Three progress lenses, refreshed every week:", 15, TEXT, True, False)])
lenses = [("PI over PI", "Are we delivering more value each Program Increment?"),
          ("Iteration over Iteration", "Is each 2-week sprint building on the last?"),
          ("Week over Week", "What concretely moved since the last report?")]
y = Inches(2.0)
for name, desc in lenses:
    rect(s, Inches(0.55), y, Inches(0.12), Inches(0.55), ACCENT)
    textbox(s, Inches(0.8), y, Inches(11.5), Inches(0.6),
            [(name, 14, ACCENT, True, False), (desc, 12, TEXT, False, False)])
    y += Inches(0.75)
# RAG legend
textbox(s, Inches(0.55), Inches(4.6), Inches(6), Inches(0.4),
        [("Status legend:", 13, TEXT, True, False)])
chip(s, Inches(0.55), Inches(5.05), "GREEN", GREEN, Inches(1.1))
textbox(s, Inches(1.8), Inches(5.05), Inches(5), Inches(0.3), [("On track", 12, TEXT, False, False)])
chip(s, Inches(0.55), Inches(5.5), "AMBER", AMBER, Inches(1.1))
textbox(s, Inches(1.8), Inches(5.5), Inches(6), Inches(0.3), [("Progressing, with watch items", 12, TEXT, False, False)])
chip(s, Inches(0.55), Inches(5.95), "RED", RED, Inches(1.1))
textbox(s, Inches(1.8), Inches(5.95), Inches(6), Inches(0.3), [("Blocked / needs intervention", 12, TEXT, False, False)])

# ---- Slide 3: Executive Summary ------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "Executive Summary", ITERATION)
chip(s, Inches(0.55), Inches(1.35), OVERALL_STATUS[0], AMBER, Inches(5.2))
textbox(s, Inches(0.55), Inches(1.85), Inches(12.2), Inches(0.9),
        [(OVERALL_STATUS[1], 12.5, TEXT, False, False)])
textbox(s, Inches(0.55), Inches(2.9), Inches(6), Inches(0.35),
        [("Wins this period", 15, GREEN, True, False)])
yy = Inches(3.3)
for w in THIS_WEEK:
    textbox(s, Inches(0.7), yy, Inches(6.0), Inches(0.7), [("+  " + w, 11, TEXT, False, False)])
    yy += Inches(0.82)
textbox(s, Inches(6.9), Inches(2.9), Inches(6), Inches(0.35),
        [("Needs leadership attention", 15, RED, True, False)])
yy = Inches(3.3)
for r in RISKS[:3]:
    textbox(s, Inches(7.05), yy, Inches(5.8), Inches(0.7),
            [("!  " + r[0] + " - " + r[3], 11, TEXT, False, False)])
    yy += Inches(0.82)

# ---- Slide 4: CO5 Deliverables Scorecard ---------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "CO5 Deliverables Scorecard", "The 3 contractual deliverables  ·  sub-items = moving / in-scope")
# Executive tiles: one card per CO5 deliverable.
tw, th = Inches(4.0), Inches(2.55)
gx, gy = Inches(0.55), Inches(1.30)
for i, (pri, name, obj, acc, rag, note, scope, moving, stuck) in enumerate(OBJECTIVES):
    cx = gx + (i % 3) * (tw + Inches(0.3))
    cy = gy + (i // 3) * (th + Inches(0.25))
    rect(s, cx, cy, tw, th, LIGHT)
    rect(s, cx, cy, tw, Inches(0.72), RAG[rag])
    textbox(s, cx + Inches(0.2), cy + Inches(0.04), tw - Inches(0.4), Inches(0.64),
            [(f"{pri}   {name}", 10.5, WHITE, True, False)], MSO_ANCHOR.MIDDLE)
    textbox(s, cx + Inches(0.2), cy + Inches(0.82), tw - Inches(0.4), Inches(0.8),
            [(str(scope), 44, DARK, True, False)])
    textbox(s, cx + Inches(1.4), cy + Inches(1.12), tw - Inches(1.5), Inches(0.6),
            [("sub-items in scope", 11, MUTED, False, False),
             (f"{moving} moving  ·  {stuck} stuck", 11, TEXT, True, False)])
    textbox(s, cx + Inches(0.2), cy + Inches(1.85), tw - Inches(0.4), Inches(0.62),
            [(note, 8.5, TEXT, False, False)])
textbox(s, Inches(0.55), Inches(6.72), Inches(12.25), Inches(0.35),
        [(STORY_CAVEAT, 8.5, MUTED, False, True)])

# ---- Slide 5: PI over PI --------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "Progress: PI over PI", "Are we delivering more value each Program Increment?")
x = Inches(0.55)
for pi, window, status, sub, desc in PI_PROGRESS:
    col = GREEN if status == "ACTIVE" else MUTED
    card = rect(s, x, Inches(1.5), Inches(6.0), Inches(3.4), LIGHT)
    rect(s, x, Inches(1.5), Inches(6.0), Inches(0.7), DARK)
    textbox(s, x + Inches(0.25), Inches(1.6), Inches(5.5), Inches(0.5),
            [(pi + "   " + status, 18, WHITE, True, False)], MSO_ANCHOR.MIDDLE)
    textbox(s, x + Inches(0.25), Inches(2.4), Inches(5.5), Inches(2.3), [
        (window, 13, ACCENT, True, False),
        (sub, 12, MUTED, False, True),
        ("", 6, TEXT, False, False),
        (desc, 13, TEXT, False, False),
    ])
    x += Inches(6.3)
# #4 placeholder: PI-over-PI comparison metric (pending until PI-3 has live data)
pending_box(s, Inches(0.55), Inches(5.05), Inches(12.2), Inches(0.95),
            "PI-over-PI comparison metric",
            "Comparison basis (objectives completed / business value delivered / stories accepted per PI) "
            "to be fixed before PI-3 begins. Only PI-2 has live data today - not yet ready for relevant data.")
textbox(s, Inches(0.55), Inches(6.15), Inches(12.2), Inches(0.7),
        [("Note: " + PI_NOTE, 11, MUTED, False, True)])

# ---- Slide 6: Iteration over Iteration (velocity now DERIVED) -------------
s = prs.slides.add_slide(BLANK)
header(s, "Progress: Iteration over Iteration", "Is each 2-week sprint building on the last?")
x = Inches(0.55)
cw = Inches(4.0)
for it, window, status, theme, outcomes in ITERATIONS:
    col = GREEN if status == "ACTIVE" else ACCENT
    rect(s, x, Inches(1.45), cw, Inches(0.78), col)
    textbox(s, x + Inches(0.2), Inches(1.5), cw - Inches(0.3), Inches(0.7),
            [("Sprint " + it, 16, WHITE, True, False),
             (window + "  ·  " + status, 10, RGBColor(0xE8, 0xEE, 0xF6), False, False)], MSO_ANCHOR.MIDDLE)
    rect(s, x, Inches(2.23), cw, Inches(3.0), LIGHT)
    textbox(s, x + Inches(0.2), Inches(2.33), cw - Inches(0.3), Inches(0.4),
            [(theme, 12, ACCENT, True, True)])
    yy = Inches(2.85)
    for o in outcomes:
        textbox(s, x + Inches(0.2), yy, cw - Inches(0.35), Inches(0.8),
                [("- " + o, 10, TEXT, False, False)])
        yy += Inches(0.58)
    x += cw + Inches(0.3)
# Sprint 2.4 velocity (close-out snapshot from Jul 8 board)
textbox(s, Inches(0.55), Inches(5.4), Inches(12.2), Inches(0.32),
        [("Sprint 2.4 delivery - derived from the Jul 8 ADO board (sprint close-out)", 12.5, DARK, True, False)])
metric_strip(s, Inches(0.55), Inches(5.75), VELOCITY_2_4)
textbox(s, Inches(0.55), Inches(6.85), Inches(12.2), Inches(0.5),
        [(VELOCITY_NOTE, 9, MUTED, False, True)])

# ---- Slide 7: Week over Week (structure reframed; delivery now populated) -
s = prs.slides.add_slide(BLANK)
header(s, "Progress: Week over Week", "What changed in status since the last report?")
textbox(s, Inches(0.55), Inches(1.35), Inches(12.2), Inches(0.6),
        [("Sprint 2.4 closed Jul 7. Reporting continues on the CO5 contractual deliverable frame. "
          "Week-over-week deltas reflect movement since the Jun 29 report.",
          11.5, MUTED, False, True)])
# Delivered this period - Sprint 2.4 close-out
textbox(s, Inches(0.55), Inches(2.05), Inches(12.2), Inches(0.32),
        [("Delivered this period (Sprint 2.4 close-out)", 14, DARK, True, False)])
metric_strip(s, Inches(0.55), Inches(2.45), VELOCITY_2_4)
# Key moves this week
textbox(s, Inches(0.55), Inches(3.7), Inches(12.2), Inches(0.32),
        [("Key moves", 14, DARK, True, False)])
moves = [
    "Sprint 2.4 closed (Jul 7) - 9 items / 10 pts accepted: SOX BA review, Airlift exports (PA/KY VMware + PA Physical), bulk CI retirement, NowAssist Duplicate CIs, IRE review spike.",
    "SOX Business Apps (1480105) Closed - CO5 D1 acceptance evidence confirmed.",
    "Data Dictionaries (Servers/Computers/Biz Apps/DB) delivered to Validation - CCB review Jul 7-20, approval target Jul 21.",
    "12 items / 18 pts in validation entering Sprint 2.5 - largest carry-over pipeline of PI-2.",
]
yy = Inches(4.1)
for m in moves:
    textbox(s, Inches(0.7), yy, Inches(12.0), Inches(0.5), [("+  " + m, 11.5, TEXT, False, False)])
    yy += Inches(0.5)
textbox(s, Inches(0.55), yy + Inches(0.1), Inches(12.2), Inches(0.5),
        [("Newly closed:  SOX BA review · Airlift exports · CI retirement execution   ·   Newly carried to 2.5:  3 audit spikes (Servers/DB/Computer)",
          11, AMBER, True, False)])

# ---- Slides 8-9: Deep dives ----------------------------------------------
for title, items in (DEEP_DIVE_1, DEEP_DIVE_2):
    s = prs.slides.add_slide(BLANK)
    header(s, title, "Status and detail by deliverable")
    yy = Inches(1.6)
    for name, rag, detail in items:
        rect(s, Inches(0.55), yy, Inches(12.2), Inches(1.45), LIGHT)
        chip(s, Inches(0.75), yy + Inches(0.2), rag, RAG[rag], Inches(1.1))
        textbox(s, Inches(2.05), yy + Inches(0.15), Inches(10.4), Inches(0.4),
                [(name, 15, DARK, True, False)])
        textbox(s, Inches(2.05), yy + Inches(0.6), Inches(10.4), Inches(0.8),
                [(detail, 12, TEXT, False, False)])
        yy += Inches(1.65)

# ---- Slide 10: Milestones -------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "Key Milestones & Upcoming Dates", "CO5 acceptance critical path")
yy = Inches(1.35)
for date, ev in MILESTONES:
    rect(s, Inches(0.55), yy, Inches(1.9), Inches(0.55), ACCENT)
    textbox(s, Inches(0.6), yy, Inches(1.8), Inches(0.55),
            [(date, 13, WHITE, True, False)], MSO_ANCHOR.MIDDLE, PP_ALIGN.CENTER)
    rect(s, Inches(2.55), yy, Inches(10.2), Inches(0.55), LIGHT)
    textbox(s, Inches(2.75), yy, Inches(9.9), Inches(0.55),
            [(ev, 12.5, TEXT, False, False)], MSO_ANCHOR.MIDDLE)
    yy += Inches(0.68)

# ---- Slide 11: Risks & Asks ----------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "Risks & Leadership Asks", "Where we need help to stay on track")
rows = len(RISKS) + 1
tbl_shape = s.shapes.add_table(rows, 4, Inches(0.55), Inches(1.4), Inches(12.2), Inches(4.6))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(2.6)
tbl.columns[1].width = Inches(1.1)
tbl.columns[2].width = Inches(4.4)
tbl.columns[3].width = Inches(4.1)
for j, h in enumerate(["Risk", "Severity", "Detail", "Ask of Leadership"]):
    cell = tbl.cell(0, j)
    cell.fill.solid(); cell.fill.fore_color.rgb = DARK
    p = cell.text_frame.paragraphs[0]; r = p.add_run(); r.text = h
    _set(r, 12, WHITE, True)
for i, (risk, sev, detail, ask) in enumerate(RISKS, start=1):
    vals = [risk, sev, detail, ask]
    for j, v in enumerate(vals):
        cell = tbl.cell(i, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT if i % 2 else WHITE
        p = cell.text_frame.paragraphs[0]; r = p.add_run(); r.text = v
        if j == 1:
            cell.fill.fore_color.rgb = RAG[sev]
            _set(r, 11, WHITE, True); p.alignment = PP_ALIGN.CENTER
        elif j == 0:
            _set(r, 11, DARK, True)
        else:
            _set(r, 10.5, TEXT, False)

# ---- Slide 12: What's Next ------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "What's Next", "CO5 acceptance, CO6 gap closure, and resuming paused streams")
yy = Inches(1.7)
for n in NEXT:
    rect(s, Inches(0.55), yy + Inches(0.05), Inches(0.12), Inches(0.5), ACCENT)
    textbox(s, Inches(0.85), yy, Inches(11.8), Inches(0.7), [(n, 14, TEXT, False, False)])
    yy += Inches(0.85)

# ---- Slide 13: Appendix ---------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "Appendix", "Story-level detail")
textbox(s, Inches(0.55), Inches(1.8), Inches(12), Inches(3), [
    ("Full traceability is maintained in the team's working notes:", 14, TEXT, True, False),
    ("", 8, TEXT, False, False),
    ("- CO5 deliverable traceability:  co5-deliverable-tracking.md", 12, ACCENT, False, False),
    ("- Objectives -> Features -> Stories:  PI-2/pi2-objectives-features-stories.md", 12, ACCENT, False, False),
    ("- Activity by Iteration (2.1 - 2.4):  PI-2/pi2-iteration-activity.md", 12, ACCENT, False, False),
    ("- Program Backlog Review notes:  PI-2/Program-Backlog-Review/", 12, ACCENT, False, False),
    ("", 8, TEXT, False, False),
    ("This deck is regenerated weekly from build_status_deck.py.", 11, MUTED, False, True),
])

def save_safe(save_fn, path):
    """Save via save_fn(path); if the target is locked (open in PowerPoint /
    Explorer preview), write to a '-new' file instead so the run never fails."""
    try:
        save_fn(path)
        print("Saved", path)
    except PermissionError:
        base, ext = os.path.splitext(path)
        alt = base + "-new" + ext
        save_fn(alt)
        print("Target locked - saved to", alt, "(close the original, then rename)")

out_pptx = os.path.join(os.path.dirname(os.path.abspath(__file__)), EDITION + ".pptx")
save_safe(prs.save, out_pptx)

# ---------------------------------------------------------------------------
# MARKDOWN (Marp) MIRROR
# ---------------------------------------------------------------------------
def md_lines():
    L = []
    L.append("---")
    L.append("marp: true")
    L.append("theme: default")
    L.append("paginate: true")
    L.append(f"title: CMDB-CSDM Delivery Status - {WEEK_LABEL}")
    L.append("---")
    L.append("")
    # title
    L.append("<!-- _class: lead -->")
    L.append("# CMDB-CSDM Delivery Status")
    L.append("### Weekly Leadership Report")
    L.append("")
    L.append(f"**{PI} · {ITERATION}**")
    L.append(f"{WEEK_LABEL}")
    L.append(f"_{AUTHOR}_")
    L.append("")
    L.append("---")
    # how to read
    L.append("## How to Read This Report")
    L.append(f"_{EDITION_NOTE}_")
    L.append("")
    L.append("Three progress lenses, refreshed every week:")
    L.append("")
    L.append("- **PI over PI** — are we delivering more value each Program Increment?")
    L.append("- **Iteration over Iteration** — is each 2-week sprint building on the last?")
    L.append("- **Week over Week** — what concretely moved since the last report?")
    L.append("")
    L.append("**Legend:** 🟢 On track  ·  🟡 Progressing, watch items  ·  🔴 Blocked / needs intervention")
    L.append("")
    L.append("---")
    # exec summary
    L.append("## Executive Summary")
    L.append(f"### 🟡 {OVERALL_STATUS[0]}")
    L.append(OVERALL_STATUS[1])
    L.append("")
    L.append("**Wins this period**")
    for w in THIS_WEEK:
        L.append(f"- ✅ {w}")
    L.append("")
    L.append("**Needs leadership attention**")
    for r in RISKS[:3]:
        L.append(f"- ⚠️ **{r[0]}** — {r[3]}")
    L.append("")
    L.append("---")
    # scorecard
    L.append("## CO5 Deliverables Scorecard")
    L.append("")
    L.append(f"**Sub-item coverage:** {OBJ_SCOPE} in scope · {OBJ_MOVING} moving · {OBJ_STUCK} need attention")
    L.append("")
    L.append("| # | CO5 Deliverable | Accept by | Sub-items | Status | Summary |")
    L.append("|---|-----------------|-----------|-----------|--------|---------|")
    dot = {"GREEN": "🟢 On track", "AMBER": "🟡 Watch", "RED": "🔴 Blocked"}
    for pri, name, obj, acc, rag, note, scope, moving, stuck in OBJECTIVES:
        L.append(f"| {pri} | {name} | {acc} | {moving}/{scope} | {dot[rag]} | {note} |")
    L.append("")
    L.append(f"> _{STORY_CAVEAT}_")
    L.append("")
    L.append("---")
    # PI over PI
    L.append("## Progress: PI over PI")
    L.append("")
    L.append("| PI | Window | Status | Focus |")
    L.append("|----|--------|--------|-------|")
    for pi, window, status, sub, desc in PI_PROGRESS:
        L.append(f"| **{pi}** | {window} | {status} ({sub}) | {desc} |")
    L.append("")
    L.append(f"> _{PI_NOTE}_")
    L.append("")
    L.append("> ⏳ _**Pending (#4):** PI-over-PI comparison metric (objectives completed / BV delivered / "
             "stories accepted per PI) to be fixed before PI-3 begins - not yet ready for relevant data._")
    L.append("")
    L.append("---")
    # iteration over iteration
    L.append("## Progress: Iteration over Iteration")
    for it, window, status, theme, outcomes in ITERATIONS:
        L.append("")
        L.append(f"**Sprint {it}** — {window} · _{status}_ · {theme}")
        for o in outcomes:
            L.append(f"  - {o}")
    L.append("")
    L.append("**Sprint 2.4 delivery — derived from the Jul 8 ADO board (sprint close-out)**")
    L.append("")
    L.append("| On board | Accepted / Done | Delivered, in validation | Active / not started |")
    L.append("|----------|-----------------|--------------------------|-----------------------|")
    L.append(f"| {VELOCITY_2_4[0][1]} | {VELOCITY_2_4[1][1]} | {VELOCITY_2_4[2][1]} | {VELOCITY_2_4[3][1]} |")
    L.append("")
    L.append(f"> _{VELOCITY_NOTE}_")
    L.append("")
    L.append("---")
    # week over week
    L.append("## Progress: Week over Week")
    L.append("")
    L.append("_Sprint 2.4 closed Jul 7. Reporting continues on the CO5 contractual deliverable frame. "
             "Week-over-week deltas reflect movement since the Jun 29 report._")
    L.append("")
    L.append("**Delivered this period (Sprint 2.4 close-out)**")
    L.append("")
    L.append("| On board | Accepted / Done | Delivered, in validation | Active / not started |")
    L.append("|----------|-----------------|--------------------------|-----------------------|")
    L.append(f"| {VELOCITY_2_4[0][1]} | {VELOCITY_2_4[1][1]} | {VELOCITY_2_4[2][1]} | {VELOCITY_2_4[3][1]} |")
    L.append("")
    L.append("**Key moves**")
    L.append("- Sprint 2.4 closed (Jul 7) - 9 items / 10 pts accepted: SOX BA review, Airlift exports (PA/KY VMware + PA Physical), bulk CI retirement, NowAssist Duplicate CIs, IRE review spike.")
    L.append("- SOX Business Apps (1480105) Closed - CO5 D1 acceptance evidence confirmed.")
    L.append("- Data Dictionaries (Servers/Computers/Biz Apps/DB) delivered to Validation - CCB review Jul 7-20, approval target Jul 21.")
    L.append("- 12 items / 18 pts in validation entering Sprint 2.5 - largest carry-over pipeline of PI-2.")
    L.append("")
    L.append("**Newly closed:** SOX BA review · Airlift exports · CI retirement execution  ·  **Newly carried to 2.5:** 3 audit spikes (Servers/DB/Computer)")
    L.append("")
    L.append("---")
    # deep dives
    for title, items in (DEEP_DIVE_1, DEEP_DIVE_2):
        L.append(f"## {title}")
        L.append("")
        for name, rag, detail in items:
            L.append(f"- {dot[rag]} **{name}** — {detail}")
        L.append("")
        L.append("---")
    # milestones
    L.append("## Key Milestones & Upcoming Dates")
    L.append("")
    L.append("| Date | Event |")
    L.append("|------|-------|")
    for date, ev in MILESTONES:
        L.append(f"| **{date}** | {ev} |")
    L.append("")
    L.append("---")
    # risks
    L.append("## Risks & Leadership Asks")
    L.append("")
    L.append("| Risk | Severity | Detail | Ask of Leadership |")
    L.append("|------|----------|--------|-------------------|")
    sev_dot = {"HIGH": "🔴 HIGH", "MEDIUM": "🟡 MED", "LOW": "🟢 LOW"}
    for risk, sev, detail, ask in RISKS:
        L.append(f"| **{risk}** | {sev_dot[sev]} | {detail} | {ask} |")
    L.append("")
    L.append("---")
    # whats next
    L.append("## What's Next")
    L.append("")
    for n in NEXT:
        L.append(f"- {n}")
    L.append("")
    L.append("---")
    # appendix
    L.append("## Appendix — Story-Level Detail")
    L.append("")
    L.append("- CO5 deliverable traceability: `co5-deliverable-tracking.md`")
    L.append("- Objectives → Features → Stories: `PI-2/pi2-objectives-features-stories.md`")
    L.append("- Activity by Iteration (2.1 - 2.4): `PI-2/pi2-iteration-activity.md`")
    L.append("- Program Backlog Review notes: `PI-2/Program-Backlog-Review/`")
    L.append("")
    L.append("_This deck is regenerated weekly from `build_status_deck.py`._")
    L.append("")
    return "\n".join(L)

out_md = os.path.join(os.path.dirname(os.path.abspath(__file__)), EDITION + ".md")
def _write_md(p):
    with open(p, "w", encoding="utf-8") as f:
        f.write(md_lines())
save_safe(_write_md, out_md)

# ---------------------------------------------------------------------------
# PERSIST THIS WEEK'S SNAPSHOT  (#1) - keyed by EDITION, so re-runs are idempotent
# ---------------------------------------------------------------------------
HISTORY[EDITION] = CURR
with open(HISTORY_PATH, "w", encoding="utf-8") as f:
    json.dump(HISTORY, f, indent=2, ensure_ascii=False)
print(f"History updated: {HISTORY_PATH} ({len(HISTORY)} edition(s))")
