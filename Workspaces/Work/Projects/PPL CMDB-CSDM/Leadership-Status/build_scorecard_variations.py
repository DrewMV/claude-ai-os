"""
Scorecard slide VARIATIONS (exploratory) - separate from the main weekly deck.
Three different ways to convey objective status, each adding a derived story-count
statistic. Counts are derived from team notes (pi2-objectives-features-stories.md),
NOT ADO-validated. 'Completed' is omitted by design until the first PROD deploy (Jun 23).

Run:  python build_scorecard_variations.py  ->  2026-06-11-scorecard-variations.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Derived counts: pri, name, in_scope, moving, stuck, one-line status reason
OBJS = [
    ("P0", "Airlift (Azure Migration)",     4, 4, 0, "On track - dev access resolved; 4 stories active"),
    ("P1", "Discovery / SCCM",              6, 3, 3, "At risk - 1 approval gate + 1 on-hold; ~6 stories not yet ID'd"),
    ("P1", "Data Certification",            1, 1, 0, "On track - in UAT, signed off; PROD Jun 23"),
    ("P2", "Service Mapping",               6, 6, 0, "On track - discovery scaling; Wave 20-21 not yet storied"),
    ("P2", "Regulatory / Qualys",           3, 3, 0, "Progressing - Qualys active; CCB/NERC-CIP not yet storied"),
    ("P3", "NowAssist AI",                  4, 0, 4, "Blocked - all 4 stories await owner reassignment"),
]
TOT_SCOPE  = sum(o[2] for o in OBJS)
TOT_MOVING = sum(o[3] for o in OBJS)
TOT_STUCK  = sum(o[4] for o in OBJS)

CAVEAT = ("Derived from team notes - not ADO-validated. 'Completed' tracking begins at first "
          "PROD deployment (Jun 23); today's view shows work in flight.")

def rag_of(stuck, moving):
    if stuck >= moving and stuck > 0:
        return "RED"
    if stuck > 0:
        return "AMBER"
    return "GREEN"

# theme
DARK   = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0x2E, 0x6D, 0xB4)
LIGHT  = RGBColor(0xF4, 0xF6, 0xF9)
TEXT   = RGBColor(0x2D, 0x2D, 0x2D)
MUTED  = RGBColor(0x6B, 0x72, 0x80)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
AMBER  = RGBColor(0xE0, 0x8A, 0x00)
RED    = RGBColor(0xC0, 0x39, 0x2B)
GREY   = RGBColor(0xC4, 0xCB, 0xD4)
RAG = {"GREEN": GREEN, "AMBER": AMBER, "RED": RED}
FONT = "Segoe UI"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
EMU_W = Inches(13.333)
BLANK = prs.slide_layouts[6]


def _set(run, size, color, bold=False, italic=False):
    run.font.size = Pt(size); run.font.color.rgb = color
    run.font.bold = bold; run.font.italic = italic; run.font.name = FONT


def rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp


def tb(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h); tf = box.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(3); tf.margin_top = tf.margin_bottom = Pt(1)
    for i, (t, s, c, b, it) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = t; _set(r, s, c, b, it)
    return box


def header(slide, title, tag):
    rect(slide, 0, 0, EMU_W, Inches(0.95), DARK)
    rect(slide, 0, Inches(0.95), EMU_W, Inches(0.05), ACCENT)
    tb(slide, Inches(0.5), Inches(0.12), Inches(9.5), Inches(0.7),
       [(title, 24, WHITE, True, False)], MSO_ANCHOR.MIDDLE)
    chip = rect(slide, Inches(10.6), Inches(0.28), Inches(2.2), Inches(0.42), ACCENT)
    ctf = chip.text_frame; ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = ctf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = tag; _set(r, 11, WHITE, True)
    tb(slide, Inches(0.5), Inches(7.08), Inches(12.3), Inches(0.35),
       [(CAVEAT, 8.5, MUTED, False, True)])


def stat_banner(slide, y):
    # portfolio headline numbers
    cards = [(str(TOT_SCOPE), "stories in scope", DARK),
             (str(TOT_MOVING), "moving", GREEN),
             (str(TOT_STUCK), "need attention", RED)]
    x = Inches(0.5); w = Inches(2.7)
    for num, lbl, col in cards:
        rect(slide, x, y, w, Inches(0.95), LIGHT)
        rect(slide, x, y, Inches(0.1), Inches(0.95), col)
        tb(slide, x + Inches(0.25), y + Inches(0.05), w - Inches(0.3), Inches(0.6),
           [(num, 30, col, True, False)])
        tb(slide, x + Inches(0.25), y + Inches(0.6), w - Inches(0.3), Inches(0.3),
           [(lbl, 11, TEXT, False, False)])
        x += w + Inches(0.25)
    return x


# ============================================================ VAR A : bars + counts
s = prs.slides.add_slide(BLANK)
header(s, "Objectives - Story Progress (Variation A)", "BARS + COUNTS")
stat_banner(s, Inches(1.2))
y = Inches(2.45)
unit = Inches(4.2) / 6.0  # 6 = max in_scope -> scale bar length to magnitude
for pri, name, scope, moving, stuck, reason in OBJS:
    tb(s, Inches(0.5), y, Inches(0.6), Inches(0.5),
       [(pri, 13, ACCENT, True, False)], MSO_ANCHOR.MIDDLE)
    tb(s, Inches(1.1), y, Inches(4.0), Inches(0.5),
       [(name, 13, TEXT, True, False)], MSO_ANCHOR.MIDDLE)
    bx = Inches(5.2); bh = Inches(0.32); by = y + Inches(0.08)
    mlen = int(moving * unit); slen = int(stuck * unit)
    if moving:
        rect(s, bx, by, mlen, bh, GREEN)
    if stuck:
        rect(s, bx + mlen, by, slen, bh, RED)
    tb(s, Inches(10.0), y, Inches(3.0), Inches(0.5),
       [(f"{scope} in scope  /  {moving} moving  /  {stuck} stuck", 10.5, TEXT, False, False)],
       MSO_ANCHOR.MIDDLE)
    y += Inches(0.62)
tb(s, Inches(5.2), y + Inches(0.05), Inches(7), Inches(0.3),
   [("Bar length = stories in scope (magnitude);  green = moving, red = needs attention.", 9.5, MUTED, False, True)])


# ============================================================ VAR B : status + succinct why
s = prs.slides.add_slide(BLANK)
header(s, "Objectives - Status Explained (Variation B)", "SELF-EXPLAINING")
stat_banner(s, Inches(1.2))
y = Inches(2.4)
for pri, name, scope, moving, stuck, reason in OBJS:
    rag = rag_of(stuck, moving)
    rect(s, Inches(0.5), y, Inches(12.3), Inches(0.68), LIGHT)
    rect(s, Inches(0.5), y, Inches(0.12), Inches(0.68), RAG[rag])
    # status chip
    chip = rect(s, Inches(0.75), y + Inches(0.16), Inches(1.15), Inches(0.36), RAG[rag])
    ctf = chip.text_frame; ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = ctf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = rag; _set(r, 10, WHITE, True)
    tb(s, Inches(2.05), y + Inches(0.04), Inches(3.1), Inches(0.6),
       [(f"{pri}  {name}", 12.5, DARK, True, False)], MSO_ANCHOR.MIDDLE)
    tb(s, Inches(5.2), y + Inches(0.04), Inches(6.0), Inches(0.6),
       [(reason, 11, TEXT, False, False)], MSO_ANCHOR.MIDDLE)
    tb(s, Inches(11.3), y + Inches(0.04), Inches(1.5), Inches(0.6),
       [(f"{moving}/{scope}", 14, RAG[rag], True, False),
        ("moving", 8, MUTED, False, False)], MSO_ANCHOR.MIDDLE, PP_ALIGN.CENTER)
    y += Inches(0.74)


# ============================================================ VAR C : executive tiles
s = prs.slides.add_slide(BLANK)
header(s, "Objectives - Executive Tiles (Variation C)", "VISUAL GRID")
tw, th = Inches(4.0), Inches(2.55)
gx, gy = Inches(0.5), Inches(1.25)
for i, (pri, name, scope, moving, stuck, reason) in enumerate(OBJS):
    rag = rag_of(stuck, moving)
    col = i % 3; row = i // 3
    x = gx + col * (tw + Inches(0.3))
    y = gy + row * (th + Inches(0.25))
    rect(s, x, y, tw, th, LIGHT)
    rect(s, x, y, tw, Inches(0.55), RAG[rag])
    tb(s, x + Inches(0.2), y + Inches(0.05), tw - Inches(0.4), Inches(0.45),
       [(f"{pri}   {name}", 13, WHITE, True, False)], MSO_ANCHOR.MIDDLE)
    # big stat
    tb(s, x + Inches(0.2), y + Inches(0.62), tw - Inches(0.4), Inches(0.9),
       [(str(scope), 44, DARK, True, False)])
    tb(s, x + Inches(1.4), y + Inches(0.95), tw - Inches(1.5), Inches(0.6),
       [("stories", 12, MUTED, False, False),
        (f"{moving} moving  ·  {stuck} stuck", 11, TEXT, True, False)])
    tb(s, x + Inches(0.2), y + Inches(1.75), tw - Inches(0.4), Inches(0.75),
       [(reason, 10, TEXT, False, False)])

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "2026-06-11-scorecard-variations.pptx")
prs.save(out)
print("Saved", out)
