"""
CMDB Audit Readiness deck - business-stakeholder view.

FOCUS on audit-required attributes, but does NOT exclude the other attributes:
the exec Scorecard shows the audit attributes only; each per-class slide shows
EVERY attribute (audit ones in colour vs. the 90% target, other tracked ones
greyed - managed but not audit-gated).

Audit attributes (confirmed):
  Business Applications : CI Owner, Classification
  Servers               : CI Owner, Location, IP Address
  Computers             : Assigned To, Location, Serial Number

Per-class Status & Remediation slides for Business Applications and Computers
are grounded in the team's own gap-remediation spikes:
  - BA:       ADO Spike 1480111 (owner Joe Dames)      - see ../Backlog/audit-dashboard-accuracy-spike.md
  - Computer: ADO Spike 1480114 (owner A. de Araujo)   - see ../Backlog/audit-dashboard-computer-spike.md
Servers is audit-ready (all 3 audit attributes at target) - brief slide, no remediation.

Data as of 7/16/2026. Compliance % is the audit metric (raw counts shift
week-over-week as the active-CI population changes).

Matches the house deck style (see ../Contract/build_co6_deliverable_deck.py).

Single source of truth. Edit CONTENT, then run:
    python build_audit_readiness_deck.py
Output:
    CMDB-Audit-Readiness-<EDITION>.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# CONTENT
# ---------------------------------------------------------------------------
EDITION   = "2026-07-16"
TITLE     = "CMDB Audit Readiness"
SUBTITLE  = "Audit-required attributes in focus - full attribute status by CI class"
TAGLINE   = "Bridging Gaps to Achieve CMDB Excellence"
AUTHOR    = "Manuel Vazquez  -  Scrum Master, CMDB-CSDM"
FOOTLABEL = "PPL CMDB-CSDM  |  Audit Readiness  |  Attribute Completeness  -  Jul 16, 2026"
TARGET    = 90

# Exec scorecard: audit attributes only. (CI class, total, [(attr, %, gap, rag)])
SCORE = [
    ("Business Applications", "2,148", [
        ("CI Owner",       100, "11",    "green"),
        ("Classification",  36, "1,371", "red"),
    ]),
    ("Servers", "3,917", [
        ("CI Owner",    100, "0",   "green"),
        ("Location",    100, "1",   "green"),
        ("IP Address",   91, "358", "green"),
    ]),
    ("Computers", "19,901", [
        ("Assigned To",     83, "3,412", "amber"),
        ("Location",        66, "6,748", "red"),
        ("Serial Number",  100, "0",     "green"),
    ]),
]

# Per-class content. audit=[(name,%,gap,rag)]  other=[(name,%)]
BA = dict(
    name="Business Applications", total="2,148", spike="ADO Spike 1480111",
    audit=[("CI Owner", 100, "11", "green"), ("Classification", 36, "1,371", "red")],
    other=[("Technical Owner Group", 89), ("Business Owner", 84), ("Approval Group", 81),
           ("Support Group", 77), ("App Code / Asset Tag", 100), ("Value Stream", 62), ("SOX Type", 100)],
    summary="1 of 2 audit attributes at target. CI Owner is effectively complete; Classification (36%) is the single audit gap - and the only audit attribute that has not moved since June.",
    plan=[
        (["Classification", "1,371 CIs  ·  36%"],
         "Governance / manual attribute - set per application by business owners; no automated discovery source. Remediate as a governed, change-controlled data story, classified by business criticality.",
         ["Joe Dames - PO", "Todd Dierksheide - CCB", "Target: TBD with CCB"]),
    ],
    note=[("Constraint:  ", "red", True),
          ("Classification remediation is gated by the 6/10 Business Application enhancement pause - mass updates must go through a governed story + change control, not inline edits. This is the root of the stall.", None, False)],
)

SRV = dict(
    name="Servers", total="3,917", spike=None,
    audit=[("CI Owner", 100, "0", "green"), ("Location", 100, "1", "green"), ("IP Address", 91, "358", "green")],
    other=[("Technical Owner Group", 100), ("Support Group", 100), ("App Code / Asset Tag", 75),
           ("Value Stream", 0), ("SOX Type", 99), ("Classification", 99)],
    summary="Audit-ready. All three audit attributes are at or above the 90% target as of 7/16 - CI Owner 100%, Location 100%, IP Address 91%. Location and IP Address were remediated this week (44% -> 100% and 43% -> 91%).",
    action="No audit remediation required - sustain current levels through ongoing discovery and governance.",
    note=[("Not audit-gated:  ", None, True),
          ("Value Stream (0%) and App Code / Asset Tag (75%) are tracked but sit outside the audit target - addressed under separate data-quality work.", None, False)],
)

COMP = dict(
    name="Computers", total="19,901", spike="ADO Spike 1480114",
    audit=[("Assigned To", 83, "3,412", "amber"), ("Location", 66, "6,748", "red"), ("Serial Number", 100, "0", "green")],
    other=[("CI Owner", 100), ("Technical Owner Group", 100), ("Support Group", 100),
           ("App Code / Asset Tag", 64), ("Operating System", 99), ("IP Address", 19)],
    summary="1 of 3 audit attributes at target (Serial Number 100%). Location (66%) and Assigned To (83%) are the open audit gaps - both improving: Location 38% -> 66% and Assigned To 71% -> 83% since early July.",
    plan=[
        (["Location", "6,748 CIs  ·  66%"],
         "Auto-populated by device discovery - gaps are discovery coverage (offline physical devices may not report). Close at source via the discovery / SCCM pipeline, not manual entry.",
         ["Anthony de Araujo", "Monica Green / Paul Becker - CCB", "On trajectory to ≥90%"]),
        (["Assigned To", "3,412 CIs  ·  83%"],
         "Discovery-driven endpoint assignment; within 7 points of target and climbing. Sustain the same pipeline remediation.",
         ["Anthony de Araujo", "On trajectory"]),
    ],
    note=[("", None, False),
          ("Retired-device cleanup reduced the active population and lifted compliance; discovery attributes (incl. IP Address, App Code) improve together. Validate physical vs. virtual coverage.", None, False)],
)

# Bottom line: (headline, body, rag)
BOTTOM = [
    ("Servers - audit-ready.",
     "All three audit attributes (CI Owner, Location, IP Address) are at or above target; Location and IP Address were remediated this week. Sustain via discovery.",
     "green"),
    ("Computers - on trajectory.",
     "Serial Number at target. Location (66%) and Assigned To (83%) are climbing via the device-discovery pipeline and retired-device cleanup (Spike 1480114). On track.",
     "amber"),
    ("Business Applications - one stalled gap.",
     "CI Owner at target. Classification (36%) is the stalled audit gap - governance/manual data blocked by the BA-enhancement pause; needs a governed story + CCB owner (Spike 1480111).",
     "red"),
]
FOCUS = ("Unblock Business Application Classification (governed story with Joe Dames + CCB / Todd Dierksheide) and set a target date; "
         "sustain the Computers device-discovery remediation for Location & Assigned To. Servers are audit-ready.")

# ---------------------------------------------------------------------------
# THEME  (matches ../Contract/build_co6_deliverable_deck.py)
# ---------------------------------------------------------------------------
DARK   = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0x2E, 0x6D, 0xB4)
LIGHT  = RGBColor(0xF4, 0xF6, 0xF9)
TEXT   = RGBColor(0x2D, 0x2D, 0x2D)
MUTED  = RGBColor(0x6B, 0x72, 0x80)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
AMBER  = RGBColor(0xE0, 0x8A, 0x00)
RED    = RGBColor(0xC0, 0x39, 0x2B)
GREY   = RGBColor(0x8A, 0x92, 0x9E)
RULE   = RGBColor(0xDD, 0xE3, 0xEA)

RAG       = {"green": GREEN, "amber": AMBER, "red": RED}
RAG_LABEL = {"green": "On target", "amber": "At risk", "red": "Critical"}

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
FONT = "Segoe UI"


def _set(run, size, color, bold=False, italic=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = FONT


def rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def textbox(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, spec in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        text, size, color, bold, italic = spec
        r = p.add_run(); r.text = text
        _set(r, size, color, bold, italic)
    return tb


def textbox_rich(slide, x, y, w, h, paragraphs, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, runs in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0); p.space_after = Pt(1)
        for (text, size, color, bold, italic) in runs:
            r = p.add_run(); r.text = text
            _set(r, size, color, bold, italic)
    return tb


def header(slide, title, subtitle=None):
    rect(slide, 0, 0, EMU_W, Inches(1.05), DARK)
    rect(slide, 0, Inches(1.05), EMU_W, Inches(0.06), ACCENT)
    textbox(slide, Inches(0.55), Inches(0.18), Inches(12.2), Inches(0.55),
            [(title, 24, WHITE, True, False)], MSO_ANCHOR.MIDDLE)
    if subtitle:
        textbox(slide, Inches(0.57), Inches(0.66), Inches(12.2), Inches(0.32),
                [(subtitle, 12, RGBColor(0xC9, 0xD6, 0xE6), False, False)], MSO_ANCHOR.MIDDLE)
    textbox(slide, Inches(0.55), Inches(7.08), Inches(10), Inches(0.32),
            [(FOOTLABEL, 9, MUTED, False, False)])


def cell(t, i, j, lines, size, color, bold=False, align=PP_ALIGN.LEFT, fill=None, anchor=None):
    c = t.cell(i, j)
    if fill is not None:
        c.fill.solid(); c.fill.fore_color.rgb = fill
    if anchor is not None:
        c.vertical_anchor = anchor
    tf = c.text_frame; tf.word_wrap = True
    if isinstance(lines, str):
        lines = [lines]
    for k, ln in enumerate(lines):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0); p.space_after = Pt(2)
        r = p.add_run(); r.text = ln
        _set(r, size, color, bold)


def legend(slide, x, y):
    for label, col in [("≥ 90%  On target", GREEN), ("80-89%  At risk", AMBER), ("< 80%  Critical", RED)]:
        rect(slide, x, y + Inches(0.03), Inches(0.16), Inches(0.16), col)
        textbox(slide, x + Inches(0.22), y - Inches(0.02), Inches(2.4), Inches(0.28),
                [(label, 10.5, TEXT, False, False)], MSO_ANCHOR.MIDDLE)
        x += Inches(2.4)
    textbox(slide, x + Inches(0.2), y - Inches(0.02), Inches(4.6), Inches(0.28),
            [("Target: ≥ 90% on every audit attribute", 10.5, MUTED, False, True)],
            MSO_ANCHOR.MIDDLE, PP_ALIGN.RIGHT)


def audit_chip(s, x, y, w, name, pct, gap, color):
    rect(s, Inches(x), Inches(y), Inches(w), Inches(0.92), LIGHT)
    rect(s, Inches(x), Inches(y), Inches(w), Inches(0.08), color)
    textbox(s, Inches(x), Inches(y + 0.13), Inches(w), Inches(0.3),
            [(name, 10, TEXT, True, False)], MSO_ANCHOR.MIDDLE, PP_ALIGN.CENTER)
    textbox(s, Inches(x), Inches(y + 0.4), Inches(w), Inches(0.36),
            [("{}%".format(pct), 25, color, True, False)], MSO_ANCHOR.MIDDLE, PP_ALIGN.CENTER)
    textbox(s, Inches(x), Inches(y + 0.72), Inches(w), Inches(0.18),
            [("Gap: {} CIs".format(gap), 8.5, MUTED, False, False)], MSO_ANCHOR.TOP, PP_ALIGN.CENTER)


def other_chip(s, x, y, w, name, pct):
    rect(s, Inches(x), Inches(y), Inches(w), Inches(0.74), LIGHT)
    rect(s, Inches(x), Inches(y), Inches(w), Inches(0.07), GREY)
    textbox(s, Inches(x), Inches(y + 0.1), Inches(w), Inches(0.34),
            [(name, 8.5, TEXT, True, False)], MSO_ANCHOR.MIDDLE, PP_ALIGN.CENTER)
    textbox(s, Inches(x), Inches(y + 0.42), Inches(w), Inches(0.28),
            [("{}%".format(pct), 12.5, MUTED, True, False)], MSO_ANCHOR.MIDDLE, PP_ALIGN.CENTER)


def attribute_bands(s, d):
    """Shared top of every per-class slide: audit chips + other chips + summary."""
    textbox_rich(s, Inches(0.55), Inches(1.16), Inches(12.2), Inches(0.26),
                 [[("AUDIT ATTRIBUTES", 11.5, DARK, True, False),
                   ("     measured against the ≥ 90% target", 10, MUTED, False, True)]])
    aw, agap = 2.8, 0.25
    for i, (name, pct, gap, rag) in enumerate(d["audit"]):
        audit_chip(s, 0.55 + i * (aw + agap), 1.46, aw, name, pct, gap, RAG[rag])
    textbox(s, Inches(0.55), Inches(2.46), Inches(12.2), Inches(0.36),
            [(d["summary"], 11, TEXT, False, False)])
    textbox_rich(s, Inches(0.55), Inches(2.9), Inches(12.2), Inches(0.24),
                 [[("OTHER TRACKED ATTRIBUTES", 11.5, GREY, True, False),
                   ("     managed - not part of the audit target", 10, MUTED, False, True)]])
    n = len(d["other"]); ogap = 0.15
    ow = (12.2 - (n - 1) * ogap) / n
    for i, (name, pct) in enumerate(d["other"]):
        other_chip(s, 0.55 + i * (ow + ogap), 3.18, ow, name, pct)
    rect(s, Inches(0.55), Inches(4.06), Inches(12.2), Inches(0.016), RULE)


def note_para(s, y, note, size=10.5):
    textbox_rich(s, Inches(0.55), Inches(y), Inches(12.2), Inches(0.5),
                 [[(txt, size, (RED if col == "red" else TEXT), bold, False) for (txt, col, bold) in note]])


# ---------------------------------------------------------------------------
# SLIDE BUILDERS
# ---------------------------------------------------------------------------
def add_title(prs, blank):
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, EMU_W, EMU_H, DARK)
    rect(s, 0, Inches(4.55), EMU_W, Inches(0.08), ACCENT)
    textbox(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(1.1), [(TITLE, 44, WHITE, True, False)])
    textbox(s, Inches(0.82), Inches(2.82), Inches(11.7), Inches(0.6), [(SUBTITLE, 19, RGBColor(0xC9, 0xD6, 0xE6), False, False)])
    textbox(s, Inches(0.82), Inches(3.42), Inches(11.7), Inches(0.5), [(TAGLINE, 15, RGBColor(0x9F, 0xB2, 0xCC), False, True)])
    textbox(s, Inches(0.82), Inches(4.75), Inches(11.7), Inches(1.5), [
        ("The attributes PPL Audit asked us to track - measured against a 90% completeness target", 16, WHITE, True, False),
        ("As of July 16, 2026  ·  25,966 configuration items across Business Applications, Servers & Computers", 14, RGBColor(0xC9, 0xD6, 0xE6), False, False),
        (AUTHOR, 12, RGBColor(0x9F, 0xB2, 0xCC), False, False),
    ])


def add_scorecard(prs, blank):
    s = prs.slides.add_slide(blank)
    header(s, "Audit Readiness Scorecard",
           "Audit-required attributes vs. the 90% target  ·  full attribute detail per class follows  ·  as of Jul 16, 2026")
    legend(s, Inches(0.55), Inches(1.28))
    nrows = 1 + sum(1 + len(a) for _, _, a in SCORE)
    t = s.shapes.add_table(nrows, 4, Inches(0.55), Inches(1.72), Inches(12.2), Inches(4.5)).table
    for j, w in enumerate([Inches(4.7), Inches(2.3), Inches(2.2), Inches(3.0)]):
        t.columns[j].width = w
    for j, (h, a) in enumerate([("Audit Attribute", PP_ALIGN.LEFT), ("Complete", PP_ALIGN.CENTER),
                                ("Gap (CIs)", PP_ALIGN.CENTER), ("Status", PP_ALIGN.CENTER)]):
        cell(t, 0, j, h, 11.5, WHITE, True, a, DARK, MSO_ANCHOR.MIDDLE)
    t.rows[0].height = Inches(0.4)
    r = 1
    for cls, total, attrs in SCORE:
        t.cell(r, 0).merge(t.cell(r, 3))
        cell(t, r, 0, "{}      ·      {} configuration items".format(cls, total),
             11, WHITE, True, PP_ALIGN.LEFT, ACCENT, MSO_ANCHOR.MIDDLE)
        t.rows[r].height = Inches(0.34); r += 1
        for i, (name, pct, gap, rag) in enumerate(attrs):
            fill = LIGHT if i % 2 == 0 else WHITE
            cell(t, r, 0, name, 11, TEXT, False, PP_ALIGN.LEFT, fill, MSO_ANCHOR.MIDDLE)
            cell(t, r, 1, "{}%".format(pct), 13, RAG[rag], True, PP_ALIGN.CENTER, fill, MSO_ANCHOR.MIDDLE)
            cell(t, r, 2, gap, 11, TEXT, False, PP_ALIGN.CENTER, fill, MSO_ANCHOR.MIDDLE)
            cell(t, r, 3, RAG_LABEL[rag].upper(), 10.5, WHITE, True, PP_ALIGN.CENTER, RAG[rag], MSO_ANCHOR.MIDDLE)
            t.rows[r].height = Inches(0.38); r += 1
    textbox(s, Inches(0.55), Inches(6.42), Inches(12.2), Inches(0.55),
            [("5 of 8 audit attributes are at target.  Servers are fully compliant; the open gaps are "
              "Business Application Classification and Computer Location / Assigned To - detailed on the following slides.",
              11, TEXT, False, False)])


def add_class_detail(prs, blank, d):
    s = prs.slides.add_slide(blank)
    header(s, d["name"] + " - Status & Remediation",
           "All attributes shown  ·  audit attributes vs. the 90% target  ·  {} CIs  ·  as of 7/16".format(d["total"]))
    attribute_bands(s, d)
    textbox(s, Inches(0.55), Inches(4.14), Inches(12.2), Inches(0.26),
            [("REMEDIATION ACTION PLAN", 11.5, ACCENT, True, False)])
    rows = 1 + len(d["plan"])
    t = s.shapes.add_table(rows, 3, Inches(0.55), Inches(4.44), Inches(12.2), Inches(0.4 * rows)).table
    for j, w in enumerate([Inches(2.6), Inches(6.7), Inches(2.9)]):
        t.columns[j].width = w
    for j, h in enumerate(["Audit gap", "Root cause & remediation approach", "Owner  ·  Target"]):
        cell(t, 0, j, h, 10.5, WHITE, True, PP_ALIGN.LEFT, DARK, MSO_ANCHOR.MIDDLE)
    t.rows[0].height = Inches(0.36)
    for i, (gap, approach, owner) in enumerate(d["plan"], start=1):
        fill = LIGHT if i % 2 else WHITE
        cell(t, i, 0, gap, 10, DARK, True, PP_ALIGN.LEFT, fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 1, approach, 9.5, TEXT, False, PP_ALIGN.LEFT, fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 2, owner, 9.5, TEXT, False, PP_ALIGN.LEFT, fill, MSO_ANCHOR.MIDDLE)
        t.rows[i].height = Inches(0.78)
    ny = 4.44 + 0.36 + 0.78 * len(d["plan"]) + 0.14
    note_para(s, ny, d["note"])
    if d.get("spike"):
        textbox(s, Inches(2.55), Inches(7.08), Inches(10.2), Inches(0.24),
                [("Remediation tracked under " + d["spike"], 8.5, MUTED, False, True)],
                MSO_ANCHOR.TOP, PP_ALIGN.RIGHT)


def add_class_brief(prs, blank, d):
    s = prs.slides.add_slide(blank)
    header(s, d["name"] + " - Status",
           "All attributes shown  ·  audit-ready  ·  {} CIs  ·  as of 7/16".format(d["total"]))
    attribute_bands(s, d)
    textbox(s, Inches(0.55), Inches(4.2), Inches(12.2), Inches(0.26),
            [("SUMMARY & STATUS", 11.5, ACCENT, True, False)])
    textbox(s, Inches(0.55), Inches(4.5), Inches(12.2), Inches(0.8), [(d["summary"], 12, TEXT, False, False)])
    rect(s, Inches(0.55), Inches(5.5), Inches(0.14), Inches(0.5), GREEN)
    textbox(s, Inches(0.8), Inches(5.5), Inches(11.9), Inches(0.5), [(d["action"], 12, DARK, True, False)])
    note_para(s, 6.35, d["note"])


def add_bottom(prs, blank):
    s = prs.slides.add_slide(blank)
    header(s, "Bottom Line", "Where each CI class stands against the audit target")
    yy = 1.5
    for head_txt, body, rag in BOTTOM:
        rect(s, Inches(0.55), Inches(yy + 0.03), Inches(0.14), Inches(0.92), RAG[rag])
        textbox(s, Inches(0.8), Inches(yy), Inches(11.9), Inches(0.35), [(head_txt, 15, DARK, True, False)])
        textbox(s, Inches(0.8), Inches(yy + 0.37), Inches(11.9), Inches(0.6), [(body, 11.5, TEXT, False, False)])
        yy += 1.18
    rect(s, Inches(0.55), Inches(yy + 0.05), Inches(12.2), Inches(0.98), LIGHT)
    rect(s, Inches(0.55), Inches(yy + 0.05), Inches(0.14), Inches(0.98), ACCENT)
    textbox_rich(s, Inches(0.82), Inches(yy + 0.16), Inches(11.8), Inches(0.78),
                 [[("Immediate focus", 12, ACCENT, True, False)],
                  [(FOCUS, 11.5, TEXT, False, False)]])


def build_deck():
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    blank = prs.slide_layouts[6]
    add_title(prs, blank)
    add_scorecard(prs, blank)
    add_class_detail(prs, blank, BA)
    add_class_brief(prs, blank, SRV)
    add_class_detail(prs, blank, COMP)
    add_bottom(prs, blank)
    return prs


def save_safe(prs, path):
    try:
        prs.save(path); print("Saved", os.path.basename(path), "-", len(prs.slides._sldIdLst), "slides")
    except PermissionError:
        base, ext = os.path.splitext(path)
        alt = base + "-new" + ext
        prs.save(alt); print("Target locked - saved to", os.path.basename(alt))


here = os.path.dirname(os.path.abspath(__file__))
save_safe(build_deck(), os.path.join(here, "CMDB-Audit-Readiness-" + EDITION + ".pptx"))
