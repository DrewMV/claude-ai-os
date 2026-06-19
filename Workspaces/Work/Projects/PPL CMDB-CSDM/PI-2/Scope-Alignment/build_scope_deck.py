"""
PI-2 Scope & Delivery Focus deck generator  (CMDB team working session).

Audience: the CMDB team doing the work and making solution decisions.
Purpose: establish the initial (contracted) scope, the scope added in PI-2,
and the delivery details (priority, live work, capacity, roadmap, risks,
and the decisions the team owns) — to give clarity and direction.

Single source of truth for this deck. Edit the CONTENT block, then run:
    python build_scope_deck.py
Output:  <EDITION>-scope-focus.pptx

Theme + helpers mirror build_status_deck.py so the two decks look like a set.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# CONTENT
# ---------------------------------------------------------------------------
EDITION   = "2026-06-18"
TITLE     = "PI-2 Scope & Delivery Focus"
SUBTITLE  = "CMDB Team Working Session"
PI        = "PI-2"
PI_WINDOW = "May 13 - Aug 4, 2026"
ITERATION = "Iteration 2.3  (Jun 10 - 23)"
AUTHOR    = "Manuel Vazquez  -  Scrum Master, CMDB-CSDM"
FOOTLABEL = "PPL CMDB-CSDM  |  PI-2  |  Scope & Delivery Focus  -  Jun 18, 2026"

PURPOSE = ("This deck aligns what we are building to what the engagement committed to. "
           "It establishes the original contracted scope, the scope added during PI-2, and "
           "where to focus our effort and solution decisions for the rest of the increment.")

# Initial contracted scope: WS#, name, what it covers, points
ORIGINAL = [
    ("WS #1", "Governance", "Data dictionary / class attributes (Servers, Computers, Business Apps, Databases); SOX / ESS-02; CCB facilitation", "27 pts"),
    ("WS #2", "Automated Data Ingestion", "SCCM / Discovery precedence; 90% coverage validation (Computers & Servers); Oracle / MS-SQL discovery; Life Cycle bulk updates", "33 pts"),
    ("WS #3", "Other Enhancements", "Evaluate ONE integration - Tanium OR Qualys; asset / service-mapping data only (no vulnerability ingestion)", "7 pts"),
    ("WS #4", "Service Mapping", "Dev 4 fully dedicated; 5 maps / sprint x 6 sprints = 30 service maps", "48 pts"),
]
ORIGINAL_TOTAL = "115 pts total  (67 general pool  +  48 Dev 4 dedicated)"
TERM_NOTE = ("All CO#5 deliverables are due June 30, with a $533,775 holdback released on acceptance.  "
             "An extension covering July / August (Sprints 5-6) is expected but NOT yet signed.")

# Added/aligned in PI-2: #, objective, verdict word, verdict color key, note
# Basis = signed Change Order #5 (Mar 31 - Jun 30, 2026) + prior COs.
ALIGNMENT = [
    ("6", "VMware -> Azure Migration", "ADDED", "RED",   "Not in CO#5 - added, critical. Needs a change order. 100% VMs discovered & mapped pre-migration."),
    ("1", "Network Gear Discovery",    "VERIFY", "AMBER","Believed contracted (prior CO) - NOT named in CO#5. Confirm against CO#1-4. >=90% network-gear quality."),
    ("2", "CI Data Certification",     "CONTRACTED", "GREEN","CO#5 Deliverable 1.2 (named). Process, intake, Business App pilot, KB docs."),
    ("3", "Service Mapping Expansion", "CONTRACTED", "GREEN","Contracted (prior CO). 30 maps target - exact commitment shape TBD."),
    ("4", "NERC-CIP & Qualys",         "PARTIAL", "AMBER","CO#5 buys Qualys REQUIREMENTS only; team is building (verify build CO). NERC-CIP explicitly EXCLUDED."),
    ("5", "NowAssist AI Embedding",    "ADDED", "RED",   "Not in CO#5 - added. ~17 pts. Needs a change order."),
]
ALIGN_SUMMARY = "Basis: signed CO#5 (Mar 31 - Jun 30, 2026)   ·   2 contracted · 1 partial · 2 added · 1 to verify"
GAPS = [
    "WS #2 Auto Data Ingestion (CO#5 Deliverable 2) - contracted and delivered at task level, but has no PI-2 objective of its own. Give it a tracked goal.",
    "SOX / ESS-02 (CO#5 Deliverable 1) - contracted, but not explicitly called out as its own PI-2 objective.",
]

# Priority board: tier, pri label, value label, objective, target, status word, status color, origin
BOARD = [
    (0, "P0", "BV 10", "VMware -> Azure Migration",        "100% in-scope VMs discovered, classified & relationship-mapped pre-migration", "At risk", "AMBER", "Added in PI-2"),
    (1, "P1", "BV 10", "Network Gear Discovery",           ">=90% of network-gear CIs meet data-quality standard by Aug 4",                "At risk", "AMBER", "Contracted? verify"),
    (1, "P1", "Core",  "Server / Computer Data Ingestion", "SCCM precedence + 90% Computer/Server coverage + class data & form updates",    "On track", "GREEN", "Original WS #2"),
    (1, "P1", "BV 8",  "CI Data Certification",            "Framework defined + Business App pilot + repeatable cadence",                   "On track", "GREEN", "Original WS #1"),
    (2, "P2", "BV 9",  "Service Mapping Expansion",        "30 maps (10/month) of business-critical apps, by criticality",                 "On track", "GREEN", "Original WS #4"),
    (2, "P2", "BV 8",  "NERC-CIP & Qualys Foundation",     "100% foundational data model + integration approach defined",                  "Blocked", "RED",   "Partial WS #3"),
    (3, "P3", "BV 7",  "NowAssist AI Embedding",           "Embed in key CMDB/ITSM touchpoints; measurable quality gain",                  "At risk", "AMBER", "Added - defer candidate"),
]

# Drill-down cards: tier, pri, name, rollup, [detail lines], status color
DRILL_P0P1 = [
    (0, "P0", "VMware -> Azure Migration", "3 Validation · 2 Active · Feat 1420613", [
        "Airlift Pre/During/Post (1418610/18/21) - Validation;  CI auto-populate 1416384 - Active",
        "Live driver: Ray's 450-server app-instance request (82/450 covered) - Risk #13",
    ], "AMBER"),
    (1, "P1", "Network Gear Discovery", "2 Validation · 1 Active · 2 New · Feat 1356646", [
        "Spikes 1402555 (Active) / 1402559 (New);  credentials 1444864 (Validation, 6 tasks still in 2.2)",
        "Dep #1 network-gear stakeholder requirements (1383487) still New",
    ], "AMBER"),
    (1, "P1", "Server / Computer Data Ingestion  (WS #2 - contracted)", "6 Validation · 4 Active", [
        "SCCM precedence 1403759 / 1403760 / 1403762 / 1403763 - all Validation",
        "Class data: 1454371, 1355167, 1455858 Active;  1407572, 1421790, 1387236 Validation",
    ], "GREEN"),
    (1, "P1", "CI Data Certification", "1 UAT · 1 Validation · 3 empty features", [
        "Dashboard 1402727 - UAT signed, PROD 6/23;  Pilot changes 1435307 - Validation",
        "Watch: Feats 1371672 / 1382404 / 1402958 are shells, no stories yet",
    ], "GREEN"),
]
DRILL_P2P3 = [
    (2, "P2", "Service Mapping Expansion", "4 maps active · 2 spikes · Dev 4 dedicated", [
        "WATT / Vault Inspection / MV90 / RIE active;  spikes 1326754, 1420634",
        "Watch: ~36 carryover map tasks stranded in 2.1/2.2;  Wave 20 & 21 empty",
    ], "GREEN"),
    (2, "P2", "NERC-CIP & Qualys Foundation", "2 Blocked · 1 Active · NERC = dep only", [
        "Qualys 1428703 / 1428704 BLOCKED - Issue 1465952 (vendor plugin approval)",
        "ESS-02 1420244 Active;  NERC-CIP = dependency 1383515 only, no build stories",
    ], "RED"),
    (3, "P3", "NowAssist AI Embedding", "1 Closed · 1 Resolved · 2 Active · Feat 1436574", [
        "1436576 Closed, 1436579 Resolved, 1436592 / 1470837 Active;  1436593 / 1436581 Ready",
        "Watch: ~17-pt admin / platform-setup prerequisite not yet storied",
    ], "AMBER"),
]
UNMAPPED = ("Not mapped to any objective - capacity leaking off-board",
            "OCM 1428659 (Active, Nora Lizenberg) - no objective home   ·   "
            "orphan stories 1472365 & 1399787 - no parent / no feature")

# Capacity model
CAP_ROWS = [
    ("Original scope (contracted)", "115 pts", "67 general pool  +  48 Dev 4 dedicated"),
    ("PI-2 additions (general pool)", "~55 pts", "Network Gear + NowAssist + VMware - absorbed into the buffer"),
    ("General-pool demand", "122 / 144", "Buffer ~22 pts - but front/back-loaded"),
]
CAP_SPRINTS = [  # sprint, window, used/avail, status word, status color
    ("S1", "May 13-23",  "10 / 24", "Under",       "GREEN"),
    ("S2", "May 27-Jun 6","19 / 24", "Under",       "GREEN"),
    ("S3", "Jun 10-20",  "23 / 24", "Near",        "AMBER"),
    ("S4", "Jun 24-Jul 4","24 / 24", "At capacity", "RED"),
    ("S5", "Jul 8-18",   "24 / 24", "At capacity", "RED"),
    ("S6", "Jul 22-Aug 1","22 / 24", "Under",       "GREEN"),
]
CAP_NOTE = ("S4-S5 run at 100% (zero buffer).  Dev code freeze Jun 27 - Jul 18 covers most of S4 + all of S5;  "
            "Test code freeze Jul 18 - Aug 15 covers all of S6.  The back half is the crunch.  "
            "Note: CO#5 term ends Jun 30 (mid-S4) - S5/S6 sit beyond the signed term, pending a not-yet-signed extension.")

# Roadmap: plan-vs-actual slippage callouts
ROADMAP_PLAN = ("Plan phases:  Initiation (S1)  ->  Foundation Build (S2-S3)  ->  Core Delivery (S4-S5)  ->  Finalization & Close (S6).  "
                "Contract reality: CO#5 ends Jun 30 (mid-S4) - everything in S5/S6 sits beyond the signed term, pending an extension.")
SLIPPAGE = [
    ("NowAssist admin 'complete by S2 (Jun 6)'", "Not started - and it gates S5 workflow embedding."),
    ("Network Gear 'CI class defined by S3 (Jun 20)'", "Dep #1 still New;  unlikely by Jun 20 (we are in S3 now)."),
    ("Service Mapping '15 maps by S3'", "Only 4 maps active;  ~36 tasks stranded in carryover."),
    ("VMware (P0, critical) parked entirely in S6", "S6 is the IP iteration AND inside the test code freeze (Jul 18 - Aug 15)."),
]

# Risks: #, description, type, status word, status color, action
RISKS = [
    ("1", "Missing / delayed HR / location data (physical computers)", "Program", "TO ESCALATE", "RED",
     "Escalate now - 'before S3' deadline is today. Blocks Dep #7 + computer-location ingestion (WS #2)."),
    ("2", "Undefined NERC-CIP requirements", "Program", "TO ESCALATE", "RED",
     "Escalate now - overdue. Blocks Dep #2 and all PI-2 #4 NERC build work."),
    ("4", "Competing priorities / unplanned Airlift", "Team", "ACCEPTED", "AMBER",
     "S4 + S5 at 100%. Reconsider: deferring NowAssist (P3) is a clean mitigation lever."),
    ("5", "ServiceNow upgrade / code-freeze windows", "Team", "ACCEPTED", "AMBER",
     "Dates are known (dev Jun 27-Jul 18, test Jul 18-Aug 15) = S4-S6. 'Adjust sprint plan' action still outstanding."),
    ("9", "Airlift / Azure - credential, IP, firewall", "Team", "MITIGATED", "GREEN",
     "Marked 'IP-only changes'. Pressure-test vs Ray's 450-server reality before trusting 'Mitigated'."),
]
RISK_FOOT = "Also open: #3 key-member dependency (Owned), #6 Service Mapping PoC, #7 Qualys field mapping, #8 Business Owners (all Mitigated)."

# Solution decisions the team owns
DECISIONS = [
    "Track Server/Computer Ingestion (WS #2) as its own goal - it is contracted and mature, but invisible at the objective level today.",
    "Set the NowAssist (P3) defer trigger now - it is the first thing to drop if S4/S5 tightens.",
    "Hold the line in S4/S5 - no new scope; route every new ask through the PO trade-off.",
    "Escalate Risks #1 and #2 this week - both are overdue and block contracted + added work.",
    "NERC-CIP needs requirements before any build - it is dependency-only today; do not commit build stories yet.",
    "Qualys is blocked on the vendor plugin (Issue 1465952) - monitor; do not burn capacity waiting on it.",
    "Re-plan VMware out of the frozen S6 where possible - sequence Airlift discovery earlier.",
    "Decide Service Mapping carryover (~36 tasks) - re-sprint or de-scope at iteration review.",
    "Clean up off-board work - OCM 1428659 and orphans 1472365 / 1399787 need an objective home or closure.",
]
NEXT = [
    "Confirm the WS #2 goal split on the board and in ADO.",
    "PO trade-off session: NowAssist defer trigger + S4/S5 capacity protection.",
    "Escalate Risks #1 (location data) and #2 (NERC requirements) - this week.",
    "Iteration-review decision on the Service Mapping carryover backlog.",
]

# ---------------------------------------------------------------------------
# THEME  (mirrors build_status_deck.py)
# ---------------------------------------------------------------------------
DARK   = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0x2E, 0x6D, 0xB4)
LIGHT  = RGBColor(0xF4, 0xF6, 0xF9)
TEXT   = RGBColor(0x2D, 0x2D, 0x2D)
MUTED  = RGBColor(0x6B, 0x72, 0x80)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
AMBER  = RGBColor(0xE0, 0x8A, 0x00)
RED    = RGBColor(0xC0, 0x39, 0x2B)
TIER   = [RGBColor(0x1A, 0x36, 0x5D), RGBColor(0x2B, 0x6C, 0xB0),
          RGBColor(0x2F, 0x80, 0xC9), RGBColor(0x4A, 0x9C, 0xE0)]
RAG = {"GREEN": GREEN, "AMBER": AMBER, "RED": RED}

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
FONT = "Segoe UI"

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def _set(run, size, color, bold=False, italic=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = FONT


def rect(slide, x, y, w, h, color, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def textbox(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, spec in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        text, size, color, bold, italic = spec
        r = p.add_run(); r.text = text
        _set(r, size, color, bold, italic)
    return tb


def header(slide, title, subtitle=None):
    rect(slide, 0, 0, EMU_W, Inches(1.05), DARK)
    rect(slide, 0, Inches(1.05), EMU_W, Inches(0.06), ACCENT)
    textbox(slide, Inches(0.55), Inches(0.18), Inches(12.2), Inches(0.55),
            [(title, 26, WHITE, True, False)], MSO_ANCHOR.MIDDLE)
    if subtitle:
        textbox(slide, Inches(0.57), Inches(0.66), Inches(12.2), Inches(0.32),
                [(subtitle, 12, RGBColor(0xC9, 0xD6, 0xE6), False, False)], MSO_ANCHOR.MIDDLE)
    textbox(slide, Inches(0.55), Inches(7.08), Inches(10), Inches(0.32),
            [(FOOTLABEL, 9, MUTED, False, False)])


def chip(slide, x, y, label, color, w=Inches(1.0), h=Inches(0.30), size=10):
    c = rect(slide, x, y, w, h, color)
    tf = c.text_frame; tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    _set(r, size, WHITE, True)
    return c


# ---- Slide 1: Title -------------------------------------------------------
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, EMU_W, EMU_H, DARK)
rect(s, 0, Inches(4.55), EMU_W, Inches(0.08), ACCENT)
textbox(s, Inches(0.8), Inches(1.9), Inches(11.7), Inches(1.1),
        [(TITLE, 44, WHITE, True, False)])
textbox(s, Inches(0.82), Inches(3.0), Inches(11.7), Inches(0.6),
        [(SUBTITLE, 22, RGBColor(0xC9, 0xD6, 0xE6), False, False)])
textbox(s, Inches(0.82), Inches(4.75), Inches(11.7), Inches(1.5), [
    (f"{PI}  ·  {ITERATION}", 16, WHITE, True, False),
    (PI_WINDOW, 14, RGBColor(0xC9, 0xD6, 0xE6), False, False),
    (AUTHOR, 12, RGBColor(0x9F, 0xB2, 0xCC), False, False),
])

# ---- Slide 2: Purpose -----------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "Why This Deck", "Clarity on what we committed to, what we added, and where to focus")
textbox(s, Inches(0.55), Inches(1.4), Inches(12.2), Inches(0.9),
        [(PURPOSE, 15, TEXT, False, False)])
flow = [("1.  Initial scope", "What the engagement was contracted to deliver - 4 workstreams."),
        ("2.  Added in PI-2", "What PI-2 layered on top - 6 objectives, 3 of them new."),
        ("3.  Direction", "Our ranked goals, live work, capacity, risks, and the calls we own.")]
y = Inches(2.7)
for name, desc in flow:
    rect(s, Inches(0.55), y, Inches(0.12), Inches(0.62), ACCENT)
    textbox(s, Inches(0.8), y, Inches(11.8), Inches(0.65),
            [(name, 15, ACCENT, True, False), (desc, 12.5, TEXT, False, False)])
    y += Inches(0.85)
textbox(s, Inches(0.55), Inches(6.2), Inches(12.2), Inches(0.5),
        [("Audience: the CMDB team doing the work and making solution decisions.", 12, MUTED, False, True)])

# ---- Slide 3: Initial Scope ----------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "Initial Scope - Contracted", "The contracted CMDB scope: 4 workstreams  (WS #1-3 = CO#5  ·  WS #4 = prior CO)")
rows = len(ORIGINAL) + 1
tshape = s.shapes.add_table(rows, 4, Inches(0.55), Inches(1.4), Inches(12.2), Inches(3.8))
t = tshape.table
for j, wdt in enumerate([Inches(1.3), Inches(2.7), Inches(6.9), Inches(1.3)]):
    t.columns[j].width = wdt
for j, h in enumerate(["#", "Workstream", "What it covers", "Effort"]):
    c = t.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = DARK
    p = c.text_frame.paragraphs[0]; r = p.add_run(); r.text = h
    _set(r, 12, WHITE, True)
for i, (ws, name, covers, pts) in enumerate(ORIGINAL, start=1):
    vals = [ws, name, covers, pts]
    for j, v in enumerate(vals):
        c = t.cell(i, j); c.fill.solid()
        c.fill.fore_color.rgb = LIGHT if i % 2 else WHITE
        p = c.text_frame.paragraphs[0]; r = p.add_run(); r.text = v
        if j == 0:
            _set(r, 11.5, ACCENT, True); p.alignment = PP_ALIGN.CENTER
        elif j == 1:
            _set(r, 12, DARK, True)
        elif j == 3:
            _set(r, 11.5, DARK, True); p.alignment = PP_ALIGN.CENTER
        else:
            _set(r, 10.5, TEXT, False)
textbox(s, Inches(0.55), Inches(5.3), Inches(12.2), Inches(0.4),
        [(ORIGINAL_TOTAL, 14, DARK, True, False)])
textbox(s, Inches(0.55), Inches(5.66), Inches(12.2), Inches(0.32),
        [("Dev 4 is walled off for Service Mapping and does not draw on the shared (general) pool.", 11, MUTED, False, True)])
# contract-term reality callout
rect(s, Inches(0.55), Inches(6.05), Inches(12.2), Inches(0.92), RGBColor(0xFB, 0xEE, 0xDB), line=AMBER)
textbox(s, Inches(0.75), Inches(6.12), Inches(11.9), Inches(0.32),
        [("Contract term: CO#5 runs through June 30, 2026", 13, RGBColor(0x7A, 0x4F, 0x10), True, False)])
textbox(s, Inches(0.75), Inches(6.45), Inches(11.9), Inches(0.45),
        [(TERM_NOTE, 10.5, RGBColor(0x5A, 0x4A, 0x2A), False, False)])

# ---- Slide 4: Added in PI-2 ----------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "PI-2 Scope vs Signed Contract (CO#5)", ALIGN_SUMMARY)
rows = len(ALIGNMENT) + 1
tshape = s.shapes.add_table(rows, 4, Inches(0.55), Inches(1.35), Inches(12.2), Inches(3.3))
t = tshape.table
for j, wdt in enumerate([Inches(0.7), Inches(3.6), Inches(1.7), Inches(6.2)]):
    t.columns[j].width = wdt
for j, h in enumerate(["#", "PI-2 Objective", "Alignment", "Note"]):
    c = t.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = DARK
    p = c.text_frame.paragraphs[0]; r = p.add_run(); r.text = h
    _set(r, 12, WHITE, True)
    if j == 2:
        p.alignment = PP_ALIGN.CENTER
for i, (num, name, verdict, vkey, note) in enumerate(ALIGNMENT, start=1):
    vals = [num, name, verdict, note]
    for j, v in enumerate(vals):
        c = t.cell(i, j); c.fill.solid()
        c.fill.fore_color.rgb = LIGHT if i % 2 else WHITE
        p = c.text_frame.paragraphs[0]; r = p.add_run(); r.text = v
        if j == 0:
            _set(r, 11, MUTED, True); p.alignment = PP_ALIGN.CENTER
        elif j == 1:
            _set(r, 11.5, DARK, True)
        elif j == 2:
            c.fill.fore_color.rgb = RAG[vkey]
            _set(r, 10.5, WHITE, True); p.alignment = PP_ALIGN.CENTER
        else:
            _set(r, 10, TEXT, False)
textbox(s, Inches(0.55), Inches(4.85), Inches(12.2), Inches(0.35),
        [("Two original commitments need a home:", 13, DARK, True, False)])
yy = Inches(5.25)
for g in GAPS:
    rect(s, Inches(0.55), yy + Inches(0.04), Inches(0.12), Inches(0.5), AMBER)
    textbox(s, Inches(0.8), yy, Inches(11.9), Inches(0.62), [(g, 11, TEXT, False, False)])
    yy += Inches(0.72)

# ---- Slide 5: Priority Board ---------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "Our Goals, Ranked", "By operational priority  ·  #1 split into added network gear vs. contracted ingestion")
y = Inches(1.32)
rh = Inches(0.72)
for tier, pri, val, name, target, status, scol, origin in BOARD:
    rect(s, Inches(0.55), y, Inches(12.2), rh, LIGHT)
    # priority chip
    rect(s, Inches(0.7), y + Inches(0.12), Inches(0.95), Inches(0.48), TIER[tier])
    textbox(s, Inches(0.7), y + Inches(0.10), Inches(0.95), Inches(0.5),
            [(pri, 14, WHITE, True, False), (val, 8.5, RGBColor(0xDC, 0xEB, 0xF8), False, False)],
            MSO_ANCHOR.MIDDLE, PP_ALIGN.CENTER)
    # name + target
    textbox(s, Inches(1.8), y + Inches(0.06), Inches(7.4), Inches(0.32),
            [(name, 13.5, DARK, True, False)])
    textbox(s, Inches(1.8), y + Inches(0.38), Inches(7.4), Inches(0.3),
            [(target, 9.5, MUTED, False, False)])
    # status chip
    chip(s, Inches(9.35), y + Inches(0.21), status, RAG[scol], Inches(1.45), Inches(0.30), 10)
    # origin
    textbox(s, Inches(10.95), y + Inches(0.21), Inches(1.7), Inches(0.32),
            [(origin, 9.5, DARK, False, False)], MSO_ANCHOR.MIDDLE)
    y += rh + Inches(0.04)

# ---- Slides 6-7: Drill-down ----------------------------------------------
def drill_slide(title, cards, with_unmapped=False):
    s = prs.slides.add_slide(BLANK)
    header(s, title, "What we are actually working - live ADO status (6/17 snapshot)")
    y = Inches(1.3)
    ch = Inches(1.2)
    for tier, pri, name, rollup, lines, scol in cards:
        rect(s, Inches(0.55), y, Inches(12.2), ch, LIGHT)
        rect(s, Inches(0.55), y, Inches(0.12), ch, RAG[scol])
        chip(s, Inches(0.78), y + Inches(0.15), pri, TIER[tier], Inches(0.7), Inches(0.30), 11)
        textbox(s, Inches(1.6), y + Inches(0.12), Inches(8.3), Inches(0.36),
                [(name, 13.5, DARK, True, False)])
        textbox(s, Inches(9.5), y + Inches(0.14), Inches(3.1), Inches(0.34),
                [(rollup, 9.5, MUTED, False, True)], MSO_ANCHOR.TOP, PP_ALIGN.RIGHT)
        yy = y + Inches(0.5)
        for ln in lines:
            textbox(s, Inches(1.6), yy, Inches(11.0), Inches(0.32),
                    [("- " + ln, 10.5, TEXT, False, False)])
            yy += Inches(0.32)
        y += ch + Inches(0.1)
    if with_unmapped:
        rect(s, Inches(0.55), y, Inches(12.2), Inches(0.78), RGBColor(0xFB, 0xEE, 0xDB), line=AMBER)
        textbox(s, Inches(0.75), y + Inches(0.08), Inches(12.0), Inches(0.32),
                [(UNMAPPED[0], 12, RGBColor(0x7A, 0x4F, 0x10), True, False)])
        textbox(s, Inches(0.75), y + Inches(0.42), Inches(12.0), Inches(0.32),
                [(UNMAPPED[1], 10, RGBColor(0x5A, 0x4A, 0x2A), False, False)])
    return s

drill_slide("Live Work - P0 / P1", DRILL_P0P1)
drill_slide("Live Work - P2 / P3", DRILL_P2P3, with_unmapped=True)

# ---- Slide 8: Capacity ----------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "Capacity Reality", "The added scope fits on paper - but the back half has no margin")
y = Inches(1.4)
for label, val, note in CAP_ROWS:
    rect(s, Inches(0.55), y, Inches(12.2), Inches(0.6), LIGHT)
    textbox(s, Inches(0.75), y, Inches(4.6), Inches(0.6),
            [(label, 12.5, DARK, True, False)], MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(5.4), y, Inches(2.0), Inches(0.6),
            [(val, 15, ACCENT, True, False)], MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(7.5), y, Inches(5.2), Inches(0.6),
            [(note, 10.5, TEXT, False, False)], MSO_ANCHOR.MIDDLE)
    y += Inches(0.68)
# sprint strip
y += Inches(0.05)
cw = Inches(2.0)
x = Inches(0.55)
for sp, win, ua, st, scol in CAP_SPRINTS:
    rect(s, x, y, cw - Inches(0.1), Inches(1.4), WHITE, line=RGBColor(0xD0, 0xD7, 0xDE))
    rect(s, x, y, cw - Inches(0.1), Inches(0.42), RAG[scol])
    textbox(s, x, y, cw - Inches(0.1), Inches(0.42),
            [(sp, 14, WHITE, True, False)], MSO_ANCHOR.MIDDLE, PP_ALIGN.CENTER)
    textbox(s, x, y + Inches(0.5), cw - Inches(0.1), Inches(0.85),
            [(win, 9, MUTED, False, False), (ua, 13, DARK, True, False), (st, 10, RAG[scol], True, False)],
            MSO_ANCHOR.TOP, PP_ALIGN.CENTER)
    x += cw
textbox(s, Inches(0.55), y + Inches(1.5), Inches(12.2), Inches(0.7),
        [(CAP_NOTE, 11, TEXT, False, False)])

# ---- Slide 9: Roadmap -----------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "Roadmap - Plan vs. Reality", "The plan is sound; several near-term gates are slipping")
textbox(s, Inches(0.55), Inches(1.32), Inches(12.2), Inches(0.7),
        [(ROADMAP_PLAN, 12, TEXT, False, False)])
textbox(s, Inches(0.55), Inches(2.12), Inches(12.2), Inches(0.35),
        [("Where plan and live status diverge:", 13, DARK, True, False)])
y = Inches(2.55)
for plan, actual in SLIPPAGE:
    rect(s, Inches(0.55), y, Inches(12.2), Inches(0.92), LIGHT)
    rect(s, Inches(0.55), y, Inches(0.12), Inches(0.92), AMBER)
    textbox(s, Inches(0.8), y + Inches(0.1), Inches(11.9), Inches(0.34),
            [("Plan:  " + plan, 12, DARK, True, False)])
    textbox(s, Inches(0.8), y + Inches(0.46), Inches(11.9), Inches(0.4),
            [("Reality:  " + actual, 11, RED, False, False)])
    y += Inches(1.02)

# ---- Slide 10: Risks ------------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "Risks & Escalations", "9 risks total  ·  only 2 need leadership action - and both are overdue")
rows = len(RISKS) + 1
tshape = s.shapes.add_table(rows, 4, Inches(0.55), Inches(1.4), Inches(12.2), Inches(4.0))
t = tshape.table
for j, wdt in enumerate([Inches(3.7), Inches(1.0), Inches(1.8), Inches(5.7)]):
    t.columns[j].width = wdt
for j, h in enumerate(["Risk", "Type", "Status", "Action"]):
    c = t.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = DARK
    p = c.text_frame.paragraphs[0]; r = p.add_run(); r.text = h
    _set(r, 12, WHITE, True)
    if j == 2:
        p.alignment = PP_ALIGN.CENTER
for i, (num, desc, typ, status, scol, action) in enumerate(RISKS, start=1):
    vals = [f"#{num}  {desc}", typ, status, action]
    for j, v in enumerate(vals):
        c = t.cell(i, j); c.fill.solid()
        c.fill.fore_color.rgb = LIGHT if i % 2 else WHITE
        p = c.text_frame.paragraphs[0]; r = p.add_run(); r.text = v
        if j == 0:
            _set(r, 10.5, DARK, True)
        elif j == 1:
            _set(r, 10, MUTED, False); p.alignment = PP_ALIGN.CENTER
        elif j == 2:
            c.fill.fore_color.rgb = RAG[scol]
            _set(r, 9.5, WHITE, True); p.alignment = PP_ALIGN.CENTER
        else:
            _set(r, 10, TEXT, False)
textbox(s, Inches(0.55), Inches(5.7), Inches(12.2), Inches(0.5),
        [(RISK_FOOT, 10.5, MUTED, False, True)])

# ---- Slide 11: Solution Decisions ----------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "Solution Decisions We Own", "The calls the team and PO need to make - this is the direction")
col_w = Inches(6.05)
xs = [Inches(0.55), Inches(6.75)]
half = (len(DECISIONS) + 1) // 2
for idx, d in enumerate(DECISIONS):
    col = 0 if idx < half else 1
    row = idx if idx < half else idx - half
    x = xs[col]
    yy = Inches(1.45) + row * Inches(1.02)
    rect(s, x, yy, Inches(0.32), Inches(0.32), ACCENT)
    textbox(s, x, yy - Inches(0.03), Inches(0.32), Inches(0.38),
            [(str(idx + 1), 12, WHITE, True, False)], MSO_ANCHOR.MIDDLE, PP_ALIGN.CENTER)
    textbox(s, x + Inches(0.42), yy - Inches(0.05), col_w - Inches(0.5), Inches(1.0),
            [(d, 11, TEXT, False, False)])

# ---- Slide 12: Focus & Next ----------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "Focus & Next Steps", "Protect committed goals; make the open calls")
rect(s, Inches(0.55), Inches(1.4), Inches(12.2), Inches(1.0), RGBColor(0xFB, 0xEE, 0xDB), line=AMBER)
textbox(s, Inches(0.78), Inches(1.5), Inches(11.9), Inches(0.34),
        [("Focus discipline", 14, RGBColor(0x7A, 0x4F, 0x10), True, False)])
textbox(s, Inches(0.78), Inches(1.86), Inches(11.9), Inches(0.5),
        [("S4-S5 at 100% capacity + freezes Jun 27 - Aug 15. Protect P0/P1; NowAssist (P3) is the first defer candidate.",
          11.5, RGBColor(0x5A, 0x4A, 0x2A), False, False)])
textbox(s, Inches(0.55), Inches(2.7), Inches(12.2), Inches(0.35),
        [("Immediate next steps:", 14, DARK, True, False)])
y = Inches(3.2)
for n in NEXT:
    rect(s, Inches(0.55), y + Inches(0.05), Inches(0.12), Inches(0.5), ACCENT)
    textbox(s, Inches(0.85), y, Inches(11.8), Inches(0.7), [(n, 13.5, TEXT, False, False)])
    y += Inches(0.8)
textbox(s, Inches(0.55), Inches(6.7), Inches(12.2), Inches(0.4),
        [("Full traceability: PI-2/pi2-objectives-features-stories.md", 11, ACCENT, False, False)])


def save_safe(path):
    try:
        prs.save(path); print("Saved", path)
    except PermissionError:
        base, ext = os.path.splitext(path)
        alt = base + "-new" + ext
        prs.save(alt); print("Target locked - saved to", alt)


out = os.path.join(os.path.dirname(os.path.abspath(__file__)), EDITION + "-scope-focus.pptx")
save_safe(out)
