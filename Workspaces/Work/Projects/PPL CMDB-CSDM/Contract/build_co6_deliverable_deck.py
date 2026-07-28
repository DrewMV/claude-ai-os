"""
CO6 Deliverables deck - built to MATCH THE SOW deliverables table.

Source of truth: the deliverables table in the AUTHORITATIVE CO6 v3
"AP00105-6 to ServiceNow Release 2 and 3 Deployment v3.docx"
(# | Workstream | Description | AC # | Acceptance Criteria | Due Date).
Each card reproduces those SOW fields (light cleanup of grammar/encoding only).

HYBRID layout - do not waste space:
  * information-heavy deliverables get a full slide
  * lighter deliverables are paired two-to-a-slide (stacked, with a divider)

Emits TWO versions for comparison:
  A "with extras"  -> Title, Index, detail slides, Timeline, Where-to-Focus
  B "pure SOW"     -> Title, Index, detail slides

NO contract-commercial data (fees live in separate SOW tables, not touched).

Single source of truth. Edit CONTENT, then run:
    python build_co6_deliverable_deck.py
Outputs:
    CO6-Deliverables-<EDITION>-A-with-extras.pptx
    CO6-Deliverables-<EDITION>-B-pure-sow.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# CONTENT  (authoritative CO6 v3)
# ---------------------------------------------------------------------------
EDITION   = "2026-07-12"
TITLE     = "CO6 Deliverables - Team Focus"
SUBTITLE  = "Per-SOW deliverables: description, acceptance criteria & target dates"
AUTHOR    = "Manuel Vazquez  -  Scrum Master, CMDB-CSDM"
FOOTLABEL = "PPL CMDB-CSDM  |  CO6 Deliverables (per SOW v3)  |  Team Focus  -  Jul 12, 2026"
PLANNING_NOTE = "Sourced from the authoritative CO6 v3.  Dates align to CO6 gates (PI-3 ends Oct 27; final gates Oct 30)."

# Deliverables index (date-ordered). (name, due, pi, focus, is_support)
DELIVERABLES = [
    ("Network Gear Discovery", "Aug 31 -> Oct 30", "PI-3",
     "Automated network-device discovery into the CMDB; 90% coverage, business-owner validated", False),
    ("Service Mapping", "Aug 31 -> Oct 30", "PI-3",
     "15 Silver-tier + Gold-tier maps; >=75% business-service inventory; app->infra and service->app layers", False),
    ("Legacy Platform Rationalization", "Aug 31 -> Oct 30", "PI-3",
     "Gap analysis of iTeam / DISCO / Cherwell / AIM; migration plan; October execution", False),
    ("CMDB Governance", "Sep 30 -> Oct 30", "PI-3",
     "Data certification (process + PROD + dashboards + training) for Computers, Servers, Databases, Network Devices", False),
    ("Qualys Integration", "Sep 30", "PI-3",
     "One-way ServiceNow -> Qualys sync of owner, support group & SOX flag, live in PROD", False),
    ("NERC-CIP Platform Strategy", "Sep 30 -> Oct 30", "PI-3",
     "Dual-instance evaluation + executive strategy package", False),
    ("ATF Strategy", "Oct 30", "PI-3",
     "Plan for Automated Test Framework rollout across in-prod ServiceNow capabilities", False),
    ("ITSM Product Management", "Monthly", "PI-2/3",
     "Dedicated ITSM Product Owner - backlog prioritization, ceremonies, governance", True),
    ("Platform Support", "Monthly", "PI-2/3",
     "2 BAU / DevOps + 4 major-enhancement developers; 40 hrs of stories per week", True),
]
DELIVERABLES_FOOT = ("At-a-glance summary (date-ordered).  Per-deliverable SOW detail follows.  "
                     "Bottom two are standing PO / BAU lanes.  " + PLANNING_NOTE)

# Per-deliverable detail. key, num, name, due_summary, desc, acs=[(bullets[], date)]
SOW = [
    dict(key="network", num="1", name="Network Gear Discovery", due="Aug 31 - Oct 30, 2026",
         desc="Automated, real-time discovery of network devices to populate and maintain the CMDB. Eliminates manual data collection, reduces configuration drift, and provides the foundation for incident impact analysis - improving response times and reducing outage risk across PPL's network.",
         acs=[
             ([
                 "Network device types (routers, switches, firewalls, load balancers, WAPs, network controllers) actively discovered using validated credentials - no failed authentications on target devices",
                 "Discovery schedules running on defined intervals; CMDB records auto-updated without manual intervention within each cycle",
             ], "Aug 31, 2026"),
             (["All discoverable mandatory CMDB attributes for network devices (per CMDB governance) fully populated via discovery - no mandatory fields blank"], "Sep 30, 2026"),
             (["90%+ of expected network-device inventory represented in the CMDB, validated against authoritative data from business owners"], "Oct 30, 2026"),
         ]),
    dict(key="mapping", num="2", name="Service Mapping", due="Aug 31 - Oct 30, 2026 (phased)",
         desc="End-to-end visibility of business services and their infrastructure dependencies in ServiceNow. Accurate maps enable faster incident triage, reliable change-impact assessment, and better outage prioritization - reducing MTTR and unintended disruptions.",
         acs=[
             ([
                 "15 priority (Silver-tier) business-app maps built - full dependencies from business application to infrastructure CIs; owner-validated; ops teams enabled",
                 ">=75% of identified business-service inventory represented in the CMDB, owner-validated  (*assumes business services defined prior to start)",
             ], "Aug 31, 2026"),
             ([
                 "15 Silver-tier maps (business application -> infrastructure): additional validation pass",
                 "Gold-tier maps (business service -> business application) built, owner-validated, teams enabled",
             ], "Sep 30, 2026"),
             ([
                 "Remaining Contractor-managed Silver-tier maps (business application -> infrastructure) built and validated",
                 "Silver-tier maps (business service -> business application) built, owner-validated, teams enabled",
             ], "Oct 30, 2026"),
         ]),
    dict(key="qualys", num="3", name="Qualys Integration", due="Sep 30, 2026",
         desc="Automates synchronization of key CMDB attributes (asset owners, support groups, SOX flags) from ServiceNow to Qualys, eliminating manual reconciliation - so vulnerability scan data stays tied to the right owners and compliance context.",
         acs=[
             (["One-way ServiceNow -> Qualys integration fully configured, tested, and live in production - owner, support group & SOX flag data auto-synchronized to Qualys on a defined schedule, no manual handoffs  (*assumes the required Qualys plug-in is available)"], "Sep 30, 2026"),
         ]),
    dict(key="governance", num="4", name="CMDB Governance", due="Sep 30 - Oct 30, 2026",
         desc="Establishes the data standards, certification processes, and compliance alignment to keep the CMDB trustworthy - accurate, current configuration data for change, incident, and regulatory decisions.",
         acs=[
             ([
                 "Data certification for Computers - process flow documented, technical build released to PROD, tracking dashboards live, training delivered",
                 "Data certification for Servers - same (process, PROD, dashboards, training)",
             ], "Sep 30, 2026"),
             ([
                 "Data certification for Databases - same (process, PROD, dashboards, training)",
                 "Data certification for Network Devices - same (process, PROD, dashboards, training)",
             ], "Oct 30, 2026"),
         ]),
    dict(key="legacy", num="5", name="Legacy Platform Rationalization", due="Aug 31 - Oct 30, 2026",
         desc="A deliberate, structured transition from legacy platforms (iTeam, DISCO, Cherwell, AIM) into ServiceNow - reducing tool sprawl and maintenance cost and enabling retirement of redundant tooling with a validated roadmap.",
         acs=[
             (["Gap analysis of functionality in iTeam, DISCO, Cherwell, and AIM vs. ServiceNow existing capabilities; existing PPL documentation reviewed and validated"], "Aug 31, 2026"),
             (["Migration plan delivered - prioritized list of functionality to migrate + an actionable roadmap with timelines"], "Sep 30, 2026"),
             (["Execution of the October efforts outlined in the migration plan (per 5.2)"], "Oct 30, 2026"),
         ]),
    dict(key="itsm", num="6", name="ITSM Product Management", due="Monthly, Jul 31 - Oct 30, 2026",
         desc="Dedicated product ownership so PPL's ServiceNow ITSM capabilities are strategically guided, backlog-prioritized, and continuously improved - keeping governance and cross-functional coordination in step with platform growth.",
         acs=[
             (["Product owner activities fully executed for the period - stakeholder engagements, active backlog prioritization in agile ceremonies, cross-functional coordination, and participation in governance / process-improvement forums"], "Monthly: Jul 31 -> Oct 30, 2026"),
         ]),
    dict(key="atf", num="7", name="ATF Strategy", due="Oct 30, 2026",
         desc="A scalable, structured approach for adopting the Automated Test Framework (ATF) across PPL's ServiceNow platform - reducing regression risk and improving deployment confidence as capabilities expand.",
         acs=[
             (["Published ATF implementation plan covering all production ServiceNow capabilities in use as of Jul 31 - defined criteria for which capabilities require ATF coverage, a sequenced rollout timeline, and a reusable implementation approach applied consistently across capabilities"], "Oct 30, 2026"),
         ]),
    dict(key="platform", num="8", name="Platform Support", due="Monthly, Jul 31 - Oct 30, 2026",
         desc="A dedicated team of 2 developers for BAU support and DevOps pipeline management (ServiceNow ITSM & CMDB), plus 4 developers for major enhancements (primarily ITSM, >40 hrs effort). Enhancement priority set by product ownership, scoped in agile ceremonies.",
         acs=[
             (["Time- and effort-tracking mechanism operational - all team members logging 40 hours of completed stories per week (holiday/PTO-adjusted); sprint allocation reports available to stakeholders"], "Jul 31, 2026"),
             (["Each team member delivers 40 hours of completed stories per sprint, with ongoing tracking confirming consistent allocation to BAU scope throughout the period"], "Monthly: Aug 31 -> Oct 30, 2026"),
         ]),
    dict(key="nerccip", num="9", name="NERC-CIP ServiceNow Platform Strategy", due="Sep 30 - Oct 30, 2026",
         desc="A leadership-ready strategy package evaluating a separate ServiceNow environment for NERC-CIP assets - enabling PPL to decide the go-forward solution for managing NERC-CIP assets in ServiceNow.",
         acs=[
             (["Evaluation of the dual-instance NERC-CIP platform strategy completed - documented high-level considerations, pros and cons, risks, and dependencies"], "Sep 30, 2026"),
             (["Executive-level NERC-CIP Platform Strategy package (e.g. PowerPoint / Word) describing go-forward solution options, tradeoffs, and high-level implementation steps"], "Oct 30, 2026"),
         ]),
]

# Detail layout: each page = 1 key (full slide) or 2 keys (paired slide).
LAYOUT = [
    ["network"],                 # full - staged discovery
    ["mapping"],                 # full - 6 ACs across 3 gates
    ["governance"],              # full - 4 CI classes
    ["legacy"],                  # full - analysis -> plan -> exec
    ["nerccip"],                 # full - leadership strategy
    ["qualys", "atf"],           # pair - single-AC items
    ["itsm", "platform"],        # pair - standing lanes
]

# Timeline milestones: (when, what's due, lands in)
MILESTONES = [
    ("Jul 31", "ITSM Product Management & Platform Support monthly lanes begin; Platform time/effort tracking operational", "PI-2 tail"),
    ("Aug 31", "Network Gear (creds + schedules)  ·  Service Mapping (15 Silver + >=75% inventory)  ·  Legacy (4-platform gap analysis)", "PI-3"),
    ("Sep 30", "Qualys live in PROD  ·  Governance data cert (Computers, Servers)  ·  Network attrs  ·  Service Mapping (Gold + Silver pass)  ·  NERC-CIP evaluation  ·  Legacy (migration plan)", "PI-3"),
    ("Oct 27", "PI-3 ends", "-"),
    ("Oct 30", "Governance data cert (Databases, Network)  ·  Network 90%  ·  Service Mapping (remaining Silver + service->app)  ·  ATF plan  ·  NERC-CIP exec package  ·  Legacy (Oct execution)", "IP / post-PI-3"),
]

# Where to focus (close)
TAKEAWAYS = [
    ("Two lanes start Jul 31.", "ITSM Product Management and Platform Support (2 BAU + 4 enhancement devs) begin monthly - stand up tracking and capacity first."),
    ("Aug 31 is the first big gate.", "Network Gear discovery live, 15 Silver-tier service maps + >=75% inventory, and the 4-platform Legacy gap analysis all land Aug 31."),
    ("Sep 30 is the heaviest gate.", "Qualys live in PROD, data certification for Computers & Servers, network attributes, Gold-tier maps, and the NERC-CIP evaluation."),
    ("Qualys is ServiceNow -> Qualys.", "It syncs owner, support group and SOX flag TO Qualys (not vulnerability ingestion) - confirm the required Qualys plug-in is available."),
    ("Service Mapping is the biggest build.", "15 Silver + Gold tiers, >=75% inventory, two mapping layers - app-owner validation gates it; assumes business services are pre-defined."),
    ("Questions?", "Jordan Yung is the CO6 SME; Joe Dames owns governance scope. Flag anything unclear to Manny."),
]

# ---------------------------------------------------------------------------
# THEME
# ---------------------------------------------------------------------------
DARK   = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0x2E, 0x6D, 0xB4)
LIGHT  = RGBColor(0xF4, 0xF6, 0xF9)
TEXT   = RGBColor(0x2D, 0x2D, 0x2D)
MUTED  = RGBColor(0x6B, 0x72, 0x80)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
AMBER  = RGBColor(0xE0, 0x8A, 0x00)
RED    = RGBColor(0xC0, 0x39, 0x2B)
GREY   = RGBColor(0x8A, 0x92, 0x9E)
PALE   = RGBColor(0xEC, 0xEE, 0xF1)
RULE   = RGBColor(0xDD, 0xE3, 0xEA)
PI2C   = RGBColor(0x6E, 0x88, 0xAE)
PI3C   = GREEN

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
FONT = "Segoe UI"


def _set(run, size, color, bold=False, italic=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = FONT


def rect(slide, x, y, w, h, color, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def textbox(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, spec in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        text, size, color, bold, italic = spec
        r = p.add_run(); r.text = text
        _set(r, size, color, bold, italic)
    return tb


def textbox_rich(slide, x, y, w, h, paragraphs, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, runs in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0); p.space_after = Pt(1)
        for (text, size, color, bold, italic) in runs:
            r = p.add_run(); r.text = text
            _set(r, size, color, bold, italic)
    return tb


def header(slide, title, subtitle=None):
    rect(slide, 0, 0, EMU_W, Inches(1.05), DARK)
    rect(slide, 0, Inches(1.05), EMU_W, Inches(0.06), ACCENT)
    textbox(slide, Inches(0.55), Inches(0.18), Inches(12.2), Inches(0.55),
            [(title, 24, WHITE, True, False)], MSO_ANCHOR.MIDDLE)
    if subtitle:
        textbox(slide, Inches(0.57), Inches(0.66), Inches(12.2), Inches(0.32),
                [(subtitle, 12, RGBColor(0xC9, 0xD6, 0xE6), False, False)], MSO_ANCHOR.MIDDLE)
    textbox(slide, Inches(0.55), Inches(7.08), Inches(10), Inches(0.32),
            [(FOOTLABEL, 9, MUTED, False, False)])


def cell(t, i, j, lines, size, color, bold=False, align=PP_ALIGN.LEFT, fill=None, anchor=None):
    c = t.cell(i, j)
    if fill is not None:
        c.fill.solid(); c.fill.fore_color.rgb = fill
    if anchor is not None:
        c.vertical_anchor = anchor
    tf = c.text_frame; tf.word_wrap = True
    if isinstance(lines, str):
        lines = [lines]
    for k, ln in enumerate(lines):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0); p.space_after = Pt(2)
        r = p.add_run(); r.text = ln
        _set(r, size, color, bold)


def ac_table(slide, x, y, w, h, acs, hsize=11, bsize=10):
    t = slide.shapes.add_table(len(acs) + 1, 2, Inches(x), Inches(y), Inches(w), Inches(h)).table
    t.columns[0].width = Inches(w - 3.0)
    t.columns[1].width = Inches(3.0)
    cell(t, 0, 0, "Acceptance criteria", hsize, WHITE, True, PP_ALIGN.LEFT, DARK)
    cell(t, 0, 1, "Targeted due date", hsize, WHITE, True, PP_ALIGN.CENTER, DARK)
    for i, (bullets, date) in enumerate(acs, start=1):
        fill = LIGHT if i % 2 else WHITE
        disp = bullets if len(bullets) == 1 else ["•  " + b for b in bullets]
        cell(t, i, 0, disp, bsize, TEXT, False, PP_ALIGN.LEFT, fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 1, date, bsize, ACCENT, True, PP_ALIGN.CENTER, fill, MSO_ANCHOR.MIDDLE)


# ---------------------------------------------------------------------------
# SLIDE BUILDERS
# ---------------------------------------------------------------------------
def add_title(prs, blank):
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, EMU_W, EMU_H, DARK)
    rect(s, 0, Inches(4.55), EMU_W, Inches(0.08), ACCENT)
    textbox(s, Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.1), [(TITLE, 42, WHITE, True, False)])
    textbox(s, Inches(0.82), Inches(2.9), Inches(11.7), Inches(0.6), [(SUBTITLE, 20, RGBColor(0xC9, 0xD6, 0xE6), False, False)])
    textbox(s, Inches(0.82), Inches(4.75), Inches(11.7), Inches(1.5), [
        ("9 deliverables, straight from the CO6 SOW", 16, WHITE, True, False),
        ("Delivery window: Jul 1 - Oct 30, 2026  (tail of PI-2 + all of PI-3)", 14, RGBColor(0xC9, 0xD6, 0xE6), False, False),
        (AUTHOR, 12, RGBColor(0x9F, 0xB2, 0xCC), False, False),
    ])
    textbox(s, Inches(0.82), Inches(6.5), Inches(11.7), Inches(0.4), [(PLANNING_NOTE, 11, RGBColor(0x7F, 0x8C, 0xA6), False, True)])


def add_index(prs, blank):
    s = prs.slides.add_slide(blank)
    header(s, "Deliverables & Dates", "At-a-glance - ordered by due date (SOW detail follows)")
    rows = len(DELIVERABLES) + 1
    t = s.shapes.add_table(rows, 4, Inches(0.55), Inches(1.4), Inches(12.2), Inches(4.7)).table
    for j, wdt in enumerate([Inches(3.0), Inches(2.0), Inches(0.9), Inches(6.3)]):
        t.columns[j].width = wdt
    for j, h in enumerate(["Deliverable", "Due", "PI", "What we focus on"]):
        cell(t, 0, j, h, 11.5, WHITE, True, PP_ALIGN.LEFT if j in (0, 3) else PP_ALIGN.CENTER, DARK)
    for i, (name, due, pi, focus, sup) in enumerate(DELIVERABLES, start=1):
        fill = PALE if sup else (LIGHT if i % 2 else WHITE)
        cell(t, i, 0, name, 10.5, MUTED if sup else DARK, True, PP_ALIGN.LEFT, fill)
        cell(t, i, 1, due, 10, ACCENT if not sup else MUTED, True, PP_ALIGN.CENTER, fill)
        cell(t, i, 2, pi, 10, TEXT, True, PP_ALIGN.CENTER, fill)
        cell(t, i, 3, focus, 9.5, TEXT if not sup else MUTED, False, PP_ALIGN.LEFT, fill)
    textbox(s, Inches(0.55), Inches(6.25), Inches(12.2), Inches(0.6), [(DELIVERABLES_FOOT, 10, TEXT, False, False)])


DETAIL_HEAD = "Deliverable Detail"
DETAIL_SUB = "SOW: description, acceptance criteria & target dates"


def card(slide, y0, d):
    """Unified deliverable card - identical formatting on full and paired slides."""
    textbox_rich(slide, Inches(0.55), Inches(y0), Inches(12.2), Inches(0.32),
                 [[("Deliverable " + d["num"] + "  -  " + d["name"], 15, DARK, True, False),
                   ("      " + d["due"], 11.5, ACCENT, True, False)]])
    rect(slide, Inches(0.55), Inches(y0 + 0.34), Inches(12.2), Inches(0.022), ACCENT)
    textbox_rich(slide, Inches(0.55), Inches(y0 + 0.42), Inches(12.2), Inches(0.6),
                 [[("Description   ", 9.5, MUTED, True, False), (d["desc"], 10, TEXT, False, False)]])
    rows = len(d["acs"]) + 1
    ac_table(slide, 0.55, y0 + 1.06, 12.2, 0.4 * rows, d["acs"], hsize=11, bsize=10)


def add_detail_full(prs, blank, d):
    s = prs.slides.add_slide(blank)
    header(s, DETAIL_HEAD, DETAIL_SUB)
    card(s, 1.3, d)


def add_detail_pair(prs, blank, dt, db):
    s = prs.slides.add_slide(blank)
    header(s, DETAIL_HEAD, DETAIL_SUB)
    card(s, 1.25, dt)
    rect(s, Inches(0.55), Inches(3.93), Inches(12.2), Inches(0.016), RULE)
    card(s, 4.0, db)


def add_timeline(prs, blank):
    s = prs.slides.add_slide(blank)
    header(s, "Timeline - Where It Lands", "Tail of PI-2 + all of PI-3  ·  staged monthly acceptance")
    axl, axr, m0, m1 = 2.0, 12.6, 6, 11

    def px(m):
        return Inches(axl + (m - m0) / (m1 - m0) * (axr - axl))

    gt, gh = 1.3, 1.35
    for m, lbl in [(6, "Jun"), (7, "Jul"), (8, "Aug"), (9, "Sep"), (10, "Oct"), (11, "Nov")]:
        rect(s, px(m), Inches(gt), Pt(1), Inches(gh), RULE)
        textbox(s, px(m) - Inches(0.3), Inches(gt + gh + 0.02), Inches(0.6), Inches(0.25), [(lbl, 9, MUTED, False, False)], align=PP_ALIGN.CENTER)
    textbox(s, Inches(0.55), Inches(gt), Inches(1.4), Inches(gh), [("Windows", 9, MUTED, True, False)], MSO_ANCHOR.MIDDLE)
    bars = [
        ("Delivery window   Jul 1 - Oct 30", 7.0, 10.935, ACCENT, WHITE),
        ("PI-2  (-> Aug 4)", 6.0, 8.097, PI2C, WHITE),
        ("PI-3  Aug 5 - Oct 27", 8.129, 10.839, PI3C, WHITE),
    ]
    bh, gap, yy = 0.32, 0.085, gt + 0.06
    for lbl, ms, me, col, tc in bars:
        x1, x2 = px(ms), px(me)
        rect(s, x1, Inches(yy), x2 - x1, Inches(bh), col)
        textbox(s, x1 + Inches(0.08), Inches(yy), (x2 - x1) - Inches(0.12), Inches(bh), [(lbl, 9.5, tc, True, False)], MSO_ANCHOR.MIDDLE)
        yy += bh + gap
    rows = len(MILESTONES) + 1
    t = s.shapes.add_table(rows, 3, Inches(0.55), Inches(3.35), Inches(12.2), Inches(3.1)).table
    for j, wdt in enumerate([Inches(1.5), Inches(8.6), Inches(2.1)]):
        t.columns[j].width = wdt
    for j, h in enumerate(["When", "What's due", "Lands in"]):
        cell(t, 0, j, h, 11, WHITE, True, PP_ALIGN.LEFT if j == 1 else PP_ALIGN.CENTER, DARK)
    for i, (when, what, lands) in enumerate(MILESTONES, start=1):
        end = when == "Oct 27"
        fill = RGBColor(0xFB, 0xEE, 0xDB) if end else (LIGHT if i % 2 else WHITE)
        cell(t, i, 0, when, 10.5, RED if end else DARK, True, PP_ALIGN.CENTER, fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 1, what, 9, RGBColor(0x5A, 0x4A, 0x2A) if end else TEXT, end, PP_ALIGN.LEFT, fill, MSO_ANCHOR.MIDDLE)
        cell(t, i, 2, lands, 9.5, ACCENT, True, PP_ALIGN.CENTER, fill, MSO_ANCHOR.MIDDLE)


def add_focus(prs, blank):
    s = prs.slides.add_slide(blank)
    header(s, "Where to Focus", "Priorities for Development + supporting roles")
    yy = 1.4
    for head, body in TAKEAWAYS:
        rect(s, Inches(0.55), Inches(yy + 0.04), Inches(0.12), Inches(0.62), ACCENT)
        textbox(s, Inches(0.82), Inches(yy), Inches(11.9), Inches(0.32), [(head, 13.5, DARK, True, False)])
        textbox(s, Inches(0.82), Inches(yy + 0.34), Inches(11.9), Inches(0.5), [(body, 11, TEXT, False, False)])
        yy += 0.92


def build_deck(include_extras):
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    blank = prs.slide_layouts[6]
    bykey = {d["key"]: d for d in SOW}
    add_title(prs, blank)
    add_index(prs, blank)
    for page in LAYOUT:
        if len(page) == 1:
            add_detail_full(prs, blank, bykey[page[0]])
        else:
            add_detail_pair(prs, blank, bykey[page[0]], bykey[page[1]])
    if include_extras:
        add_timeline(prs, blank)
        add_focus(prs, blank)
    return prs


def save_safe(prs, path):
    try:
        prs.save(path); print("Saved", os.path.basename(path), "-", len(prs.slides._sldIdLst), "slides")
    except PermissionError:
        base, ext = os.path.splitext(path)
        alt = base + "-new" + ext
        prs.save(alt); print("Target locked - saved to", os.path.basename(alt))


here = os.path.dirname(os.path.abspath(__file__))
save_safe(build_deck(True),  os.path.join(here, "CO6-Deliverables-" + EDITION + "-A-with-extras.pptx"))
save_safe(build_deck(False), os.path.join(here, "CO6-Deliverables-" + EDITION + "-B-pure-sow.pptx"))
