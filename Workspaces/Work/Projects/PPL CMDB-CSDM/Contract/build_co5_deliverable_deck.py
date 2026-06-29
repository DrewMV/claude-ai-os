"""
CO5 Deliverable deck generator  (one slide per CO5 deliverable).

Audience: PO / delivery leadership working session to fine-tune what we report
per contractual deliverable. Each deliverable slide maps the CO5 sub-items to
the related PI-2 ADO features + stories, their status, and the CO6 (draft) effect.

Single source of truth for this deck. Edit the CONTENT block, then run:
    python build_co5_deliverable_deck.py
Output:  CO5-Deliverables-<EDITION>.pptx

Theme mirrors build_scope_deck.py / build_status_deck.py so the decks look like a set.
Story status from the 6/19 ADO grid (see ../co5-deliverable-tracking.md).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# CONTENT
# ---------------------------------------------------------------------------
EDITION   = "2026-06-20"
TITLE     = "CO5 Deliverables - Story Validation"
SUBTITLE  = "Change Order #5  ·  per-deliverable feature & story mapping"
AUTHOR    = "Manuel Vazquez  -  Scrum Master, CMDB-CSDM"
FOOTLABEL = "PPL CMDB-CSDM  |  CO5 Deliverables  |  ADO 6/19 snapshot  -  Jun 20, 2026"

CONTRACT_NOTE = ("CO5 (signed 2026-05-26)  ·  term Mar 31 - Jun 30, 2026  ·  "
                 "$533,775 June holdback released on final acceptance.  "
                 "PI-2 runs to Aug 4 - stories in Iter 2.4/2.5/2.6 land after the 6/30 deadline.")

# Overview rows: deliverable, CO5 due, stories mapped, hard gaps, CO6 (draft) date
OVERVIEW = [
    ("1  Governance",                 "Jun 30", "~17", "G1 Databases · G5 CCB",  "Jul 31"),
    ("2  Automated Data Ingestion",   "Jun 30", "~9",  "G2/G3 90% coverage",      "Jul 31 / Sep 30 / Oct 30"),
    ("3  Other Enhancement (Qualys)", "Jun 30", "5",   "G4 BLOCKED",              "Oct 27"),
]
OVERVIEW_NOTE = ("Hard gaps (G1-G5) = a CO5 acceptance item with NO story, or fully blocked.  "
                 "CO6 (draft, unsigned) re-baselines all but G5 - contingent on signature before July 1.")

# Each deliverable: (number, name, due_strip, rows)
# rows = (sub-item, PI-2 feature(s), related stories [ID - status], status_word, status_color)
GREEN_K, AMBER_K, RED_K, GREY_K = "GREEN", "AMBER", "RED", "GREY"

D1_STRIP = ("Due Jun 30 (CO5)   ·   CO6 (draft) re-baselines to Jul 31   ·   "
            "Acceptance: only Data Cert is near acceptance;  Databases (G1) & CCB (G5) are hard gaps")
D1_ROWS = [
    ("1.1 Data Dictionary - Servers", "1354797 Server Class Data & Form",
     "1421790 Gap Analysis (Validation) · 1454371 Extract w/o Support Grp (Active) · 1355167 spike (Active)",
     "In progress", AMBER_K),
    ("1.1 Data Dictionary - Computers", "1354794 Computer Class Data & Form",
     "1313400 Asset Tags (Resolved) · 1399547 Shared Devices (Validation) · 1455858 Virtual Location (Active)",
     "In progress", AMBER_K),
    ("1.1 Data Dictionary - Business Apps", "BA data-population cluster (no single feature)",
     "1475582 Support Grps (Refine Rdy) · 1478286 Classification (Refine Rdy) · 1475584/85 (New, Manuel) · 1474892 (New)",
     "Early", AMBER_K),
    ("1.1 Data Dictionary - Databases  [G1]", "- none -",
     "NO STORY IN ADO",
     "GAP", RED_K),
    ("1.2 Data Certification", "1247179 · 1402979 · 1402958 · 1371672",
     "1402727 Dashboard (UAT signed, PROD 6/23) · 1435307 Pilot Changes (Validation) · 1402962/76/80/84/85 (Defining)",
     "Near acceptance", GREEN_K),
    ("1.3 ESS-02 Alignment", "under 1406668 (Removed)",
     "1420244 ESS-02 Policy & CMDB Alignment (Active - spike, analysis only)",
     "Analysis only", AMBER_K),
    ("1.4 SOX BA Review  (SOX only, not NERC/CIP)", "1354797 / governance",
     "1438967 (Closed/approved) · 1455827 SOX Notify (Refine Rdy, Iter 2.5) · 1399787 Full Access Grp (Validation)",
     "Partial", AMBER_K),
    ("1.5 Monthly CCB + future-PI backlog  [G5]", "1406668 - REMOVED in ADO",
     "1406672/683/687 May/Jun/Jul CCB (New, under Removed feature) · 1339116 (Refine Rdy)",
     "Gap - broken evidence", RED_K),
]

D2_STRIP = ("Due Jun 30 (CO5)   ·   CO6 (draft): precedence/discovery -> Jul 31, 90% coverage -> Sep 30 / Oct 30   ·   "
            "Acceptance: precedence in Validation;  both 90% goals (G2/G3) have no story")
D2_ROWS = [
    ("2.1 Computers - SCCM/Discovery precedence", "1354794 / SCCM Computer Class",
     "1348712/716/717 (Resolved) · 1348715 Last-Seen (Validation)", "On track", GREEN_K),
    ("2.1 Computers - bulk Life Cycle Stage/Status", "Computer Class 1354794",
     "1402790 Retired Computers Lifecycle (Closed)", "Done", GREEN_K),
    ("2.1 Computers - 90% coverage validation  [G2]", "- none -",
     "NO MEASUREMENT STORY", "GAP", RED_K),
    ("2.2 Servers - SCCM/Discovery precedence", "1356826 SCCM Server Class Precedence",
     "1403759 / 1403760 / 1403762 / 1403763 (all Validation, approval gate cleared)", "Validation", AMBER_K),
    ("2.2 Servers - credentials (enables discovery)", "1356646 / 1354797",
     "1444864 Fix Creds Svr/DB/Net (Validation; 6 child tasks still Active) · 1459721 SNMP/MID (Validation)",
     "In progress", AMBER_K),
    ("2.2 Servers - 90% non-NERC-CIP coverage  [G3]", "- none -",
     "NO MEASUREMENT STORY", "GAP", RED_K),
    ("2.3 Databases - enhanced Discovery (MS-SQL/Oracle)", "1356646 / credential work",
     "Oracle DB cred tasks under 1444864 (in progress, task-level only)", "Task-level only", AMBER_K),
]
D2_FOOT = ("SOX indicators for Servers & Databases are manually maintained - explicitly EXCLUDED from "
           "automated ingestion per CO5.  CO5-documented RISK: no scalable credential-distribution solution for target CIs.")

D3_STRIP = ("Due Jun 30 (CO5)   ·   CO6 (draft) ESCALATES evaluate -> full one-way PROD integration, due Oct 27   ·   "
            "Acceptance: BLOCKED (G4)")
D3_ROWS = [
    ("3  Evaluate ONE integration (Tanium OR Qualys)\n     -> team chose Qualys",
     "Integration Qualys & ServiceNow\n(CMDB Data Read-Only) (PI2)",
     "1428703 Install Plugin P1 (BLOCKED) · 1428704 Configure Integration P2 (BLOCKED) · "
     "1234585 Data Scope spike (Ready, no parent) · 1465952 Plugin Replacement (Closed - may unblock) · "
     "1383519 Qualys Dev/Support dep (Active)",
     "BLOCKED", RED_K),
]
D3_DETAIL = [
    "Blocker: Qualys replaced the plugin version; previous uninstalled, new one requested - pending vendor approval (Issue 1465952).",
    "1465952 now shows Closed in the 6/19 grid - MAY unblock 1428703/1428704, but story states were NOT auto-flipped. Confirm with Rich Santillo / Stan.",
    "CO5 only requires EVALUATE + document requirements for ONE integration - a lighter bar than full implementation.",
    "Tanium is the contractual alternative, but no Tanium stories exist - a pivot this late is unrealistic for 6/30.",
    "CO6 (draft) escalates this to a full one-way Qualys->ServiceNow PROD integration, due Oct 27 (PI-3).",
]

# ---------------------------------------------------------------------------
# THEME  (mirrors build_scope_deck.py)
# ---------------------------------------------------------------------------
DARK   = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0x2E, 0x6D, 0xB4)
LIGHT  = RGBColor(0xF4, 0xF6, 0xF9)
TEXT   = RGBColor(0x2D, 0x2D, 0x2D)
MUTED  = RGBColor(0x6B, 0x72, 0x80)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
AMBER  = RGBColor(0xE0, 0x8A, 0x00)
RED    = RGBColor(0xC0, 0x39, 0x2B)
GREY   = RGBColor(0x8A, 0x92, 0x9E)
RAG = {"GREEN": GREEN, "AMBER": AMBER, "RED": RED, "GREY": GREY}

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
FONT = "Segoe UI"

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def _set(run, size, color, bold=False, italic=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = FONT


def rect(slide, x, y, w, h, color, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def textbox(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, spec in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        text, size, color, bold, italic = spec
        r = p.add_run(); r.text = text
        _set(r, size, color, bold, italic)
    return tb


def header(slide, title, subtitle=None):
    rect(slide, 0, 0, EMU_W, Inches(1.05), DARK)
    rect(slide, 0, Inches(1.05), EMU_W, Inches(0.06), ACCENT)
    textbox(slide, Inches(0.55), Inches(0.18), Inches(12.2), Inches(0.55),
            [(title, 26, WHITE, True, False)], MSO_ANCHOR.MIDDLE)
    if subtitle:
        textbox(slide, Inches(0.57), Inches(0.66), Inches(12.2), Inches(0.32),
                [(subtitle, 12, RGBColor(0xC9, 0xD6, 0xE6), False, False)], MSO_ANCHOR.MIDDLE)
    textbox(slide, Inches(0.55), Inches(7.08), Inches(10), Inches(0.32),
            [(FOOTLABEL, 9, MUTED, False, False)])


def cell(t, i, j, lines, size, color, bold=False, align=PP_ALIGN.LEFT, fill=None):
    c = t.cell(i, j)
    if fill is not None:
        c.fill.solid(); c.fill.fore_color.rgb = fill
    tf = c.text_frame; tf.word_wrap = True
    if isinstance(lines, str):
        lines = [lines]
    for k, ln in enumerate(lines):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        _set(r, size, color, bold)


# ---- Slide 1: Title -------------------------------------------------------
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, EMU_W, EMU_H, DARK)
rect(s, 0, Inches(4.55), EMU_W, Inches(0.08), ACCENT)
textbox(s, Inches(0.8), Inches(1.9), Inches(11.7), Inches(1.1),
        [(TITLE, 42, WHITE, True, False)])
textbox(s, Inches(0.82), Inches(3.0), Inches(11.7), Inches(0.6),
        [(SUBTITLE, 20, RGBColor(0xC9, 0xD6, 0xE6), False, False)])
textbox(s, Inches(0.82), Inches(4.75), Inches(11.7), Inches(1.5), [
    ("3 deliverables  ·  one slide each  ·  PI-2 features & stories", 16, WHITE, True, False),
    ("CO5  ·  Mar 31 - Jun 30, 2026", 14, RGBColor(0xC9, 0xD6, 0xE6), False, False),
    (AUTHOR, 12, RGBColor(0x9F, 0xB2, 0xCC), False, False),
])

# ---- Slide 2: Overview ----------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "CO5 Deliverables - Overview", "3 contractual deliverables, all due Jun 30  ·  CO6 (draft) is the re-baselining vehicle")
rows = len(OVERVIEW) + 1
tshape = s.shapes.add_table(rows, 5, Inches(0.55), Inches(1.45), Inches(12.2), Inches(2.4))
t = tshape.table
for j, wdt in enumerate([Inches(3.7), Inches(1.4), Inches(1.8), Inches(2.9), Inches(2.4)]):
    t.columns[j].width = wdt
for j, h in enumerate(["Deliverable", "CO5 Due", "Stories", "Hard gaps", "CO6 (draft) date"]):
    cell(t, 0, j, h, 12, WHITE, True, PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER, DARK)
for i, (name, due, stories, gaps, co6) in enumerate(OVERVIEW, start=1):
    fill = LIGHT if i % 2 else WHITE
    cell(t, i, 0, name, 12, DARK, True, PP_ALIGN.LEFT, fill)
    cell(t, i, 1, due, 11.5, RED, True, PP_ALIGN.CENTER, fill)
    cell(t, i, 2, stories, 11.5, TEXT, False, PP_ALIGN.CENTER, fill)
    cell(t, i, 3, gaps, 10.5, RED, True, PP_ALIGN.CENTER, fill)
    cell(t, i, 4, co6, 11, GREEN, True, PP_ALIGN.CENTER, fill)
textbox(s, Inches(0.55), Inches(4.25), Inches(12.2), Inches(0.6),
        [(OVERVIEW_NOTE, 11.5, TEXT, False, False)])
# holdback / timeline callout
rect(s, Inches(0.55), Inches(5.4), Inches(12.2), Inches(1.25), RGBColor(0xFB, 0xEE, 0xDB), line=AMBER)
textbox(s, Inches(0.75), Inches(5.5), Inches(11.9), Inches(0.32),
        [("Governing constraint", 13, RGBColor(0x7A, 0x4F, 0x10), True, False)])
textbox(s, Inches(0.75), Inches(5.85), Inches(11.9), Inches(0.7),
        [(CONTRACT_NOTE, 11, RGBColor(0x5A, 0x4A, 0x2A), False, False)])

# ---- Deliverable slides ---------------------------------------------------
def deliverable_slide(title, strip, rows_data, table_h, foot=None, detail=None):
    s = prs.slides.add_slide(BLANK)
    header(s, title)
    # info strip
    rect(s, Inches(0.55), Inches(1.22), Inches(12.2), Inches(0.5), RGBColor(0xE9, 0xF0, 0xF8), line=ACCENT)
    textbox(s, Inches(0.72), Inches(1.24), Inches(11.9), Inches(0.46),
            [(strip, 10.5, DARK, False, False)], MSO_ANCHOR.MIDDLE)
    # table
    rows = len(rows_data) + 1
    tshape = s.shapes.add_table(rows, 4, Inches(0.55), Inches(1.85), Inches(12.2), table_h)
    t = tshape.table
    for j, wdt in enumerate([Inches(3.3), Inches(2.7), Inches(4.7), Inches(1.5)]):
        t.columns[j].width = wdt
    for j, h in enumerate(["CO5 sub-item", "PI-2 Feature(s)", "Related stories (ID - status)", "Status"]):
        cell(t, 0, j, h, 11, WHITE, True, PP_ALIGN.LEFT if j < 3 else PP_ALIGN.CENTER, DARK)
    for i, (sub, feat, stories, status, scol) in enumerate(rows_data, start=1):
        fill = LIGHT if i % 2 else WHITE
        cell(t, i, 0, sub.split("\n"), 9.5, DARK, True, PP_ALIGN.LEFT, fill)
        cell(t, i, 1, feat.split("\n"), 9, TEXT, False, PP_ALIGN.LEFT, fill)
        cell(t, i, 2, stories, 9, TEXT, False, PP_ALIGN.LEFT, fill)
        cell(t, i, 3, status, 9.5, WHITE, True, PP_ALIGN.CENTER, RAG[scol])
    yb = Inches(1.85) + table_h + Inches(0.12)
    if foot:
        rect(s, Inches(0.55), yb, Inches(12.2), Inches(0.6), RGBColor(0xFB, 0xEE, 0xDB), line=AMBER)
        textbox(s, Inches(0.72), yb + Inches(0.04), Inches(11.9), Inches(0.54),
                [(foot, 9.5, RGBColor(0x5A, 0x4A, 0x2A), False, False)], MSO_ANCHOR.MIDDLE)
    if detail:
        textbox(s, Inches(0.55), yb, Inches(12.2), Inches(0.3),
                [("Detail:", 12, DARK, True, False)])
        yy = yb + Inches(0.38)
        for d in detail:
            rect(s, Inches(0.55), yy + Inches(0.04), Inches(0.1), Inches(0.36), ACCENT)
            textbox(s, Inches(0.78), yy, Inches(11.9), Inches(0.5), [(d, 10.5, TEXT, False, False)])
            yy += Inches(0.52)
    return s

deliverable_slide("Deliverable 1 - Governance", D1_STRIP, D1_ROWS, Inches(5.0))
deliverable_slide("Deliverable 2 - Automated Data Ingestion", D2_STRIP, D2_ROWS, Inches(4.4), foot=D2_FOOT)
deliverable_slide("Deliverable 3 - Other Enhancement", D3_STRIP, D3_ROWS, Inches(1.5), detail=D3_DETAIL)


# ---------------------------------------------------------------------------
# GOVERNANCE DEEP-DIVE  (section after slide 5 - one slide per CO5 sub-item)
# Slides 1-5 above are unchanged.
# ---------------------------------------------------------------------------
CO6_GREEN_FILL = RGBColor(0xE6, 0xF3, 0xEC)
CO6_GREEN_TEXT = RGBColor(0x1E, 0x5A, 0x38)

# (first-col header, contract text, rows[(label, feature, stories, status, color)], table_h, co6_text, notes[])
GOV = [
    dict(
        title="Governance 1.1 - Data Dictionary / Class Attributes",
        subtitle="Define class attributes for the four in-scope CI classes",
        col0="CI Class",
        contract=("Data Dictionary / Class Attribute definitions for: Servers (Windows/Linux), "
                  "Computers (Physical/Virtual), Business Applications, Databases."),
        rows=[
            ("Servers (Win/Linux)", "1354797 Server Class Data & Form",
             "1421790 Gap Analysis (Validation) · 1454371 Extract Servers w/o Support Grp (Active) · 1355167 Migrate BA to Svc Instance (Active)",
             "In progress", AMBER_K),
            ("Computers (Phys/Virtual)", "1354794 Computer Class Data & Form",
             "1313400 Asset Tags (Resolved) · 1399547 Shared Devices (Validation) · 1455858 Virtual Location (Active) · 1407572 Auto-Populate Owner (Validation)",
             "In progress", AMBER_K),
            ("Business Applications", "BA data-population cluster (no single feature)",
             "1475582 Support Grps (Refine Rdy) · 1478286 Classification (Refine Rdy) · 1475584 Tech Owner Grp (New, Manuel) · 1475585 Approval Grps (New, Manuel) · 1474892 Value Stream (New)",
             "Early", AMBER_K),
            ("Databases  [G1]", "- none -",
             "NO STORY IN ADO - hard gap", "GAP", RED_K),
        ],
        table_h=Inches(2.7),
        co6=("CO6 (draft) 'CMDB Governance' lists Databases explicitly -> re-baselines the full data dictionary "
             "to Jul 31. Closes hard gap G1 (if signed)."),
        notes=[
            "Databases is the only CI class with zero story coverage - the clean contractual gap (G1).",
            "Business App data-population is mostly New / Refinement Ready (Iter 2.4) - lands after the 6/30 deadline.",
        ],
    ),
    dict(
        title="Governance 1.2 - Data Certification",
        subtitle="Strongest-covered governance sub-item; pilot deploying to PROD 6/23",
        col0="CO5 component",
        contract=("Process Definition; Intake for New Data Certifications; End-to-end process flow overview; "
                  "Pilot Release for the Business Application CI Class; Supporting User Documentation / KB Articles."),
        rows=[
            ("Process Definition & framework", "1247179 Pilot Functionality · 1402979 Impl Support",
             "1435307 Data Cert Pilot Changes (Validation)", "On track", GREEN_K),
            ("Intake for new certifications", "1371672 Governance Model (New BA Requests)",
             "1455827 SOX Notify on Ownership Change (Refine Rdy, Iter 2.5)", "Partial", AMBER_K),
            ("E2E flow + pilot execution", "1402958 Impl Planning",
             "1402962 Kickoff · 1402976 Execute Policies · 1402980 Monitor · 1402984 Office Hours · 1402985 Feedback (Defining/New)",
             "In progress", AMBER_K),
            ("Pilot Release - Business App CI Class", "1402979 Impl Support",
             "1402727 Dashboard (UAT signed 6/10, PROD 6/23) · 1402959 Pilot Group (Closed)",
             "Near acceptance", GREEN_K),
            ("Supporting KB / User Docs", "- none -",
             "NO STORY - owner TBD", "Soft gap", AMBER_K),
        ],
        table_h=Inches(3.0),
        co6="CO6 (draft) expands the pilot from Business App only to ALL CI classes -> Jul 31.",
        notes=[
            "Pilot Release (Business App CI class) is the holdback-critical item: PROD deploy 6/23 -> Todd validates 6/24-25 - acceptance evidence just before 6/30.",
            "KB Articles / User Documentation has no story and no owner - confirm before acceptance.",
        ],
    ),
    dict(
        title="Governance 1.3 - ESS-02 Alignment",
        subtitle="Cyber Security & Compliance alignment - analysis only today",
        col0="Area",
        contract="ESS-02 Alignment with Cyber Security and Compliance.",
        rows=[
            ("ESS-02 policy & CMDB alignment", "under 1406668 (Removed)",
             "1420244 ESS-02 Policy & CMDB Alignment (Active - spike, 0 pts, Joe Dames)",
             "Analysis only", AMBER_K),
        ],
        table_h=Inches(0.9),
        co6="CO6 (draft) carries ESS-02 under 'CMDB Governance' -> Jul 31.",
        notes=[
            "Single Active spike, no implementation stories - confirm whether alignment = analysis-only (spike sufficient) or requires a deliverable artifact by deadline.",
            "Open actions: engage Jason Dubreuil (Cyber/Compliance CCB member) on ESS-02; document the CMP.",
            "Parent feature 1406668 is Removed in ADO - the same dead feature as the CCB stories (see 1.5).",
            "'ESS-02' is not expanded in the contract - confirm the definition (likely 'Enterprise Security Standard') with Jason Dubreuil.",
        ],
    ),
    dict(
        title="Governance 1.4 - SOX Business Application Review",
        subtitle="Ensure SOX BAs identified in CMDB + correct update-rights controls",
        col0="Area",
        contract=("Support review of SOX Business Applications - ensure properly identified in CMDB and rules in place "
                  "so the right users can update them. Applies to SOX data ONLY - not NERC/CIP compliance indicators."),
        rows=[
            ("SOX custom-fields governance", "1354797 / governance",
             "1438967 PMDB App Service SOX & DR Customization (Closed/approved; edit rights -> Cyber Risk & DR)",
             "Done", GREEN_K),
            ("SOX ownership-change notification", "1371672",
             "1455827 Notification for SOX Team on Ownership Changes (Refine Rdy, Iter 2.5)",
             "At risk - 2.5", AMBER_K),
            ("CI-owner access governance", "(no feature home)",
             "1399787 Add all CI-Owners to CMDB Full Access Group (Validation)",
             "In progress", AMBER_K),
        ],
        table_h=Inches(2.1),
        co6="CO6 (draft) carries SOX review under 'CMDB Governance' -> Jul 31. Same SOX-only boundary (NERC/CIP excluded).",
        notes=[
            "Contractual boundary: SOX data ONLY. NERC-CIP (Stream B, feature 1370224) is separate and NOT part of this acceptance item.",
            "1455827 sits in Iter 2.5 (post-6/30) with an open Joe decision: PI-2 or PI-3? Resolve to avoid a deadline miss.",
        ],
    ),
    dict(
        title="Governance 1.5 - Monthly CCB + Future-PI Backlog",
        subtitle="Cadence is running, but the ADO evidence is broken (hard gap G5)",
        col0="Area",
        contract="Facilitate Monthly CCB Meetings and create a backlog for future PIs.",
        rows=[
            ("Monthly CCB facilitation  [G5]", "1406668 CMDB Governance & Monthly CCB - REMOVED in ADO",
             "1406672 May CCB · 1406683 June CCB · 1406687 July CCB (all New, Joe; under Removed feature)",
             "Gap - broken evidence", RED_K),
            ("CI-owner governance (leavers)", "1355888 Operational Monitoring",
             "1339116 Monitor / Update CI Owners Who Left PPL (Refine Rdy)", "In progress", AMBER_K),
            ("Backlog for future PIs", "-", "artifact not yet tracked", "TBD", GREY_K),
        ],
        table_h=Inches(2.1),
        co6="NOT carried as a CO6 deliverable - only in the CMDB PO role text. Remains a CO5-only obligation.",
        notes=[
            "CCB cadence IS happening: CCB meeting 6/16, CMDB Audit Sync 6/18, Manuel attending - but the feature that proves it (1406668) is marked Removed in ADO.",
            "Action (high priority): un-Remove 1406668 OR re-parent the 3 CCB stories to a live feature, so the deliverable is demonstrable for acceptance.",
            "Confirm the 'backlog for future PIs' artifact exists and is tracked.",
        ],
    ),
]


def section_divider(title, subtitle):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, EMU_W, EMU_H, DARK)
    rect(s, 0, Inches(3.75), EMU_W, Inches(0.08), ACCENT)
    textbox(s, Inches(0.8), Inches(2.7), Inches(11.7), Inches(1.0),
            [(title, 40, WHITE, True, False)])
    textbox(s, Inches(0.82), Inches(4.0), Inches(11.7), Inches(0.6),
            [(subtitle, 18, RGBColor(0xC9, 0xD6, 0xE6), False, False)])
    return s


def gov_slide(spec):
    s = prs.slides.add_slide(BLANK)
    header(s, spec["title"], spec["subtitle"])
    # contract strip (amber if forward-looking / not CO5)
    fwd = spec.get("forward", False)
    strip_fill = RGBColor(0xFB, 0xEE, 0xDB) if fwd else RGBColor(0xE9, 0xF0, 0xF8)
    strip_line = AMBER if fwd else ACCENT
    strip_text = RGBColor(0x7A, 0x4F, 0x10) if fwd else DARK
    rect(s, Inches(0.55), Inches(1.22), Inches(12.2), Inches(0.62), strip_fill, line=strip_line)
    textbox(s, Inches(0.72), Inches(1.25), Inches(11.9), Inches(0.58),
            [(spec.get("contract_label", "CO5 requires:  ") + spec["contract"], 10.5, strip_text, False, False)],
            MSO_ANCHOR.MIDDLE)
    # table
    table_h = spec["table_h"]
    rows = len(spec["rows"]) + 1
    tshape = s.shapes.add_table(rows, 4, Inches(0.55), Inches(2.0), Inches(12.2), table_h)
    t = tshape.table
    for j, wdt in enumerate([Inches(2.9), Inches(2.9), Inches(4.9), Inches(1.5)]):
        t.columns[j].width = wdt
    for j, h in enumerate([spec["col0"], "PI-2 Feature(s)", "Related stories (ID - status)", "Status"]):
        cell(t, 0, j, h, 11, WHITE, True, PP_ALIGN.LEFT if j < 3 else PP_ALIGN.CENTER, DARK)
    for i, (label, feat, stories, status, scol) in enumerate(spec["rows"], start=1):
        fill = LIGHT if i % 2 else WHITE
        cell(t, i, 0, label, 9.5, DARK, True, PP_ALIGN.LEFT, fill)
        cell(t, i, 1, feat, 9, TEXT, False, PP_ALIGN.LEFT, fill)
        cell(t, i, 2, stories, 9, TEXT, False, PP_ALIGN.LEFT, fill)
        cell(t, i, 3, status, 9.5, WHITE, True, PP_ALIGN.CENTER, RAG[scol])
    yb = Inches(2.0) + table_h + Inches(0.16)
    # CO6 effect strip
    rect(s, Inches(0.55), yb, Inches(12.2), Inches(0.5), CO6_GREEN_FILL, line=GREEN)
    textbox(s, Inches(0.72), yb + Inches(0.02), Inches(11.9), Inches(0.46),
            [(spec.get("co6_label", "CO6 effect:  ") + spec["co6"], 10.5, CO6_GREEN_TEXT, False, False)],
            MSO_ANCHOR.MIDDLE)
    # notes
    yy = yb + Inches(0.66)
    textbox(s, Inches(0.55), yy, Inches(12.2), Inches(0.3),
            [("Notes / actions:", 11.5, DARK, True, False)])
    yy += Inches(0.36)
    for n in spec["notes"]:
        rect(s, Inches(0.55), yy + Inches(0.04), Inches(0.1), Inches(0.34), ACCENT)
        textbox(s, Inches(0.78), yy, Inches(11.9), Inches(0.5), [(n, 10, TEXT, False, False)])
        yy += Inches(0.46)
    return s


section_divider("Governance - Deep Dive", "Deliverable 1 broken down per CO5 sub-item  (1.1 - 1.5)")
for spec in GOV:
    gov_slide(spec)

# ---------------------------------------------------------------------------
# AUTOMATED DATA INGESTION DEEP-DIVE  (Deliverable 2 by sub-item) + CO6 look-ahead
# ---------------------------------------------------------------------------
D2_GOV = [
    dict(
        title="Data Ingestion 2.1 - Computers",
        subtitle="SCCM/Discovery precedence + bulk life-cycle + 90% coverage",
        col0="CO5 component",
        contract=("SCCM/Discovery data precedence reconciliation & updates; bulk updates for Life Cycle Stage & Status "
                  "(process defined by Asset Mgmt team); coverage validation - goal of 90% devices managed."),
        rows=[
            ("SCCM/Discovery precedence", "1354794 Computer Class / SCCM Computer Class Precedence",
             "1348712 Domain&Network (Resolved) · 1348716 HW&SW (Resolved) · 1348717 De-Prioritize SCCM (Resolved) · 1348715 Last-Seen (Validation)",
             "On track", GREEN_K),
            ("Bulk Life Cycle Stage/Status", "Computer Class 1354794",
             "1402790 Update Lifecycle for Retired Computers (Closed)", "Done", GREEN_K),
            ("90% coverage validation  [G2]", "- none -",
             "NO MEASUREMENT STORY", "GAP", RED_K),
        ],
        table_h=Inches(2.1),
        co6=("CO6 'CI Coverage - Computers' -> precedence & life-cycle Jul 31, 90% coverage validation Sep 30. "
             "Closes G2 (if signed)."),
        notes=[
            "CO5 RISK: no reliable source to compare physical devices; only on-network devices are discovered.",
            "Life Cycle process is an external dependency on the Asset Management team - not a team-owned story.",
            "90% coverage has no measurement story (G2) - CO6 gives it a Sep 30 target.",
        ],
    ),
    dict(
        title="Data Ingestion 2.2 - Servers",
        subtitle="SCCM/Discovery precedence + credentials + 90% non-NERC-CIP coverage",
        col0="CO5 component",
        contract=("SCCM/Discovery data precedence reconciliation & updates; coverage validation - goal of 90% "
                  "non-NERC-CIP devices discovered; SOX indicators manually maintained (excluded from automation)."),
        rows=[
            ("SCCM/Discovery precedence", "1356826 SCCM Server Class Precedence",
             "1403763 HW&SW (Validation) · 1403760 Domain&Network (Validation) · 1403762 Last-Seen (Validation) · 1403759 Remove SCCM Attrs (Validation)",
             "Validation", AMBER_K),
            ("Credentials (enables discovery)", "1356646 / 1354797",
             "1444864 Fix Creds Svr/DB/Net (Validation; 6 child tasks still Active) · 1459721 SNMP/MID test (Validation)",
             "In progress", AMBER_K),
            ("90% non-NERC-CIP coverage  [G3]", "- none -",
             "NO MEASUREMENT STORY", "GAP", RED_K),
            ("SOX indicators", "n/a",
             "Manually maintained - EXCLUDED from automation per CO5", "Excluded", GREY_K),
        ],
        table_h=Inches(2.7),
        co6=("CO6 'CI Coverage - Servers' -> precedence & enhanced DB discovery Jul 31, 90% coverage Oct 30. "
             "Closes G3 (if signed)."),
        notes=[
            "SCCM precedence baseline RISK (Risk #12): no complete field-by-field precedence map - rework risk across stories.",
            "Credential work gates discovery for servers AND databases; 6 tasks still Active in Iter 2.2.",
            "90% non-NERC-CIP coverage has no measurement story (G3).",
        ],
    ),
    dict(
        title="Data Ingestion 2.3 - Databases",
        subtitle="Enhanced Discovery for MS-SQL & Oracle - credential-bound",
        col0="CO5 component",
        contract=("Enhanced Discovery support for MS-SQL and Oracle databases (RISK - no scalable solution for the "
                  "distribution of credentials to target CIs); SOX indicators manually maintained (excluded)."),
        rows=[
            ("Enhanced Discovery (MS-SQL/Oracle)", "1356646 / credential work",
             "Oracle DB cred tasks under 1444864 (GWIZ_MON_ORACLE_DB cred; PA Oracle DB discovery test) - task-level only",
             "Task-level only", AMBER_K),
            ("SOX indicators", "n/a",
             "Manually maintained - EXCLUDED from automation per CO5", "Excluded", GREY_K),
        ],
        table_h=Inches(1.5),
        co6="CO6 'CI Coverage - Servers' carries enhanced DB Discovery (MS-SQL/Oracle) -> Jul 31 (same RISK language).",
        notes=[
            "CO5 RISK (verbatim): no scalable solution for the distribution of credentials to target CIs - your active credential/SolarWinds workstream.",
            "No discrete 'enhanced DB Discovery' story - it lives only as credential tasks. Acceptance needs a demonstrable deliverable.",
            "Discovery now surfacing Oracle DBs / listeners (6/9) after credential fixes - momentum is real.",
        ],
    ),
    # CO6 look-ahead (NOT a CO5 deliverable) - the bridge into CO6
    dict(
        title="Look-Ahead: Network Gear Discovery",
        subtitle="NOT a CO5 deliverable - PI-2 Objective 1, formalized as a CO6 (draft) deliverable",
        col0="Component",
        forward=True,
        contract_label="CO6 (draft) requires:  ",
        contract=("Enable & configure Network Gear Discovery; devices discovered & populated in CMDB. "
                  "Staged acceptance: creds configured/validated (excl. OT) & schedules active Aug 31; mandatory attrs Sep 30; "
                  "90% coverage (business-owner validated) Oct 30."),
        rows=[
            ("Network device coverage", "1356646 Network Device Coverage Reconciliation - Group 1",
             "1402555 Identify CI Classes (Active spike) · 1402559 Compare CI Class Sources (New spike) · 1402572/574/575 discovery & device config + rerun (New)",
             "In progress", AMBER_K),
            ("Credentials / SNMP", "1356646",
             "1444864 Fix Creds (Validation) · 1459721 SNMP/MID test (Validation)", "In progress", AMBER_K),
            ("Stakeholder requirements", "dependency",
             "1383487 Network-gear stakeholder requirements to assess (New, Dan Carabelas)", "New", AMBER_K),
        ],
        table_h=Inches(2.1),
        co6_label="Staged acceptance:  ",
        co6="CO6 (draft) Deliverable #1 - Aug 31 / Sep 30 / Oct 30 (PI-3). 90% coverage business-owner validated by Oct 30. OT explicitly excluded.",
        notes=[
            "Network gear is a PI-2 objective (Obj 1) but was NOT a CO5 deliverable - CO6 makes it contractual.",
            "Story scaffolding exists (spikes + config stories under 1356646) but no contractual acceptance is met yet.",
            "Shown here as a precursor/bridge to CO6 - do not report as current CO5 scope.",
        ],
    ),
]

section_divider("Automated Data Ingestion - Deep Dive",
                "Deliverable 2 by sub-item (Computers / Servers / Databases)  +  a CO6 look-ahead: Network Gear")
for spec in D2_GOV:
    gov_slide(spec)

# ---------------------------------------------------------------------------
# CO6 NET-NEW SCOPE  (short, TENTATIVE - based on the unsigned 3/27 draft)
# ---------------------------------------------------------------------------
# (deliverable, what it adds, due (draft), story status today, note)
CO6_NEW = [
    ("Network Gear Discovery", "Discover & populate network devices in CMDB (excl. OT); 90% coverage",
     "Aug 31 - Oct 30", "Scaffolding under 1356646 (spikes + config stories)", "See slide 16 look-ahead"),
    ("Service Mapping", "E2E maps for 10 priority business apps down to infra CIs; owner-validated",
     "Aug 31 - Oct 30", "Wave features 1355866/68/71 + per-app stories", "Was PI-2 obj; now contractual"),
    ("Legacy Platform Rationalization", "Analysis + migration plan: iTeam, DISCO/Cherwell, AIM",
     "Aug 31 - Oct 30", "Only iTeam import 1452028", "Cherwell/AIM 'if applicable' at PI-3"),
    ("ITSM Product Management", "New ITSM Product Owner role + ITSM workstream",
     "Monthly Jul 31 - Oct 30", "None - owner TBD", "Outside CMDB; ties to CR6 role seg"),
    ("ATF Strategy", "Automated Test Framework rollout plan across in-prod capabilities",
     "Oct 31", "None - greenfield", "Net-new analysis deliverable"),
    ("Platform Support", "Dedicated BAU/DevOps team; 40 hrs stories/wk per member",
     "Monthly Aug 31 - Oct 30", "None - team size TBD ('XXX')", "Ongoing capacity, not finite"),
]
CO6_NEW_NOTE = ("Net-new != free: ITSM PM + Platform Support are ongoing-capacity commitments, not finite deliverables.  "
                "CO6 also RE-BASELINES the CO5 items (Governance, CI Coverage, Qualys) - those live in the gap-closure / "
                "timeline-shift maps, not here.")

section_divider("CO6 (Draft) - Net-New Scope",
                "TENTATIVE - based on the unsigned 3/27 draft.  To be updated once the signed contract is reviewed.")

s = prs.slides.add_slide(BLANK)
header(s, "CO6 Net-New Deliverables", "Brand-new scope with no CO5 antecedent  -  6 items, mostly PI-3")
# TENTATIVE banner
rect(s, Inches(0.55), Inches(1.2), Inches(12.2), Inches(0.42), RGBColor(0xFB, 0xEE, 0xDB), line=AMBER)
textbox(s, Inches(0.72), Inches(1.22), Inches(11.9), Inches(0.38),
        [("TENTATIVE  -  DRAFT / UNSIGNED (3/27).  Numbers, dates and scope will change once the signed contract is reviewed.",
          11, RGBColor(0x7A, 0x4F, 0x10), True, False)], MSO_ANCHOR.MIDDLE)
rows = len(CO6_NEW) + 1
tshape = s.shapes.add_table(rows, 5, Inches(0.55), Inches(1.78), Inches(12.2), Inches(3.7))
t = tshape.table
for j, wdt in enumerate([Inches(2.6), Inches(3.5), Inches(1.7), Inches(2.5), Inches(1.9)]):
    t.columns[j].width = wdt
for j, h in enumerate(["CO6 Deliverable", "What it adds", "Due (draft)", "Story status today", "Note"]):
    cell(t, 0, j, h, 10.5, WHITE, True, PP_ALIGN.LEFT, DARK)
for i, (name, adds, due, story, note) in enumerate(CO6_NEW, start=1):
    fill = LIGHT if i % 2 else WHITE
    cell(t, i, 0, name, 10, DARK, True, PP_ALIGN.LEFT, fill)
    cell(t, i, 1, adds, 9, TEXT, False, PP_ALIGN.LEFT, fill)
    cell(t, i, 2, due, 9, AMBER, True, PP_ALIGN.LEFT, fill)
    cell(t, i, 3, story, 9, TEXT, False, PP_ALIGN.LEFT, fill)
    cell(t, i, 4, note, 8.5, MUTED, False, PP_ALIGN.LEFT, fill)
textbox(s, Inches(0.55), Inches(5.65), Inches(12.2), Inches(0.9),
        [(CO6_NEW_NOTE, 10.5, TEXT, False, False)])

# ---------------------------------------------------------------------------
# SPRINTS vs 7/31  (timeline reality for the CO6 re-baselined date)
# ---------------------------------------------------------------------------
# (iteration, window, relation to 7/31, freeze overlay, freeze color, highlight-relation)
SPRINT_ROWS = [
    ("2.3", "Jun 10 - 23", "Active now - last fully-open sprint", "none", GREEN_K, False),
    ("2.4", "Jun 24 - Jul 7", "CO5 6/30 deadline here · clean dev window only Jun 24 - Jul 3",
     "Dev freeze from Jul 4", AMBER_K, False),
    ("2.5", "Jul 8 - 21", "Runway to 7/31 - squeezed both ends",
     "Dev freeze <= Jul 18; Test freeze >= Jul 18", RED_K, False),
    ("2.6 (IP)", "Jul 22 - Aug 4", "7/31 LANDS HERE - IP iteration, not a delivery sprint",
     "Test freeze Jul 18 - Aug 15", RED_K, True),
]
SPRINT_TAKE = ("7/31 sits in the IP iteration, bracketed by two freezes.  The last realistic window to land a 7/31 "
               "deliverable is end of Iter 2.4 (~Jul 3, before the Dev freeze) - so CO6's 7/31 work must be "
               "substantially complete in Iter 2.3-2.4, not 'sometime in July.'")
SPRINT_NOTE = ("The 90% coverage items (Sep 30 / Oct 30) are the ones genuinely given room - PI-3, past the test freeze.  "
               "Open question: is the IP iteration (2.6) expected to carry CO6 acceptance, or should 7/31 work be pulled "
               "into 2.4/2.5 planning?")

s = prs.slides.add_slide(BLANK)
header(s, "Sprints vs the 7/31 Date", "CO6 re-dates Governance & CI Coverage to Jul 31 - but the sprint calendar tells a tighter story")
rows = len(SPRINT_ROWS) + 1
tshape = s.shapes.add_table(rows, 4, Inches(0.55), Inches(1.35), Inches(12.2), Inches(2.7))
t = tshape.table
for j, wdt in enumerate([Inches(1.4), Inches(1.9), Inches(5.4), Inches(3.5)]):
    t.columns[j].width = wdt
for j, h in enumerate(["Iteration", "Window", "Relation to CO6 7/31", "Freeze overlay"]):
    cell(t, 0, j, h, 11.5, WHITE, True, PP_ALIGN.LEFT, DARK)
for i, (it, win, rel, frz, fcol, hot) in enumerate(SPRINT_ROWS, start=1):
    fill = LIGHT if i % 2 else WHITE
    cell(t, i, 0, it, 11, DARK, True, PP_ALIGN.LEFT, fill)
    cell(t, i, 1, win, 10.5, TEXT, False, PP_ALIGN.LEFT, fill)
    cell(t, i, 2, rel, 10, RED if hot else TEXT, hot, PP_ALIGN.LEFT, fill)
    cell(t, i, 3, frz, 9.5, WHITE, True, PP_ALIGN.LEFT, RAG[fcol])
# takeaway callout
rect(s, Inches(0.55), Inches(4.35), Inches(12.2), Inches(1.15), RGBColor(0xFB, 0xEE, 0xDB), line=AMBER)
textbox(s, Inches(0.75), Inches(4.44), Inches(11.9), Inches(0.32),
        [("Timeline reality", 13, RGBColor(0x7A, 0x4F, 0x10), True, False)])
textbox(s, Inches(0.75), Inches(4.78), Inches(11.9), Inches(0.66),
        [(SPRINT_TAKE, 11, RGBColor(0x5A, 0x4A, 0x2A), False, False)])
# secondary note
textbox(s, Inches(0.55), Inches(5.75), Inches(12.2), Inches(0.9),
        [(SPRINT_NOTE, 10.5, TEXT, False, False)])
textbox(s, Inches(0.55), Inches(6.7), Inches(12.2), Inches(0.32),
        [("Freeze source: Dev clone Jun 27 / upgrade Jul 4 (CHG70100870)  ·  Test freeze CHG70100865  -  see PI-2/Memory Risks #10/#11",
          9, MUTED, False, True)])


def save_safe(path):
    try:
        prs.save(path); print("Saved", path)
    except PermissionError:
        base, ext = os.path.splitext(path)
        alt = base + "-new" + ext
        prs.save(alt); print("Target locked - saved to", alt)


out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "CO5-Deliverables-" + EDITION + ".pptx")
save_safe(out)
